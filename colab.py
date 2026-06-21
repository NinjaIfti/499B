# ╔══════════════════════════════════════════════════════════════╗
# ║         LectureForge v1.1  —  Backend (API-only mode)          ║
# ║  Run in Google Colab with T4 GPU                            ║
# ║  UI → open quizforge_ui.html locally on your PC            ║
# ╚══════════════════════════════════════════════════════════════╝
#
# RUN ORDER:
#   1. Run Cell 1 (Install)  →  Runtime ▸ Restart Runtime
#   2. Set NGROK_AUTH_TOKEN below
#   3. Run Cell 2 (this file) — everything else is automatic
#   4. Copy the printed ngrok URL into quizforge_ui.html

# ══════════════════════════════════════════════════════════════
# CONFIG
# ══════════════════════════════════════════════════════════════



NGROK_AUTH_TOKEN = "3CYotz762dE9Q22K2Od7hH57ngW_5ipi8fNnpayeiS9ytkQ14"   # get free token at ngrok.com
OLLAMA_MODEL     = "gemma3:12b"              # primary model — summaries, chat, grading
OLLAMA_FALLBACK  = "gemma3:4b"               # fallback if primary fails
QUIZ_MODEL       = "gemma3:4b"               # fast model for quiz generation
QUIZ_CTX         = 16384                     # smaller context = faster quiz inference
OLLAMA_BASE      = "http://localhost:11434"
UPLOAD_DIR       = "/content/uploads"
OUTPUT_DIR       = "/content/outputs"
CLIPS_DIR        = "/content/outputs/clips"



# Google Drive persistence — survives Colab disconnects
DRIVE_DIR        = "/content/drive/MyDrive/LectureForge"
SESSIONS_DIR     = "/content/drive/MyDrive/LectureForge/sessions"
OCR_CACHE_DIR    = "/content/drive/MyDrive/LectureForge/ocr_cache"

# EasyOCR languages — add more codes from https://www.jaided.ai/easyocr/
OCR_LANGUAGES    = ["en", "bn"]          # English + Bengali

# ══════════════════════════════════════════════════════════════
# IMPORTS
# ══════════════════════════════════════════════════════════════



!pip install -q torch torchvision torchaudio --index-url https://download.pytorch.org/whl/cpu
!pip install -q easyocr
!pip install -q flask pyngrok faster-whisper sentence-transformers faiss-cpu pillow numpy werkzeug
!pip install -q python-docx python-pptx PyMuPDF
!apt-get install -qq ffmpeg curl
!pip install -q flask-cors
!apt-get install -qq zstd        # ← add this line
!curl -fsSL https://ollama.com/install.sh | sh
print("✅ All packages installed. Now go to Runtime → Restart Runtime, then run from Cell 2.")



import os, json, re, time, copy, difflib, threading, subprocess, shutil, hashlib
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed



import cv2
import numpy as np
import easyocr
import faiss
import requests as http_requests
from faster_whisper import WhisperModel
from sentence_transformers import SentenceTransformer
from flask import Flask, request, jsonify, send_file, Response
from flask_cors import CORS
from werkzeug.utils import secure_filename



# Mount Google Drive for persistent storage
from google.colab import drive
drive.mount('/content/drive', force_remount=False)



for d in [UPLOAD_DIR, OUTPUT_DIR, CLIPS_DIR, DRIVE_DIR, SESSIONS_DIR, OCR_CACHE_DIR]:
    os.makedirs(d, exist_ok=True)



ALLOWED_VIDEO = {'mp4', 'avi', 'mov', 'mkv', 'webm'}
ALLOWED_DOCS  = {'pdf', 'docx', 'pptx', 'ppt'}



print("✅ Imports & config ready.")
print(f"✅ Google Drive mounted — sessions saved to {SESSIONS_DIR}")

# ══════════════════════════════════════════════════════════════
# APP STATE — single place for all runtime data
# ══════════════════════════════════════════════════════════════



class AppState:
    """Keeps every piece of mutable runtime data in one object
    instead of scattered globals."""

    def __init__(self):
        self.lock         = threading.Lock()
        self.active_model = OLLAMA_MODEL
        self.video_path   = ""
        self.doc_path     = ""
        self.suggested_questions = []
        self.session_id   = ""
        self.agent_log    = []
        self.student_performance = {}    # topic -> list of 0.0/1.0 scores
        self._quiz_agent_thinking = []   # thinking log for manual/adaptive quiz agents
        self._quiz_thinking_session = 0  # bumps each time a fresh agentic generation starts
        self._quiz_thinking_active  = False  # true while the agentic pipeline is running
        self.quiz_history = []          # list of graded quiz submissions saved per session
        self.generated_quiz_runs = []   # snapshots of each generation (async + upload) for UI history
        self.quiz_generation_jobs = {} # quiz generation jobs for live thinking
        self.quiz_generation_lock = threading.Lock()

        self.status = {
            "state": "idle", "stage": "Idle", "pct": 0, "error": "",
            "summary": "", "transcript": "", "board_text": "", "lecture_timeline": "",
        }

        self.quiz = {
            "mcq": [], "tf": [], "fill": [], "short": [],
            "flash": [], "hints": {},
        }

    def update_status(self, **kw):
        with self.lock:
            self.status.update(kw)

    def reset_status(self):
        with self.lock:
            self.status.update(
                state="running", stage="Upload received…", pct=2, error="",
                summary="", transcript="", board_text="", lecture_timeline="",
            )
            self._quiz_agent_thinking   = []
            self._quiz_thinking_active  = False
            self._quiz_thinking_session += 1
            self.generated_quiz_runs = []  # new material — clear past-set history

    # ── Persistence ───────────────────────────────────────────

    def _make_session_id(self, filename):
        name = Path(filename).stem
        short_hash = hashlib.md5(f"{name}{time.time()}".encode()).hexdigest()[:6]
        date_str = datetime.now().strftime("%Y%m%d_%H%M")
        return f"{date_str}_{name}_{short_hash}"

    def save_session(self, source_filename):
        """Dump full session to Google Drive, including performance and source file."""
        with self.lock:
            if not self.session_id:
                self.session_id = self._make_session_id(source_filename)

            folder = os.path.join(SESSIONS_DIR, self.session_id)
            os.makedirs(folder, exist_ok=True)

            payload = {
                "session_id":            self.session_id,
                "source_file":           source_filename,
                "saved_at":              datetime.now().isoformat(),
                "summary":               self.status.get("summary", ""),
                "transcript":            self.status.get("transcript", ""),
                "board_text":            self.status.get("board_text", ""),
                "lecture_timeline":      self.status.get("lecture_timeline", ""),
                "suggested_questions":   self.suggested_questions,
                "quiz":                  self.quiz,
                "quiz_history":         self.quiz_history,
                "generated_quiz_runs":   self.generated_quiz_runs,
                "student_performance":   self.student_performance,
            }

            session_path = os.path.join(folder, "session.json")
            with open(session_path, "w", encoding="utf-8") as f:
                json.dump(payload, f, indent=2, ensure_ascii=False)

            # copy text artifacts
            for fname in ["summary.txt", "transcript.txt", "board_text.txt",
                           "lecture_timeline.txt", "board_entries.json", "transcript_segments.json"]:
                src = os.path.join(OUTPUT_DIR, fname)
                if os.path.exists(src):
                    shutil.copy2(src, os.path.join(folder, fname))

            # copy the original uploaded file so /stream/video works after restore
            source_path = self.video_path or self.doc_path
            if source_path and os.path.exists(source_path):
                dest = os.path.join(folder, Path(source_path).name)
                if not os.path.exists(dest):
                    shutil.copy2(source_path, dest)

        print(f"💾 Session saved → {folder}")

    def load_session(self, session_id):
        """Restore a previously saved session from Drive."""
        folder = os.path.join(SESSIONS_DIR, session_id)
        session_path = os.path.join(folder, "session.json")
        if not os.path.exists(session_path):
            print(f"  Session {session_id} not found.")
            return False

        try:
            with open(session_path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except (json.JSONDecodeError, OSError) as e:
            print(f"  Corrupt session file: {e}")
            return False

        with self.lock:
            self.session_id          = session_id
            self.agent_log           = []
            self.suggested_questions = data.get("suggested_questions", [])
            self.quiz                = data.get("quiz", self.quiz)
            self.quiz_history       = data.get("quiz_history", [])
            self.generated_quiz_runs = data.get("generated_quiz_runs", [])
            self.status.update(
                state="done", stage="Restored from Drive ✓", pct=100, error="",
                summary=data.get("summary", ""),
                transcript=data.get("transcript", ""),
                board_text=data.get("board_text", ""),
                lecture_timeline=data.get("lecture_timeline", ""),
            )

            # merge performance history (accumulate, don't replace)
            for topic, scores in data.get("student_performance", {}).items():
                if topic not in self.student_performance:
                    self.student_performance[topic] = []
                self.student_performance[topic].extend(scores)

            # restore video/doc path from the saved copy
            self.video_path = ""
            self.doc_path   = ""
            source_file = data.get("source_file", "")
            if source_file:
                saved_copy = os.path.join(folder, source_file)
                if os.path.exists(saved_copy):
                    ext = Path(source_file).suffix.lower().lstrip('.')
                    if ext in ALLOWED_VIDEO:
                        self.video_path = saved_copy
                    elif ext in ALLOWED_DOCS:
                        self.doc_path = saved_copy

        # copy text files back to OUTPUT_DIR so /download endpoints work
        for fname in ["summary.txt", "transcript.txt", "board_text.txt",
                       "lecture_timeline.txt", "board_entries.json", "transcript_segments.json"]:
            src = os.path.join(folder, fname)
            if os.path.exists(src):
                shutil.copy2(src, os.path.join(OUTPUT_DIR, fname))

        with self.lock:
            if not self.status.get("lecture_timeline") and self.video_path:
                tl = rebuild_lecture_timeline_from_disk()
                if tl:
                    self.status["lecture_timeline"] = tl

        print(f"✅ Session restored: {session_id}")
        return True

    @staticmethod
    def list_sessions():
        sessions = []
        if not os.path.isdir(SESSIONS_DIR):
            return sessions
        for name in sorted(os.listdir(SESSIONS_DIR), reverse=True):
            sp = os.path.join(SESSIONS_DIR, name, "session.json")
            if os.path.exists(sp):
                try:
                    with open(sp, "r", encoding="utf-8") as f:
                        meta = json.load(f)
                    sessions.append({
                        "session_id":  name,
                        "source_file": meta.get("source_file", ""),
                        "saved_at":    meta.get("saved_at", ""),
                        "has_summary": bool(meta.get("summary")),
                        "mcq_count":   len(meta.get("quiz", {}).get("mcq", [])),
                        "tf_count":    len(meta.get("quiz", {}).get("tf", [])),
                        "flash_count": len(meta.get("quiz", {}).get("flash", [])),
                    })
                except Exception:
                    pass
        return sessions



state = AppState()

# ══════════════════════════════════════════════════════════════
# KEEP-ALIVE  (prevents Colab disconnecting after 90 min idle)
# ══════════════════════════════════════════════════════════════



def _keep_alive():
    while True:
        try:
            with open('/content/.heartbeat', 'w') as f:
                f.write(str(time.time()))
        except Exception:
            pass
        time.sleep(45)



threading.Thread(target=_keep_alive, daemon=True).start()
print("✅ Keep-alive thread started.")

# ══════════════════════════════════════════════════════════════
# LAZY MODEL LOADERS
# ══════════════════════════════════════════════════════════════



_ocr_reader = _whisper_model = _embed_model = None



def _gpu_available():
    """Check once whether CUDA is usable."""
    try:
        import torch
        return torch.cuda.is_available()
    except ImportError:
        return False



def get_ocr():
    global _ocr_reader
    if _ocr_reader is None:
        gpu = _gpu_available()
        print(f"  Loading EasyOCR ({OCR_LANGUAGES}) on {'GPU' if gpu else 'CPU'}…")
        _ocr_reader = easyocr.Reader(OCR_LANGUAGES, gpu=gpu)
    return _ocr_reader



def get_whisper():
    global _whisper_model
    if _whisper_model is None:
        # always use small — fits next to Ollama on T4; medium/large OOM easily
        for device, ctype in [("cuda", "float16"), ("cpu", "int8")]:
            if device == "cuda" and not _gpu_available():
                continue
            try:
                print(f"  Loading Whisper 'small' on {device.upper()} ({ctype})…")
                _whisper_model = WhisperModel("small", device=device, compute_type=ctype)
                print(f"  ✅ Whisper 'small' on {device.upper()}")
                return _whisper_model
            except Exception as e:
                print(f"  ⚠ Whisper small on {device} failed: {e}")
                _whisper_model = None
        raise RuntimeError("Could not load Whisper small on GPU or CPU")
    return _whisper_model



def get_embedder():
    global _embed_model
    if _embed_model is None:
        print("  Loading SentenceTransformer all-MiniLM-L6-v2…")
        _embed_model = SentenceTransformer('all-MiniLM-L6-v2')
    return _embed_model



print("✅ Lazy loaders defined.")

# ══════════════════════════════════════════════════════════════
# OLLAMA SETUP  (start server + pull model)
# ══════════════════════════════════════════════════════════════



def start_ollama():
    env = {
        **os.environ,
        "OLLAMA_KEEP_ALIVE":    "-1",
        "CUDA_VISIBLE_DEVICES": "0",
        "OLLAMA_GPU_OVERHEAD":  "500000000",
    }
    subprocess.Popen(
        ["ollama", "serve"],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, env=env,
    )
    time.sleep(5)
    print("✅ Ollama server started (GPU mode).")



def pull_model(name):
    print(f"  Pulling {name}… (~7 GB first time)")
    r = subprocess.run(["ollama", "pull", name], capture_output=False)
    if r.returncode != 0:
        return False
    # Pre-warm with 32 K context window
    subprocess.run(
        ["ollama", "run", name, "--", "hi"],
        capture_output=True, text=True,
        env={**os.environ, "OLLAMA_NUM_CTX": "32768"},
    )
    print(f"✅ {name} ready | context: 32K tokens")
    return True



def setup_ollama():
    start_ollama()
    if pull_model(OLLAMA_MODEL):
        state.active_model = OLLAMA_MODEL
    elif pull_model(OLLAMA_FALLBACK):
        state.active_model = OLLAMA_FALLBACK
    else:
        print("❌ No model available — check internet connection.")
    # pull the fast quiz model if it's a different model
    if QUIZ_MODEL and QUIZ_MODEL != state.active_model:
        print(f"  Pulling fast quiz model ({QUIZ_MODEL})…")
        pull_model(QUIZ_MODEL)
    print(f"✅ Active model: {state.active_model} | Quiz model: {QUIZ_MODEL or state.active_model}")



setup_ollama()

# ══════════════════════════════════════════════════════════════
# LLM WRAPPER
# ══════════════════════════════════════════════════════════════



def call_ollama(prompt, system="You are an expert educational AI assistant.",
                json_mode=False, max_retries=3, timeout=300,
                model=None, num_ctx=None):
    """Send a prompt to the local Ollama instance and return the response text."""
    payload = {
        "model":   model or state.active_model,
        "prompt":  prompt,
        "system":  system,
        "stream":  False,
        "options": {"num_ctx": num_ctx or 32768, "temperature": 0.2},
    }
    if json_mode:
        payload["format"] = "json"

    for attempt in range(1, max_retries + 1):
        try:
            resp = http_requests.post(
                f"{OLLAMA_BASE}/api/generate", json=payload, timeout=timeout,
            )
            resp.raise_for_status()
            text = resp.json().get("response", "").strip()
            if json_mode:
                text = re.sub(r'^```(?:json)?\s*', '', text)
                text = re.sub(r'\s*```$', '', text)
            return text
        except Exception as e:
            if attempt == max_retries:
                raise
            print(f"  LLM attempt {attempt} failed: {e} — retrying…")
            time.sleep(2)



def call_ollama_json(prompt, system="You are an expert educational AI assistant.",
                     fallback=None, model=None, num_ctx=None):
    """call_ollama with JSON mode. Returns parsed dict/list, or *fallback* on error."""
    try:
        raw = call_ollama(prompt, system=system, json_mode=True,
                          model=model, num_ctx=num_ctx)
        return json.loads(raw)
    except Exception as e:
        print(f"  JSON parse error: {e}")
        return fallback if fallback is not None else {}



def call_ollama_json_quiz(prompt, key="questions"):
    """Try fast quiz model first; if JSON is empty or invalid, retry with primary model."""
    fb = {key: []}
    data = call_ollama_json(prompt, fallback=fb, model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
    if not isinstance(data, dict):
        data = dict(fb)
    items = data.get(key) or []
    if items:
        return data
    print(f"  ⚠ Empty or bad '{key}' from {QUIZ_MODEL} — retrying with {state.active_model}…")
    data2 = call_ollama_json(prompt, fallback=fb)
    return data2 if isinstance(data2, dict) else fb



def call_ollama_json_list_quiz(prompt, fallback_list):
    """JSON list (e.g. suggested questions): try quiz model, then primary."""
    data = call_ollama_json(prompt, fallback=fallback_list,
                            model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
    if isinstance(data, list) and data:
        return data
    print(f"  ⚠ Empty list from {QUIZ_MODEL} — retrying with {state.active_model}…")
    data2 = call_ollama_json(prompt, fallback=fallback_list)
    return data2 if isinstance(data2, list) else fallback_list



print("✅ LLM wrapper ready.")

# ══════════════════════════════════════════════════════════════
# SHARED HELPERS
# ══════════════════════════════════════════════════════════════



def fmt_ts(seconds):
    """Convert seconds → 'MM:SS' string."""
    m, s = divmod(int(seconds), 60)
    return f"{m:02d}:{s:02d}"



def nearest_timestamp(segment_list, index):
    """Pick the segment timestamp closest to a question's ordinal index."""
    if not segment_list:
        return ""
    seg_idx = min(
        int(index / max(len(segment_list), 1) * len(segment_list)),
        len(segment_list) - 1,
    )
    return segment_list[seg_idx].get("start_str", "")



def fill_missing_timestamps(questions, segments=None, board_entries=None):
    """Back-fill source_timestamp on questions that don't have one."""
    for i, q in enumerate(questions):
        if q.get("source_timestamp"):
            continue
        if segments:
            q["source_timestamp"] = nearest_timestamp(segments, i)
        elif board_entries:
            q["source_timestamp"] = board_entries[min(i, len(board_entries) - 1)]["timestamp_str"]



def compute_grade(score, total):
    """Return (percentage, letter_grade) from a raw score."""
    pct = round(score / max(total, 1) * 100)
    if   pct >= 90: grade = "A"
    elif pct >= 80: grade = "B"
    elif pct >= 70: grade = "C"
    elif pct >= 60: grade = "D"
    else:           grade = "F"
    return pct, grade



def save_text(filename, content):
    """Write a text file into OUTPUT_DIR."""
    with open(f"{OUTPUT_DIR}/{filename}", "w", encoding="utf-8") as f:
        f.write(content)



def save_json(filename, data):
    """Write a JSON file into OUTPUT_DIR."""
    with open(f"{OUTPUT_DIR}/{filename}", "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

# ══════════════════════════════════════════════════════════════
# VIDEO EXTRACTION  (OCR + Whisper)
# ══════════════════════════════════════════════════════════════



def extract_frames(video_path, interval_sec=5):
    cap   = cv2.VideoCapture(video_path)
    fps   = cap.get(cv2.CAP_PROP_FPS) or 25.0
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    dur   = total / max(fps, 1)
    # shorter lectures → denser sampling so slides are not missed
    if dur > 3600:
        interval_sec = max(interval_sec, 10)
    step  = max(1, int(fps * interval_sec))
    frames, n = [], 0
    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        if n % step == 0:
            frames.append((n / fps, frame))
        n += 1
    cap.release()
    print(f"  Sampled {len(frames)} frames | ~{total/fps:.0f}s duration")
    return frames



def classify_frame(frame):
    """Returns (type, preprocessed_frame).  Types: slide / blackboard / whiteboard / scene."""
    gray   = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
    mean_v = np.mean(gray)
    std_v  = np.std(gray)

    def sharpen(f):
        k = np.array([[0, -1, 0], [-1, 5, -1], [0, -1, 0]])
        return cv2.filter2D(f, -1, k)

    def resize960(f):
        h, w = f.shape[:2]
        return cv2.resize(
            f, (960, int(h * 960 / w)),
            interpolation=cv2.INTER_CUBIC if w < 960 else cv2.INTER_AREA,
        )

    def clahe_color(f, clip=2.8):
        lab = cv2.cvtColor(f, cv2.COLOR_BGR2LAB)
        l, a, b = cv2.split(lab)
        cl = cv2.createCLAHE(clipLimit=clip, tileGridSize=(8, 8)).apply(l)
        return cv2.cvtColor(cv2.merge([cl, a, b]), cv2.COLOR_LAB2BGR)

    # Black / dark-theme slides & OLED decks — must run BEFORE chalk invert (invert ruins navy UI)
    if mean_v < 130 and std_v >= 16:
        enh = clahe_color(frame, clip=3.6)
        k = np.array([[0, -0.5, 0], [-0.5, 3, -0.5], [0, -0.5, 0]])
        enh = cv2.filter2D(enh, -1, k)
        return 'slide', resize960(enh)

    # Uniform dark chalk wall (little edge detail — not a busy dark slide)
    if mean_v < 80 and std_v < 16:
        inv = cv2.bitwise_not(frame)
        lab = cv2.cvtColor(inv, cv2.COLOR_BGR2LAB)
        l, a, b = cv2.split(lab)
        cl = cv2.createCLAHE(clipLimit=3.5, tileGridSize=(8, 8)).apply(l)
        enh = cv2.cvtColor(cv2.merge([cl, a, b]), cv2.COLOR_LAB2BGR)
        return 'blackboard', resize960(enh)

    # Mid-brightness projected slides (beige template, lit room)
    if std_v >= 20 and 78 <= mean_v <= 195:
        enh = clahe_color(frame, clip=3.0)
        k = np.array([[0, -0.5, 0], [-0.5, 3, -0.5], [0, -0.5, 0]])
        enh = cv2.filter2D(enh, -1, k)
        return 'slide', resize960(enh)

    # Bright classic slides (white / light gray backgrounds)
    if mean_v > 152 and std_v > 16:
        g = cv2.cvtColor(sharpen(frame), cv2.COLOR_BGR2GRAY)
        _, p = cv2.threshold(g, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        return 'slide', resize960(cv2.cvtColor(p, cv2.COLOR_GRAY2BGR))

    # Physical whiteboards / light walls with writing
    if 65 <= mean_v <= 200:
        enh = clahe_color(frame, clip=2.5)
        k = np.array([[0, -0.5, 0], [-0.5, 3, -0.5], [0, -0.5, 0]])
        enh = cv2.filter2D(enh, -1, k)
        return 'whiteboard', resize960(enh)

    return 'scene', frame



def frames_are_similar(f1, f2, threshold=0.97):
    g1 = cv2.resize(cv2.cvtColor(f1, cv2.COLOR_BGR2GRAY), (160, 90))
    g2 = cv2.resize(cv2.cvtColor(f2, cv2.COLOR_BGR2GRAY), (160, 90))
    # Fast mean-absolute-diff pre-filter — avoids expensive matchTemplate for
    # frames that are obviously identical (mad < 2) or obviously different (mad > 35)
    mad = np.mean(np.abs(g1.astype(np.float32) - g2.astype(np.float32)))
    if mad < 2.0:
        return True
    if mad > 35.0:
        return False
    r = cv2.matchTemplate(
        g1.astype(np.float32), g2.astype(np.float32), cv2.TM_CCORR_NORMED,
    )
    return r[0][0] > threshold



def _ocr_frame(reader, proc, ftype):
    """Run EasyOCR with thresholds tuned per frame type."""
    try:
        kw = dict(detail=0, paragraph=True, min_size=8,
                  text_threshold=0.45, low_text=0.25, width_ths=0.9)
        if ftype == 'slide':
            # dark-theme slides: light text on dark bg needs lower thresholds
            kw.update(min_size=8, text_threshold=0.38, low_text=0.22, width_ths=0.85)
        elif ftype == 'scene':
            kw.update(min_size=6, text_threshold=0.35, low_text=0.2, width_ths=0.7)
        else:
            kw.update(min_size=12, text_threshold=0.52, low_text=0.28, width_ths=0.85)
        parts = []
        for t in reader.readtext(proc, **kw):
            s = (t[1] if isinstance(t, (list, tuple)) and len(t) > 1 else str(t)).strip()
            if len(s) > 2:
                parts.append(s)
        return " | ".join(parts)
    except Exception as ex:
        print(f"\n  OCR error: {ex}")
        return ""



BOARD_TEXT_EMPTY = "No significant board/slide text detected."



# ── OCR Drive cache helpers ──────────────────────────────────────

def _video_ocr_hash(video_path, chunk_size=65536):
    """MD5 of first 4 MB of the video — fast proxy for file identity."""
    h = hashlib.md5()
    read = 0
    with open(video_path, "rb") as f:
        while read < 4 * 1024 * 1024:
            buf = f.read(min(chunk_size, 4 * 1024 * 1024 - read))
            if not buf:
                break
            h.update(buf)
            read += len(buf)
    # also fold in file size so renamed copies are treated differently
    h.update(str(os.path.getsize(video_path)).encode())
    return h.hexdigest()

def _ocr_cache_path(video_path):
    return os.path.join(OCR_CACHE_DIR, f"{_video_ocr_hash(video_path)}.json")

def _load_ocr_cache(video_path):
    """Return (board_text, entries) from Drive cache, or None if missing/corrupt."""
    try:
        path = _ocr_cache_path(video_path)
        if not os.path.isfile(path):
            return None
        with open(path, "r", encoding="utf-8") as f:
            cached = json.load(f)
        entries    = cached.get("entries", [])
        board_text = cached.get("board_text", BOARD_TEXT_EMPTY)
        print(f"  ✅ OCR cache hit — {len(entries)} entries loaded from Drive (skipping re-scan)")
        return board_text, entries
    except Exception as e:
        print(f"  ⚠ OCR cache read failed: {e} — will re-run OCR")
        return None

def _save_ocr_cache(video_path, board_text, entries):
    """Persist OCR result to Drive for future runs."""
    try:
        path = _ocr_cache_path(video_path)
        with open(path, "w", encoding="utf-8") as f:
            json.dump({"board_text": board_text, "entries": entries},
                      f, ensure_ascii=False)
        print(f"  💾 OCR result cached → {path}")
    except Exception as e:
        print(f"  ⚠ OCR cache save failed (non-fatal): {e}")



def format_visual_lecture_text(board_text, board_entries, max_chars=14000):
    """Flatten OCR: slides (light/dark), whiteboards, chalkboards, on-screen text — one string for LLM/RAG."""
    lines = []
    for e in (board_entries or []):
        tag = (e.get("tag") or e.get("frame_type") or "Visual").strip()
        ts  = e.get("timestamp_str", "")
        tx  = (e.get("text") or "").strip()
        if len(tx) < 2:
            continue
        lines.append(f"[{ts}] [{tag}] {tx}")
    blob = "\n".join(lines)
    bt = (board_text or "").strip()
    if bt and bt != BOARD_TEXT_EMPTY and bt not in blob:
        blob = (bt + "\n" + blob).strip() if blob else bt
    elif not blob and bt and bt != BOARD_TEXT_EMPTY:
        blob = bt
    return blob[:max_chars]



def pack_video_sources(board_text, transcript, board_entries, max_total=9000):
    """Single block: all visual OCR + speech — used by summary, quiz, flashcards."""
    vis = format_visual_lecture_text(board_text, board_entries, max(2000, max_total - 2000))
    tr  = (transcript or "").strip()
    tr  = tr[: max(800, max_total - len(vis) - 80)]
    parts = []
    if vis:
        parts.append(f"=== SLIDES + BOARDS (OCR: white/black themes, on-screen) ===\n{vis}")
    if tr:
        parts.append(f"=== SPEECH (TRANSCRIPT) ===\n{tr}")
    return "\n\n".join(parts)[:max_total] if parts else tr



def load_board_entries_from_disk():
    path = os.path.join(OUTPUT_DIR, "board_entries.json")
    if not os.path.isfile(path):
        return []
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError):
        return []



def build_lecture_timeline_text(segment_list, board_entries):
    """Speech + slide/board OCR in time order (one view, like a full transcript)."""
    events = []
    for s in segment_list or []:
        t = s.get("start_sec")
        if t is None:
            continue
        ts = s.get("start_str", "")
        txt = (s.get("text") or "").strip()
        if not txt:
            continue
        events.append((float(t), f"[{ts}] (Speech) {txt}"))
    for e in board_entries or []:
        t = e.get("timestamp_sec")
        if t is None:
            continue
        tag = e.get("tag", "Visual")
        ts = e.get("timestamp_str", "")
        txt = (e.get("text") or "").strip()
        if not txt:
            continue
        events.append((float(t), f"[{ts}] [{tag}] {txt}"))
    events.sort(key=lambda x: x[0])
    return "\n\n".join(line for _, line in events) if events else ""



def rebuild_lecture_timeline_from_disk():
    sp = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    bp = os.path.join(OUTPUT_DIR, "board_entries.json")
    segs, boards = [], []
    try:
        if os.path.isfile(sp):
            with open(sp, "r", encoding="utf-8") as f:
                segs = json.load(f)
    except (json.JSONDecodeError, OSError):
        pass
    try:
        if os.path.isfile(bp):
            with open(bp, "r", encoding="utf-8") as f:
                boards = json.load(f)
    except (json.JSONDecodeError, OSError):
        pass
    return build_lecture_timeline_text(segs, boards)



def extract_board_text(video_path):
    """Run EasyOCR on sampled frames.  Returns (board_text_str, entries_list).
    Results are cached on Google Drive by video hash — re-uploads of the same
    file skip the expensive OCR pass entirely."""
    cached = _load_ocr_cache(video_path)
    if cached is not None:
        board_text, entries = cached
        save_text("board_text.txt",
                  f"=== BOARD & SLIDE TEXT ===\n(from cache)\n\n{board_text}")
        save_json("board_entries.json", entries)
        return board_text, entries

    reader    = get_ocr()
    frames    = extract_frames(video_path)
    entries   = []
    type_log  = {t: 0 for t in ('blackboard', 'whiteboard', 'slide', 'scene')}
    prev_text = ""
    prev_frm  = None
    skipped   = 0

    LABEL_MAP = {'blackboard': 'Board', 'whiteboard': 'Board', 'slide': 'Slide', 'scene': 'On-screen'}

    def process_frame(ts, frame):
        nonlocal prev_text, entries, type_log
        ftype, proc = classify_frame(frame)
        type_log[ftype] = type_log.get(ftype, 0) + 1

        h, w = frame.shape[:2]
        rw = max(w, 1)
        resized = cv2.resize(frame, (960, int(h * 960 / rw)))
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        eq = cv2.cvtColor(cv2.equalizeHist(gray), cv2.COLOR_GRAY2BGR)
        h2, w2 = eq.shape[:2]
        eq_r = cv2.resize(eq, (960, int(h2 * 960 / max(w2, 1))))

        candidates = [(ftype, proc)]
        if ftype == 'scene':
            candidates.append(('scene', resized))
            candidates.append(('scene', eq_r))
        else:
            candidates.append(('slide', resized))
            candidates.append(('slide', eq_r))

        text = ""
        for c_ft, c_im in candidates:
            t = _ocr_frame(reader, c_im, c_ft)
            if t:
                text = t
                break
        if not text:
            return
        if difflib.SequenceMatcher(None, prev_text, text).ratio() < 0.78:
            ts_str = fmt_ts(ts)
            entries.append({
                "timestamp_sec": round(ts, 1),
                "timestamp_str": ts_str,
                "frame_type": ftype,
                "tag": LABEL_MAP.get(ftype, 'On-screen'),
                "text": text,
            })
            prev_text = text

    for i, (ts, frame) in enumerate(frames):
        print(f"  OCR frame {i+1}/{len(frames)} @{ts:.0f}s  skipped={skipped}", end='\r')

        if prev_frm is not None and frames_are_similar(prev_frm, frame, threshold=0.97):
            skipped += 1
            prev_frm = frame
            continue
        prev_frm = frame
        process_frame(ts, frame)

    # second pass: denser sampling if almost nothing found (talking-head + small text)
    if len(entries) < 2 and frames:
        print("\n  OCR sparse — second pass every 3s including scene-like frames…")
        frames2 = extract_frames(video_path, interval_sec=3)
        prev_frm = None
        for ts, frame in frames2:
            if prev_frm is not None and frames_are_similar(prev_frm, frame, threshold=0.97):
                prev_frm = frame
                continue
            prev_frm = frame
            process_frame(ts, frame)

    board_text = "\n".join(
        f"[{e['timestamp_str']}] {e['tag']}  {e['text']}" for e in entries
    ) or BOARD_TEXT_EMPTY

    save_text("board_text.txt",
              f"=== BOARD & SLIDE TEXT ===\nFrame types: {type_log}\n\n{board_text}")
    save_json("board_entries.json", entries)

    print(f"\n  OCR done — {len(entries)} unique states | {type_log}")
    _save_ocr_cache(video_path, board_text, entries)
    return board_text, entries



def transcribe_video(video_path):
    """Transcribe audio with Faster-Whisper.  Returns (transcript_str, segment_list)."""
    audio_tmp = f"{OUTPUT_DIR}/_audio.wav"
    print("  Extracting audio…")
    subprocess.run([
        'ffmpeg', '-y', '-i', video_path,
        '-vn', '-acodec', 'pcm_s16le', '-ar', '16000', '-ac', '1', audio_tmp,
    ], capture_output=True, check=False)

    if not os.path.exists(audio_tmp):
        return "Audio extraction failed.", []

    model = get_whisper()
    print("  Transcribing…")
    gen, info = model.transcribe(
        audio_tmp, beam_size=5, language=None,
        vad_filter=True,
        vad_parameters={"min_silence_duration_ms": 500},
        word_timestamps=False,
    )
    print(f"  Detected: {info.language} (p={info.language_probability:.2f})")

    segs, lines = [], []
    for seg in gen:
        ts = fmt_ts(seg.start)
        segs.append({
            "start_sec": round(seg.start, 1),
            "end_sec":   round(seg.end,   1),
            "start_str": ts,
            "end_str":   fmt_ts(seg.end),
            "text":      seg.text.strip(),
        })
        lines.append(f"[{ts}]  {seg.text.strip()}")

    transcript = "\n".join(lines) or "No speech detected."
    try:
        os.remove(audio_tmp)
    except OSError:
        pass

    save_text("transcript.txt", "=== LECTURE TRANSCRIPT ===\n\n" + transcript)
    save_json("transcript_segments.json", segs)

    print("  Transcript saved.")
    return transcript, segs



print("✅ Video extraction ready.")

# ══════════════════════════════════════════════════════════════
# DOCUMENT EXTRACTION  (PDF / DOCX / PPTX)
# ══════════════════════════════════════════════════════════════



def extract_pdf(path):
    import fitz
    doc = fitz.open(path)
    entries = []
    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            entries.append({"page": i + 1, "text": text, "source": "pdf"})
    doc.close()
    return entries



def extract_docx(path):
    from docx import Document
    doc = Document(path)
    entries, buf, sec = [], [], 1
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        if 'heading' in para.style.name.lower() and buf:
            entries.append({"section": sec, "text": " ".join(buf), "source": "docx"})
            buf = []
            sec += 1
        buf.append(t)
    if buf:
        entries.append({"section": sec, "text": " ".join(buf), "source": "docx"})
    return entries



def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    entries = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            ph = shape.placeholder_format
            if ph and ph.idx == 0:                       # idx 0 = title placeholder
                parts.insert(0, shape.text_frame.text.strip())
            else:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        if parts:
            entries.append({"slide": i + 1, "text": "\n".join(parts), "source": "pptx"})
    return entries



DOC_EXTRACTORS = {
    'pdf':  extract_pdf,
    'docx': extract_docx,
    'pptx': extract_pptx,
    'ppt':  extract_pptx,
}



print("✅ Document extraction ready.")

# ══════════════════════════════════════════════════════════════
# RAG ENGINE  (FAISS-based retrieval)
# ══════════════════════════════════════════════════════════════



class LectureRAG:
    """Encodes text chunks with SentenceTransformer, indexes with FAISS,
    and retrieves top-k relevant chunks for a query."""

    def __init__(self):
        self.embedder   = None
        self.index      = None
        self.chunks     = []
        self.timestamps = []

    def build(self, entries, text_key="text", ts_key="timestamp_str", source_key="tag"):
        self.embedder   = get_embedder()
        self.chunks     = [e[text_key] for e in entries]
        self.timestamps = [
            {"timestamp": e.get(ts_key, ""), "source": e.get(source_key, "")}
            for e in entries
        ]
        vecs = self.embedder.encode(
            self.chunks, batch_size=64, show_progress_bar=False,
        ).astype('float32')
        self.index = faiss.IndexFlatIP(vecs.shape[1])
        faiss.normalize_L2(vecs)
        self.index.add(vecs)
        print(f"  RAG index built — {len(self.chunks)} chunks.")

    def query(self, question, top_k=5):
        if not self.index:
            return "No index built yet.", []
        vec = self.embedder.encode([question]).astype('float32')
        faiss.normalize_L2(vec)
        _, idxs = self.index.search(vec, top_k)
        ctx = "\n\n".join(self.chunks[i] for i in idxs[0] if i < len(self.chunks))
        tss = [self.timestamps[i] for i in idxs[0] if i < len(self.timestamps)]
        return ctx, tss



rag = LectureRAG()



def build_rag(entries, text_key="text", ts_key="timestamp_str", src_key="tag"):
    """Rebuild the global RAG index from a new set of entries."""
    global rag
    rag = LectureRAG()
    rag.build(entries, text_key=text_key, ts_key=ts_key, source_key=src_key)



print("✅ RAG engine ready.")

# ══════════════════════════════════════════════════════════════
# ANALYSIS  (Summary + Exam Hints + Topic Ranking)
# ══════════════════════════════════════════════════════════════



def generate_summary(board_text, transcript, segment_list, board_entries):
    """Merge OCR (slides + boards, all themes) + speech into one summary."""
    visual = format_visual_lecture_text(board_text, board_entries, 10000)
    merged = (
        "Instructions: Combine everything below. OCR lines are from projector slides, "
        "digital whiteboards, chalkboards, and on-screen text (any color theme).\n\n"
        f"=== SLIDES + BOARDS (OCR) ===\n{visual}\n\n"
        f"=== SPEECH TRANSCRIPT ===\n{transcript[:10000]}"
    )
    return call_ollama(
        f"Create a comprehensive, structured lecture summary:\n\n{merged}",
        system=(
            "You are an expert educational content summarizer. "
            "Treat slide/board OCR and spoken words as one lecture — cite both where useful."
        ),
    )



EMPHASIS_SIGNALS = [
    "this is important", "remember", "this will be on", "exam", "key",
    "critical", "note", "highlight", "must know", "definition",
]



def detect_exam_hints(segment_list, board_entries, chunk_size=30):
    """Two-pass LLM analysis: find emphasis in transcript + key phrases on board."""
    all_hints = []

    def has_signal(seg):
        lower = seg.get("text", "").lower()
        return any(sig in lower for sig in EMPHASIS_SIGNALS)

    # Pass 1 — transcript emphasis moments
    for i in range(0, len(segment_list), chunk_size):
        chunk    = segment_list[i:i + chunk_size]
        filtered = [s for s in chunk if has_signal(s)]
        if not filtered:
            continue
        text = "\n".join(f"[{s['start_str']}] {s['text']}" for s in filtered)
        prompt = (
            "Identify sentences where the teacher explicitly signals exam importance. "
            "Return JSON list: [{\"text\",\"timestamp\",\"emphasis_type\",\"signal_phrase\"}]\n\n"
            + text
        )
        result = call_ollama_json(prompt, fallback=[], model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
        if isinstance(result, list):
            all_hints.extend(result)

    # Pass 2 — board/slide key phrases
    if board_entries:
        board_text = "\n".join(
            f"[{e['timestamp_str']}] {e['text'][:200]}" for e in board_entries[:40]
        )
        prompt = (
            "Extract key terms and formulas from these board/slide notes. "
            "Return JSON list: [{\"text\",\"timestamp\",\"emphasis_type\",\"signal_phrase\"}]\n\n"
            + board_text
        )
        result = call_ollama_json(prompt, fallback=[], model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
        if isinstance(result, list):
            all_hints.extend(result)

    return all_hints



def analyze_topics(summary, hints):
    """Ask LLM to rank topics by exam likelihood."""
    hint_text = "\n".join(h.get("text", "")[:200] for h in hints[:20])
    prompt = (
        "Based on the summary and teacher emphasis moments below, identify the 5-8 "
        "most exam-likely topics. "
        "Return JSON: {\"ai_important_topics\":[{\"topic\",\"confidence\",\"reason\","
        "\"evidence\",\"suggested_note\"}],\"overall_exam_focus\":\"string\"}\n\n"
        f"SUMMARY:\n{summary[:3000]}\n\nHINTS:\n{hint_text}"
    )
    return call_ollama_json(prompt, fallback={}, model=QUIZ_MODEL, num_ctx=QUIZ_CTX)



print("✅ Analysis functions ready.")

# ══════════════════════════════════════════════════════════════
# QUIZ ENGINE  —  Generators
# ══════════════════════════════════════════════════════════════



# ── Shared JSON schemas for all quiz types ──────────────────────
QUIZ_SCHEMAS = {
    "mcq":   ('{"questions":[{"question":"","options":{"A":"","B":"","C":"","D":""},'
               '"correct_answer":"A","explanation":"","topic":"","difficulty":"",'
               '"bloom_level":"","source_timestamp":""}]}'),
    "tf":    ('{"questions":[{"statement":"","answer":true,"explanation":"",'
               '"topic":"","difficulty":"","source_timestamp":""}]}'
               '  NOTE: answer must be JSON boolean true or false, never a string.'),
    "fill":  ('{"questions":[{"question":"The ___ is responsible for X.","answer":"term","hint":"starts with T",'
               '"topic":"","difficulty":"","source_timestamp":""}]}'),
    "short": ('{"questions":[{"question":"","model_answer":"","key_points":[],'
               '"marks":4,"topic":"","difficulty":"","source_timestamp":""}]}'),
}


def generate_mcq(board_text, transcript, summary, num_questions=12,
                 doc_text="", board_entries=None, segment_list=None,
                 difficulty_hint=""):
    board_entries = board_entries or []
    segment_list  = segment_list or []
    if doc_text:
        content = doc_text[:8000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 8000)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate {num_questions} high-quality MCQ from this lecture content.\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"question\",\"options\":{\"A\",\"B\",\"C\",\"D\"},"
        "\"correct_answer\",\"explanation\",\"topic\",\"difficulty\",\"bloom_level\","
        "\"source_timestamp\",\"source_type\"}]}\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, segments=segment_list)
    return qs



def generate_tf(transcript, summary, doc_text="", segment_list=None,
                difficulty_hint="", board_text="", board_entries=None):
    segment_list = segment_list or []
    if doc_text:
        content = doc_text[:6000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 7500)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate 10 True/False questions.\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"statement\",\"answer\",\"explanation\","
        "\"topic\",\"difficulty\",\"source_timestamp\"}]}\n"
        "RULES: 'answer' MUST be a JSON boolean (true or false) — never a string. "
        "Mix true and false answers roughly equally.\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, segments=segment_list)
    return qs



def generate_fill(transcript, summary, doc_text="", board_entries=None,
                  difficulty_hint="", board_text=""):
    board_entries = board_entries or []
    if doc_text:
        content = doc_text[:6000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 7500)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate 10 fill-in-the-blank questions.\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"question\",\"answer\",\"hint\","
        "\"topic\",\"difficulty\",\"source_timestamp\"}]}\n"
        "RULES: Each 'question' MUST contain exactly '___' (three underscores) "
        "to mark the blank. 'answer' is the word/phrase that fills the blank. "
        "'hint' gives a clue without revealing the answer.\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, board_entries=board_entries)
    return qs



def generate_short(transcript, summary, doc_text="", segment_list=None,
                   difficulty_hint="", board_text="", board_entries=None):
    segment_list = segment_list or []
    if doc_text:
        content = doc_text[:6000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 7500)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate 8 short-answer questions (4-6 marks each).\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"question\",\"model_answer\",\"key_points\","
        "\"marks\",\"topic\",\"difficulty\",\"source_timestamp\"}]}\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, segments=segment_list)
    return qs



def generate_flashcards(summary, doc_text="", board_text="", board_entries=None):
    if doc_text:
        content = (doc_text + "\n\n" + summary[:2500])[:6000]
    else:
        vis = format_visual_lecture_text(board_text, board_entries or [], 4000)
        parts = []
        if vis:
            parts.append(f"=== SLIDES + BOARDS (OCR) ===\n{vis}")
        parts.append(f"=== SUMMARY ===\n{summary[:3500]}")
        content = "\n\n".join(parts)[:6500]
    prompt = (
        "Create 15 flashcards for key terms/concepts from slides, boards, and summary. "
        "STRICT JSON: {\"cards\":[{\"front\",\"back\",\"topic\",\"difficulty\"}]}\n\n"
        f"CONTENT:\n{content}"
    )
    data = call_ollama_json_quiz(prompt, key="cards")
    return data.get("cards", []) if isinstance(data, dict) else []

# ══════════════════════════════════════════════════════════════
# QUIZ ENGINE  —  Graders
# ══════════════════════════════════════════════════════════════



def grade_mcq(questions, answers):
    results, score = [], 0
    for i, q in enumerate(questions):
        user    = str(answers.get(str(i), "")).strip().upper()
        correct = str(q.get("correct_answer", "")).strip().upper()
        ok = user == correct
        if ok:
            score += 1
        results.append({**q, "user_ans": user, "correct_ans": correct, "correct": ok})

    pct, grade = compute_grade(score, len(questions))
    return {"results": results, "score": score, "total": len(questions),
            "percentage": pct, "grade": grade}



def _norm_tf(v):
    """Normalise any truthy/falsy representation to the strings 'true' or 'false'."""
    if isinstance(v, bool):
        return "true" if v else "false"
    s = str(v).strip().lower()
    if s in ("1", "yes", "true",  "t"):  return "true"
    if s in ("0", "no",  "false", "f"):  return "false"
    return s  # pass-through so exact mismatch is still caught



def grade_tf(questions, answers):
    results, score = [], 0
    for i, q in enumerate(questions):
        user    = _norm_tf(answers.get(str(i), ""))
        correct = _norm_tf(q.get("answer", False))
        ok = user == correct
        if ok:
            score += 1
        results.append({**q, "user_ans": user, "correct_ans": correct, "correct": ok})

    pct, grade = compute_grade(score, len(questions))
    return {"results": results, "score": score, "total": len(questions),
            "percentage": pct, "grade": grade}



def grade_fill(questions, answers):
    results, score = [], 0
    for i, q in enumerate(questions):
        user    = str(answers.get(str(i), "")).strip().lower()
        correct = str(q.get("answer", "")).strip().lower()

        # Substring check only when correct answer is long enough to be meaningful
        # (avoids "a" matching inside any word, "DNA" matching "banana", etc.)
        substring_ok = len(correct) >= 4 and correct in user

        # Word-boundary check: correct answer appears as a distinct word/phrase in user text
        user_words   = re.split(r'\W+', user)
        word_ok      = correct in user_words

        # Fuzzy similarity as the final fallback
        fuzzy_ok     = difflib.SequenceMatcher(None, user, correct).ratio() > 0.80

        ok = word_ok or substring_ok or fuzzy_ok
        if ok:
            score += 1
        results.append({**q, "user_ans": user, "correct_ans": correct, "correct": ok})

    pct, grade = compute_grade(score, len(questions))
    return {"results": results, "score": score, "total": len(questions),
            "percentage": pct, "grade": grade}



def grade_short(questions, answers):
    """AI-powered short-answer grading — sends each Q+A pair to the LLM."""
    results, total_score, total_marks = [], 0, 0
    for i, q in enumerate(questions):
        user_ans = str(answers.get(str(i), "")).strip()
        marks    = int(q.get("marks", 4))
        total_marks += marks

        if not user_ans:
            results.append({**q, "user_ans": "", "score": 0, "max_marks": marks,
                            "correct": False, "correct_ans": q.get("model_answer", ""),
                            "feedback": "No answer provided.", "key_points_hit": []})
            continue

        prompt = (
            f"Grade this short answer out of {marks} marks.\n"
            f"Question: {q.get('question', '')}\n"
            f"Model answer: {q.get('model_answer', '')}\n"
            f"Key points expected: {json.dumps(q.get('key_points', []))}\n"
            f"Student answer: {user_ans}\n\n"
            f"Return JSON: {{\"score\": 0-{marks}, \"feedback\": \"...\", "
            f"\"key_points_hit\": [\"...\"]}}. Be fair but strict."
        )
        graded = call_ollama_json(prompt, fallback={"score": 0, "feedback": "Grading failed."})
        sc = min(int(graded.get("score", 0)), marks)
        total_score += sc
        results.append({
            **q, "user_ans": user_ans, "score": sc, "max_marks": marks,
            "correct": sc >= marks * 0.5,
            "correct_ans": q.get("model_answer", ""),
            "feedback": graded.get("feedback", ""),
            "key_points_hit": graded.get("key_points_hit", []),
        })

    pct, grade = compute_grade(total_score, total_marks)
    return {"results": results, "score": total_score, "total": total_marks,
            "percentage": pct, "grade": grade}



def explain_wrong(wrong_questions):
    """Generate AI explanations for any list of wrong questions."""
    if not wrong_questions:
        return {}
    text = "\n".join(
        f"Q: {q.get('question', q.get('statement', ''))}  "
        f"Correct: {q.get('correct_answer', q.get('correct_ans', q.get('answer', '')))}"
        for q in wrong_questions[:8]
    )
    prompt = (
        "For each wrong question below, explain the underlying concept simply "
        "so the student can learn from their mistake.\n"
        "Return JSON: {\"explanations\":{\"<question_text>\":\"<explanation>\"}}\n\n" + text
    )
    result = call_ollama_json(prompt, fallback={"explanations": {}},
                                model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
    return result.get("explanations", {}) if isinstance(result, dict) else {}



print("✅ Quiz engine ready.")

# ══════════════════════════════════════════════════════════════
# AGENT FRAMEWORK
# ══════════════════════════════════════════════════════════════



class AgentContext:
    """Shared memory for a single pipeline run.  Every agent reads/writes
    here instead of passing dozens of parameters between functions."""

    def __init__(self, content_type, file_path):
        self.content_type  = content_type        # "video" | "document"
        self.file_path     = file_path
        self.source_name   = Path(file_path).name

        # populated by extraction agents
        self.board_text    = ""
        self.board_entries = []
        self.transcript    = ""
        self.segment_list  = []
        self.doc_text      = ""
        self.doc_entries   = []

        # populated by analysis / quiz agents
        self.summary       = ""
        self.quiz_counts   = {"mcq": 12, "tf": 10, "fill": 10, "short": 8}
        self.difficulty_pref = None              # set by DifficultyAdapter

        # agent decision log
        self.agent_log     = []



class BaseAgent:
    """Every specialist agent inherits from this.
    Provides the agentic loop: should_run → plan → execute → validate → retry."""

    name        = "BaseAgent"
    max_retries = 2

    def should_run(self, ctx):
        """Return False to skip this agent entirely."""
        return True

    def plan(self, ctx):
        """Return a short string describing what this agent will do."""
        return "execute default task"

    def execute(self, ctx):
        """Do the actual work.  Must be implemented by subclasses."""
        raise NotImplementedError

    def validate(self, ctx):
        """Check output quality.  Return (ok: bool, issues: list[str])."""
        return True, []

    def run(self, ctx):
        """Full agent loop with retry on validation failure."""
        if not self.should_run(ctx):
            self._log(ctx, "SKIP", "Not needed for this content type")
            return

        self._log(ctx, "PLAN", self.plan(ctx))

        last_issues = []
        for attempt in range(1, self.max_retries + 1):
            try:
                self.execute(ctx)
                ok, issues = self.validate(ctx)
                if ok:
                    self._log(ctx, "DONE", f"Completed on attempt {attempt}")
                    return
                last_issues = issues
                self._log(ctx, "RETRY", f"Quality issues (attempt {attempt}): {issues}")
            except Exception as e:
                self._log(ctx, "ERROR", f"Attempt {attempt} failed: {e}")
                if attempt == self.max_retries:
                    self._log(ctx, "FAIL", f"Gave up after {self.max_retries} attempts: {e}")
                    return

        self._log(ctx, "DONE", f"Completed with warnings: {last_issues}")

    def _log(self, ctx, level, msg):
        entry = {
            "agent": self.name, "level": level,
            "time": datetime.now().isoformat(), "message": msg,
        }
        ctx.agent_log.append(entry)
        icons = {"SKIP": "⏭", "PLAN": "📋", "DONE": "✅",
                 "RETRY": "🔄", "ERROR": "❌", "FAIL": "💀"}
        print(f"  [{self.name}] {icons.get(level, '')} {msg}")



print("✅ Agent framework ready.")

# ══════════════════════════════════════════════════════════════
# SPECIALIST AGENTS
# ══════════════════════════════════════════════════════════════



class ExtractorAgent(BaseAgent):
    """Runs EasyOCR on video frames to extract board/slide text."""
    name = "Extractor"

    def should_run(self, ctx):
        return ctx.content_type == "video"

    def plan(self, ctx):
        return f"OCR scan frames from {ctx.source_name}"

    def execute(self, ctx):
        ctx.board_text, ctx.board_entries = extract_board_text(ctx.file_path)

    def validate(self, ctx):
        if not ctx.board_entries:
            return True, ["no board text found (may be camera-only video)"]
        return True, []



class TranscriberAgent(BaseAgent):
    """Transcribes audio from video using Faster-Whisper."""
    name = "Transcriber"

    def should_run(self, ctx):
        return ctx.content_type == "video"

    def plan(self, ctx):
        return f"Whisper transcription of {ctx.source_name}"

    def execute(self, ctx):
        global _whisper_model
        try:
            ctx.transcript, ctx.segment_list = transcribe_video(ctx.file_path)
        except Exception as e:
            err = str(e).lower()
            if "out of memory" in err or "cuda" in err:
                self._log(ctx, "RETRY", f"GPU OOM — resetting Whisper to try smaller model")
                _whisper_model = None  # force re-init with fallback cascade
                ctx.transcript, ctx.segment_list = transcribe_video(ctx.file_path)
            else:
                raise

    def validate(self, ctx):
        if not ctx.transcript or ctx.transcript == "No speech detected.":
            return False, ["empty transcript — audio extraction may have failed"]
        if len(ctx.transcript) < 50:
            return False, [f"transcript too short ({len(ctx.transcript)} chars)"]
        return True, []



class DocParserAgent(BaseAgent):
    """Extracts text from PDF / DOCX / PPTX documents."""
    name = "DocParser"

    def should_run(self, ctx):
        return ctx.content_type == "document"

    def plan(self, ctx):
        ext = Path(ctx.file_path).suffix.lower().lstrip('.')
        return f"Parse {ext.upper()} document: {ctx.source_name}"

    def execute(self, ctx):
        ext = Path(ctx.file_path).suffix.lower().lstrip('.')
        extractor = DOC_EXTRACTORS.get(ext)
        if not extractor:
            ctx.doc_entries = []
            ctx.doc_text = ""
            return
        try:
            ctx.doc_entries = extractor(ctx.file_path)
        except Exception as e:
            self._log(ctx, "ERROR", f"Failed to parse {ext.upper()}: {e}")
            ctx.doc_entries = []
        ctx.doc_text = "\n\n".join(e['text'] for e in ctx.doc_entries)

    def validate(self, ctx):
        if not ctx.doc_entries:
            return False, ["no text extracted — file may be corrupt or password-protected"]
        if len(ctx.doc_text) < 100:
            return False, [f"extracted text too short ({len(ctx.doc_text)} chars)"]
        return True, []



class SummaryAgent(BaseAgent):
    """Generates a structured summary using the LLM.  Validates length and quality."""
    name = "Summarizer"

    def plan(self, ctx):
        src = "board text + transcript" if ctx.content_type == "video" else "document text"
        return f"Generate structured summary from {src}"

    def execute(self, ctx):
        if ctx.content_type == "video":
            ctx.summary = generate_summary(
                ctx.board_text, ctx.transcript, ctx.segment_list, ctx.board_entries,
            )
        else:
            ctx.summary = call_ollama(
                f"Create a detailed, structured summary of this document:\n\n"
                f"{ctx.doc_text[:12000]}",
                system="You are an expert educational content summarizer.",
            )
        save_text("summary.txt", ctx.summary)

    def validate(self, ctx):
        issues = []
        if len(ctx.summary) < 200:
            issues.append(f"summary too short ({len(ctx.summary)} chars)")
        low_quality = ["i don't have enough", "i cannot", "no content provided"]
        if any(phrase in ctx.summary.lower() for phrase in low_quality):
            issues.append("summary appears to be a refusal or low-quality response")
        return len(issues) == 0, issues



class RAGBuilderAgent(BaseAgent):
    """Builds the FAISS search index from extracted content."""
    name = "RAGBuilder"

    def plan(self, ctx):
        return "Build FAISS vector index for RAG retrieval"

    def execute(self, ctx):
        if ctx.content_type == "video":
            combined = list(ctx.segment_list) + [
                {
                    "timestamp_str": e["timestamp_str"],
                    "text": f"[{e.get('tag', 'Visual')}] {e['text']}",
                    "start_str": e["timestamp_str"],
                    "tag": e.get("tag", "Visual"),
                }
                for e in (ctx.board_entries or [])
            ]
            build_rag(combined, ts_key="start_str", src_key="tag")
        else:
            build_rag(ctx.doc_entries, text_key="text")



class HintsAgent(BaseAgent):
    """Detects teacher emphasis moments and ranks topics by exam likelihood."""
    name = "HintsDetector"

    def should_run(self, ctx):
        return bool(ctx.summary)

    def plan(self, ctx):
        return f"Analyze {len(ctx.segment_list)} segments + summary for exam emphasis signals"

    def execute(self, ctx):
        hints     = detect_exam_hints(ctx.segment_list or [], ctx.board_entries or [])
        ai_topics = analyze_topics(ctx.summary, hints)
        state.quiz["hints"] = {
            "teacher_hints": hints,
            "ai_analysis": ai_topics,
            "total_hint_sentences": len(hints),
        }



class QuizAgent(BaseAgent):
    """Runs the per-question agentic pipeline (Plan → Generate ONE → Validate
    → Retry → Next) for every quiz type during initial processing.
    Streams thinking events to state._quiz_agent_thinking so the UI can
    render the live Agent Workspace while the upload is being processed."""
    name        = "QuizGenerator"
    max_retries = 1

    # Quiz types we generate during processing, in execution order
    _types_in_order = ("mcq", "tf", "fill", "short")

    def plan(self, ctx):
        c = ctx.quiz_counts
        return (f"Agentic pipeline for {c['mcq']} MCQ, {c['tf']} T/F, "
                f"{c['fill']} Fill-blank, {c['short']} Short-answer "
                f"(per-question validate + retry)")

    def execute(self, ctx):
        c       = ctx.quiz_counts
        target  = "medium"
        # one content snippet shared across all quiz types
        if ctx.content_type == "video":
            content = pack_video_sources(ctx.board_text, ctx.transcript,
                                         ctx.board_entries or [], 7500)
        else:
            content = (ctx.doc_text or "")[:7500]
        summary = ctx.summary or state.status.get("summary", "")

        # start a fresh thinking session for the UI
        state._quiz_thinking_session += 1
        state._quiz_thinking_active   = True
        state._quiz_agent_thinking    = []
        log = state._quiz_agent_thinking

        def push(entry):
            log.append(entry)

        # propagate ctx fields used by _run_quiz_agentic_loop
        ctx._user_instructions = ""
        ctx._mode              = "auto"
        ctx._agent_name        = self.name

        push({
            "agent": "Orchestrator",
            "step":  "auto_quiz_start",
            "phase": "PLAN",
            "verdict": "PASS",
            "reason": (
                f"Auto-quiz starting: {c['mcq']} MCQ · {c['tf']} TF · "
                f"{c['fill']} Fill · {c['short']} Short"
            ),
            "data": {"counts": c, "target_difficulty": target},
        })

        for qt in self._types_in_order:
            num = c.get(qt, 0)
            if num <= 0:
                continue
            ctx._quiz_count = num
            ctx._quiz_type  = qt
            ctx._quiz_plan  = None  # let the planner re-plan per type

            push({
                "agent":  self.name,
                "step":   "type_start",
                "phase":  "PLAN",
                "verdict":"PASS",
                "reason": f"Beginning agentic generation for {num} {qt.upper()} question(s)…",
                "data":   {"quiz_type": qt, "count": num},
            })

            try:
                result = _run_quiz_agentic_loop(
                    ctx, qt, content, summary, "",
                    target_difficulty=target,
                    thinking_cb=push,
                    max_attempts_per_q=3,
                )
                fill_missing_timestamps(
                    result,
                    segments=ctx.segment_list,
                    board_entries=ctx.board_entries,
                )
                state.quiz[qt] = result
                _append_generated_quiz_run("upload", qt, "auto", target, result)
                push({
                    "agent":  self.name,
                    "step":   "type_done",
                    "phase":  "DONE",
                    "verdict":"DONE",
                    "reason": f"{qt.upper()} ready: {len(result)} question(s)",
                    "data":   {"quiz_type": qt, "total": len(result)},
                })
            except Exception as e:
                push({
                    "agent":  self.name,
                    "step":   "type_error",
                    "phase":  "DONE",
                    "verdict":"FAIL",
                    "reason": f"{qt.upper()} crashed: {e} — keeping previous results",
                })
                state.quiz.setdefault(qt, [])

        push({
            "agent":  "Orchestrator",
            "step":   "auto_quiz_complete",
            "phase":  "DONE",
            "verdict":"DONE",
            "reason": (
                f"All quiz types ready: "
                f"MCQ={len(state.quiz.get('mcq', []))} · "
                f"TF={len(state.quiz.get('tf', []))} · "
                f"Fill={len(state.quiz.get('fill', []))} · "
                f"Short={len(state.quiz.get('short', []))}"
            ),
        })
        state._quiz_thinking_active = False

    def validate(self, ctx):
        issues = []
        mcqs = state.quiz.get("mcq", [])
        if not mcqs:
            issues.append("MCQ pipeline produced 0 questions")
            if state.quiz.get("tf") or state.quiz.get("fill") or state.quiz.get("short"):
                return True, issues
            return False, issues

        # The agentic loop already enforces uniqueness; this is a safety net.
        texts = [q.get("question", "") for q in mcqs]
        dupes = len(texts) - len(set(texts))
        if dupes > 0:
            seen, unique = set(), []
            for q in mcqs:
                t = q.get("question", "")
                if t not in seen:
                    seen.add(t); unique.append(q)
            state.quiz["mcq"] = unique
            issues.append(f"removed {dupes} duplicate MCQ(s)")

        bad_keys = 0
        for q in mcqs:
            opts = q.get("options", {})
            ans  = q.get("correct_answer", "")
            if opts and ans not in opts:
                bad_keys += 1
        if bad_keys > len(mcqs) // 2:
            issues.append(f"{bad_keys}/{len(mcqs)} MCQs have broken answer keys")
            return False, issues
        elif bad_keys > 0:
            issues.append(f"{bad_keys} MCQ(s) have mismatched answer key — non-critical")

        tfs = state.quiz.get("tf", [])
        if tfs:
            seen, unique = set(), []
            for q in tfs:
                t = q.get("statement", "")
                if t not in seen:
                    seen.add(t); unique.append(q)
            state.quiz["tf"] = unique

        return True, issues

# ══════════════════════════════════════════════════════════════
# TRUE AGENTIC QUIZ SYSTEM
# Planner → Generator → Validator → Refiner loop
# All agents push live thinking entries via thinking_cb.
# ══════════════════════════════════════════════════════════════


# ── Quiz Planner Agent ──────────────────────────────────────────

class QuizPlannerAgent(BaseAgent):
    """Reads user instructions + student performance + content summary
    and produces a structured plan (topics, counts, difficulties)
    that downstream Generator/Validator/Refiner agents follow."""
    name = "QuizPlannerAgent"
    max_retries = 2

    def plan(self, ctx):
        return "Analyzing instructions and performance to build a quiz plan"

    def execute(self, ctx):
        quiz_type  = getattr(ctx, '_quiz_type', 'mcq')
        num        = getattr(ctx, '_quiz_count', 10)
        mode       = getattr(ctx, '_mode', 'adaptive')
        difficulty = getattr(ctx, '_manual_difficulty', 'medium')
        user_instr = getattr(ctx, '_user_instructions', '')
        plan_mode  = getattr(ctx, '_plan_mode', 'auto')
        manual_topics = [t.strip() for t in (getattr(ctx, '_manual_topics', []) or []) if str(t).strip()]
        perf       = difficulty_adapter.get_all_stats()
        weak       = difficulty_adapter.get_weak_topics()

        # Manual topic mode: bypass LLM planner and force the selected topics.
        if plan_mode == "manual_topics" and manual_topics:
            topics = manual_topics[:8]
            base   = num // len(topics)
            rem    = num - base * len(topics)
            counts = {t: base for t in topics}
            counts[topics[0]] += rem
            diffs  = {t: (difficulty if mode == "manual" else "medium") for t in topics}
            ctx._quiz_plan = {
                "topics": topics,
                "per_topic_count": counts,
                "per_topic_difficulty": diffs,
                "bloom_targets": ["Remember", "Understand", "Apply"],
                "focus_notes": "Manual topic plan selected by user",
            }
            return

        perf_summary = ""
        if perf:
            lines = [
                f"  - {t}: {int(s['average']*100)}% ({s['level']})"
                for t, s in list(perf.items())[:8]
            ]
            perf_summary = "Student performance history:\n" + "\n".join(lines)

        summary_snip = (ctx.summary or state.status.get("summary", ""))[:2000]

        mode_note = (
            f"Target difficulty: {difficulty.upper()} (uniform, manual mode)"
            if mode == "manual"
            else "Mode: ADAPTIVE — assign lower difficulty to weak topics, higher to strong ones"
        )

        prompt = (
            f"You are a quiz planning agent. Decide how to distribute {num} "
            f"{quiz_type.upper()} questions across topics for a student.\n\n"
            f"Quiz type: {quiz_type.upper()}\n"
            f"{mode_note}\n"
            f"Total questions: {num}\n"
            f"Student instructions: {user_instr or 'None'}\n\n"
            f"{perf_summary}\n\n"
            f"Lecture summary (excerpt):\n{summary_snip}\n\n"
            "Return STRICT JSON (no extra text):\n"
            '{"topics":["topic1","topic2"],'
            '"per_topic_count":{"topic1":5,"topic2":5},'
            '"per_topic_difficulty":{"topic1":"easy","topic2":"medium"},'
            '"bloom_targets":["Remember","Apply"],'
            '"focus_notes":"what to emphasize or avoid based on instructions"}\n\n'
            "Rules:\n"
            f"- Sum of per_topic_count MUST equal {num}\n"
            "- ADAPTIVE: weak topics get more questions at easier difficulty\n"
            f"- MANUAL: all topics use {difficulty} difficulty\n"
            "- 'focus on X' in instructions → allocate more to X\n"
            "- 'avoid Y' in instructions → set Y count to 0\n"
            "- Use 2-4 topics maximum"
        )

        plan = call_ollama_json(
            prompt,
            system="You are an expert quiz planning agent. Return only valid JSON.",
            fallback=None,
            model=QUIZ_MODEL, num_ctx=QUIZ_CTX,
        )

        if not isinstance(plan, dict) or not plan.get("topics"):
            plan = self._fallback_plan(num, mode, difficulty, weak)

        # Ensure per_topic_count sums to num
        counts = plan.get("per_topic_count", {})
        if counts:
            total = sum(counts.values())
            if total != num:
                tlist = list(counts.keys())
                counts[tlist[-1]] = max(1, counts[tlist[-1]] + (num - total))
                plan["per_topic_count"] = counts

        ctx._quiz_plan = plan

    def _fallback_plan(self, num, mode, difficulty, weak):
        topics = list(weak.keys())[:3] if weak else ["Core Concepts", "Key Principles", "Applications"]
        topics = topics[:min(len(topics), 3)]
        base   = num // len(topics)
        rem    = num - base * len(topics)
        counts = {t: base for t in topics}
        counts[topics[0]] += rem
        diffs  = {
            t: (difficulty if mode == "manual"
                else ("easy" if weak.get(t, 1.0) < 0.4 else "medium"))
            for t in topics
        }
        return {
            "topics": topics,
            "per_topic_count": counts,
            "per_topic_difficulty": diffs,
            "bloom_targets": ["Remember", "Understand", "Apply"],
            "focus_notes": "Fallback plan — planner LLM returned invalid JSON",
        }

    def validate(self, ctx):
        if not getattr(ctx, '_quiz_plan', None) or not ctx._quiz_plan.get("topics"):
            return False, ["Planner produced no valid plan"]
        return True, []


# ── Quiz Validator Agent ─────────────────────────────────────────

class QuizValidatorAgent(BaseAgent):
    """Validates a batch of generated questions on 3 dimensions:
    (1) difficulty accuracy, (2) content grounding, (3) instruction compliance.
    Pushes per-question verdicts as thinking entries."""
    name = "QuizValidatorAgent"

    def validate_batch(self, questions, quiz_type, content_snippet,
                       target_difficulty, user_instructions, thinking_cb=None):
        """Returns (passed_questions, failed_with_reasons)."""
        if not questions:
            return [], []

        q_texts = []
        for i, q in enumerate(questions[:10]):
            txt = q.get("question") or q.get("statement") or ""
            ans = str(
                q.get("correct_answer") or q.get("answer") or q.get("model_answer") or ""
            )[:80]
            q_texts.append(f"{i+1}. Q: {txt[:200]}\n   A: {ans}")

        batch_text = "\n\n".join(q_texts)
        instr_note = (
            f"\nUser instructions to comply with: {user_instructions}"
            if user_instructions else ""
        )

        prompt = (
            f"You are a quiz quality validator. For each question apply 3 checks:\n"
            f"1. DIFFICULTY: Is it genuinely {target_difficulty.upper()}?\n"
            f"2. GROUNDED: Does it relate to the content snippet below?\n"
            f"3. INSTRUCTION: Does it follow the user's instructions?{instr_note}\n\n"
            f"Content snippet:\n{content_snippet[:700]}\n\n"
            f"Questions to validate:\n{batch_text}\n\n"
            'Return STRICT JSON:\n'
            '{"verdicts":[{"num":1,"difficulty":"PASS","grounded":"PASS",'
            '"instruction":"PASS","overall":"PASS","reason":"brief 1-line fix hint if FAIL"}]}'
        )

        result = call_ollama_json(
            prompt,
            system="You are a strict quiz validator. Return only JSON.",
            fallback={"verdicts": []},
            model=QUIZ_MODEL, num_ctx=QUIZ_CTX,
        )
        verdicts = result.get("verdicts", []) if isinstance(result, dict) else []

        passed, failed = [], []
        for i, q in enumerate(questions):
            v = next((v for v in verdicts if v.get("num") == i + 1), None)
            if v:
                overall = str(v.get("overall", "PASS")).upper()
                reason  = v.get("reason", "")
                entry = {
                    "agent":        self.name,
                    "step":         "validate_question",
                    "verdict":      overall,
                    "question_num": i + 1,
                    "reason":       reason,
                    "checks": {
                        "difficulty":  v.get("difficulty",  "PASS"),
                        "grounded":    v.get("grounded",    "PASS"),
                        "instruction": v.get("instruction", "PASS"),
                    },
                }
                if thinking_cb:
                    thinking_cb(entry)
                if overall == "PASS":
                    passed.append(q)
                else:
                    failed.append({**q, "_fail_reason": reason})
            else:
                # No verdict returned for this question → keep it
                passed.append(q)

        return passed, failed


# ── Quiz Refiner Agent ───────────────────────────────────────────

class QuizRefinerAgent(BaseAgent):
    """Re-writes questions that failed validation with targeted correction prompts.
    Only re-generates the failures, not the whole quiz."""
    name = "QuizRefinerAgent"

    def refine_batch(self, failed_questions, quiz_type, content_snippet,
                     target_difficulty, user_instructions, thinking_cb=None):
        """Returns refined question list replacing the failed ones."""
        if not failed_questions:
            return []

        schema     = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
        instr_note = f"\nStudent instructions: {user_instructions}" if user_instructions else ""

        items = []
        for i, q in enumerate(failed_questions):
            txt    = q.get("question") or q.get("statement") or ""
            reason = q.get("_fail_reason", "does not meet quality standard")
            items.append(f"{i+1}. [FAILED] Q: {txt[:200]}\n   Fix needed: {reason}")

        if thinking_cb:
            thinking_cb({
                "agent":   self.name,
                "step":    "refining",
                "verdict": "OVERRIDE",
                "reason":  f"Rewriting {len(failed_questions)} failed question(s) with targeted fixes…",
            })

        prompt = (
            f"Rewrite each failed question to fix the stated issue.\n\n"
            f"Requirements:\n"
            f"- Difficulty: {target_difficulty.upper()}\n"
            f"- Must be grounded in lecture content: {content_snippet[:500]}\n"
            f"{instr_note}\n\n"
            f"Failed questions ({len(failed_questions)}):\n"
            + "\n\n".join(items) + "\n\n"
            f"Return STRICT JSON with exactly {len(failed_questions)} rewritten question(s):\n{schema}"
        )

        data    = call_ollama_json_quiz(prompt, key="questions")
        refined = data.get("questions", []) if isinstance(data, dict) else []

        if thinking_cb:
            thinking_cb({
                "agent":   self.name,
                "step":    "refined",
                "verdict": "DONE",
                "reason":  f"Produced {len(refined)} refined replacement(s)",
            })

        return refined


# ── Per-question agentic loop ──────────────────────────────────
# Plan -> for each slot: Generate ONE -> Validate -> Retry until pass -> Next.
# Every transition pushes a structured event so the UI can render real-time.

DIFF_DESCRIPTIONS = {
    "easy":   "tests basic recall, definitions, single-concept understanding",
    "medium": "requires understanding + application, may combine 2 concepts",
    "hard":   "requires deep analysis, multi-step reasoning, edge cases, synthesis",
}


def _validate_one_question(quiz_type, question, content_snip, target_difficulty,
                            user_instr, existing_questions=None):
    """Run the validator on a SINGLE question.
    Returns (passed: bool, fail_reason: str, checks: dict)."""
    existing_questions = existing_questions or []
    qtxt = question.get("question") or question.get("statement") or ""
    ans  = str(
        question.get("correct_answer") or question.get("answer") or
        question.get("model_answer") or ""
    )[:120]

    # Fill-blank: enforce blank marker before doing anything else
    if quiz_type == "fill" and "___" not in qtxt:
        return False, "fill question missing '___' blank marker", {
            "uniqueness": "PASS", "difficulty": "PASS",
            "grounded": "PASS", "instruction": "PASS",
        }

    # cheap deterministic uniqueness check first
    qtxt_low = qtxt.strip().lower()
    if qtxt_low:
        for ex in existing_questions:
            ex_text = (ex.get("question") or ex.get("statement") or "").strip().lower()
            if not ex_text:
                continue
            if ex_text == qtxt_low:
                return False, "duplicate of an existing question", {
                    "uniqueness": "FAIL", "difficulty": "PASS",
                    "grounded": "PASS", "instruction": "PASS",
                }
            ratio = difflib.SequenceMatcher(None, ex_text, qtxt_low).ratio()
            if ratio > 0.85:
                return False, f"too similar ({int(ratio*100)}%) to an existing question", {
                    "uniqueness": "FAIL", "difficulty": "PASS",
                    "grounded": "PASS", "instruction": "PASS",
                }

    diff_desc  = DIFF_DESCRIPTIONS.get(target_difficulty, "")
    instr_note = (
        f"\nMust comply with student instructions: {user_instr}"
        if user_instr else ""
    )

    prompt = (
        "You are a strict but fair quiz quality validator. Apply 3 checks to ONE question:\n"
        f"1. DIFFICULTY: Is it genuinely {target_difficulty.upper()}? ({diff_desc})\n"
        "2. GROUNDED: Is it answerable from the content snippet below?\n"
        f"3. INSTRUCTION: Does it follow the student instructions?{instr_note}\n\n"
        f"Content snippet:\n{content_snip[:700]}\n\n"
        "Question to validate:\n"
        f"Q: {qtxt[:300]}\n"
        f"A: {ans}\n\n"
        "Return STRICT JSON:\n"
        '{"difficulty":"PASS|FAIL","grounded":"PASS|FAIL",'
        '"instruction":"PASS|FAIL","overall":"PASS|FAIL",'
        '"reason":"1-line concrete fix hint if FAIL"}'
    )

    result = call_ollama_json(
        prompt,
        system="You are a strict but fair quiz validator. Return only JSON.",
        fallback={"overall": "PASS"},
        model=QUIZ_MODEL, num_ctx=QUIZ_CTX,
    )
    if not isinstance(result, dict):
        return True, "", {
            "difficulty": "PASS", "grounded": "PASS",
            "instruction": "PASS", "uniqueness": "PASS",
        }

    overall = str(result.get("overall", "PASS")).upper()
    reason  = result.get("reason", "")
    checks  = {
        "difficulty":  str(result.get("difficulty",  "PASS")).upper(),
        "grounded":    str(result.get("grounded",    "PASS")).upper(),
        "instruction": str(result.get("instruction", "PASS")).upper(),
        "uniqueness":  "PASS",
    }
    # Enforce hard blocking checks deterministically:
    # - grounded must pass
    # - instruction must pass only when user gave instructions
    hard_fail = []
    if checks["grounded"] == "FAIL":
        hard_fail.append("grounded")
    if user_instr and checks["instruction"] == "FAIL":
        hard_fail.append("instruction")
    if checks["uniqueness"] == "FAIL":
        hard_fail.append("uniqueness")
    if hard_fail:
        return False, reason or f"failed checks: {', '.join(hard_fail)}", checks

    # Difficulty mismatch can create long retry loops (especially HARD mode).
    # Treat it as non-blocking if the question is grounded and instruction-safe.
    if checks["difficulty"] == "FAIL":
        if target_difficulty == "hard":
            return True, "difficulty borderline; accepted (grounded + instruction-safe)", checks
        return False, reason or "difficulty does not match target", checks

    return overall == "PASS", reason, checks


def _generate_one_question_with_retry(
    quiz_type, topic, difficulty, content_snip, summary,
    user_instr, existing_questions, question_idx,
    max_attempts=3, thinking_cb=None, agent_name="QuizGenerator",
):
    """Generate ONE question, validate it, retry until pass or max attempts.
    Pushes per-question events so the UI can render an attempt-by-attempt view."""
    qid                  = f"q{question_idx + 1}"
    last_failure_reason  = ""
    last_question        = None
    schema               = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
    diff_desc            = DIFF_DESCRIPTIONS.get(difficulty, "")
    strategy_tracks      = {
        "mcq": [
            "Concept check style: test one core concept precisely with plausible distractors.",
            "Scenario style: short practical situation and ask the best next reasoning choice.",
            "Compare-and-justify style: ask which option is best and why others are weaker.",
        ],
        "tf": [
            "Direct claim style: one clear factual statement from content.",
            "Conditional claim style: statement involving if/when constraints from lecture.",
            "Misconception trap style: plausible but subtle statement requiring careful reading.",
        ],
        "fill": [
            "Keyword recall style: blank out a key term or phrase from lecture concepts.",
            "Process step style: blank out a critical step/condition in a method.",
            "Contrast style: blank that distinguishes two similar concepts.",
        ],
        "short": [
            "Explain-why style: ask for concise reasoning of a concept decision.",
            "Apply-to-case style: apply concept to a concrete mini-case.",
            "Trade-off style: discuss a clear trade-off grounded in lecture points.",
        ],
    }
    strategies = strategy_tracks.get(quiz_type, strategy_tracks["mcq"])

    for attempt in range(1, max_attempts + 1):
        strategy = strategies[(attempt - 1) % len(strategies)]
        if thinking_cb:
            thinking_cb({
                "agent":        agent_name,
                "step":         "question_generating",
                "phase":        "GENERATE",
                "question_id":  qid,
                "question_idx": question_idx,
                "topic":        topic,
                "difficulty":   difficulty,
                "attempt":      attempt,
                "max_attempts": max_attempts,
                "feedback":     last_failure_reason if attempt > 1 else "",
                "verdict":      "RETRY" if attempt > 1 else "PASS",
                "reason": (
                    f"Generating Q{question_idx + 1} on '{topic}' "
                    f"at {difficulty.upper()} difficulty…"
                    if attempt == 1
                    else f"Retry {attempt}/{max_attempts} — fixing: {last_failure_reason}"
                ),
                "strategy": strategy,
            })

        feedback_block = ""
        if last_failure_reason and attempt > 1:
            feedback_block = (
                f"\nPREVIOUS ATTEMPT FAILED: {last_failure_reason}\n"
                "You MUST fix this exact issue in the new question.\n"
            )

        existing_summary = ""
        if existing_questions:
            existing_texts = [
                (q.get("question") or q.get("statement", ""))[:90]
                for q in existing_questions
            ]
            existing_summary = (
                "\nDO NOT duplicate or paraphrase these accepted questions:\n- "
                + "\n- ".join(existing_texts) + "\n"
            )

        instr_line = f"\nStudent instructions: {user_instr}\n" if user_instr else ""

        prompt = (
            f"Generate exactly ONE {quiz_type.upper()} question wrapped in a JSON array.\n"
            f"Topic: '{topic}'\n"
            f"Difficulty: {difficulty.upper()} — {diff_desc}\n"
            f"Question strategy for this attempt: {strategy}\n"
            f"{feedback_block}"
            f"{existing_summary}"
            f"{instr_line}"
            f"Output STRICT JSON in this exact shape (a 'questions' array containing exactly 1 object):\n"
            f"{schema}\n\n"
            f"CONTENT:\n{content_snip[:5500]}\n\n"
            f"SUMMARY:\n{summary[:1200]}"
        )

        # Per-question generation uses the primary model first, then falls back
        # to the quiz JSON helper (which has additional parser/repair logic).
        data = call_ollama_json(prompt, fallback={"questions": []})
        qs   = data.get("questions", []) if isinstance(data, dict) else []
        if not qs:
            repaired = call_ollama_json_quiz(prompt, key="questions")
            qs = repaired.get("questions", []) if isinstance(repaired, dict) else []

        if not qs:
            last_failure_reason = "model returned no question"
            if thinking_cb:
                thinking_cb({
                    "agent":       "QuizValidatorAgent",
                    "step":        "question_validated",
                    "phase":       "VALIDATE",
                    "question_id": qid,
                    "verdict":     "FAIL",
                    "attempt":     attempt,
                    "reason":      last_failure_reason,
                })
            continue

        q = qs[0]
        if not q.get("topic"):      q["topic"]      = topic
        if not q.get("difficulty"): q["difficulty"] = difficulty
        last_question = q

        passed, fail_reason, checks = _validate_one_question(
            quiz_type, q, content_snip, difficulty, user_instr,
            existing_questions=existing_questions,
        )

        qtxt = q.get("question") or q.get("statement") or ""

        if passed:
            if thinking_cb:
                thinking_cb({
                    "agent":       "QuizValidatorAgent",
                    "step":        "question_validated",
                    "phase":       "VALIDATE",
                    "question_id": qid,
                    "verdict":     "PASS",
                    "attempt":     attempt,
                    "checks":      checks,
                    "reason":      "All quality checks passed",
                })
                thinking_cb({
                    "agent":            agent_name,
                    "step":             "question_finalized",
                    "phase":            "GENERATE",
                    "question_id":      qid,
                    "question_idx":     question_idx,
                    "verdict":          "DONE",
                    "status":           "PASS",
                    "total_attempts":   attempt,
                    "topic":            topic,
                    "difficulty":       difficulty,
                    "question_preview": qtxt[:160],
                    "reason":           f"Q{question_idx + 1} accepted on attempt {attempt}",
                })
            return q

        last_failure_reason = fail_reason or "validation failed"
        if thinking_cb:
            thinking_cb({
                "agent":       "QuizValidatorAgent",
                "step":        "question_validated",
                "phase":       "VALIDATE",
                "question_id": qid,
                "verdict":     "FAIL",
                "attempt":     attempt,
                "checks":      checks,
                "reason":      last_failure_reason,
            })

    # Max attempts reached:
    # - keep the last generated attempt (attempt #3) even if validator failed.
    # - this preserves strict requested difficulty and avoids pipeline stop.
    if last_question is not None:
        if thinking_cb:
            thinking_cb({
                "agent":            agent_name,
                "step":             "question_finalized",
                "phase":            "GENERATE",
                "question_id":      qid,
                "question_idx":     question_idx,
                "verdict":          "DONE",
                "status":           "PASS",
                "total_attempts":   max_attempts,
                "topic":            topic,
                "difficulty":       difficulty,
                "question_preview": (last_question.get("question") or last_question.get("statement") or "")[:160],
                "reason":           f"Attempt {max_attempts} accepted after max retries. Last issue: {last_failure_reason}",
            })
        return last_question

    # Rare case: model produced no question in all attempts. Do a strict JSON recovery
    # and accept the recovered output as final attempt content.
    recovery_prompt = (
        f"Return exactly ONE valid {quiz_type.upper()} question in STRICT JSON:\n"
        f"{schema}\n"
        f"Topic: {topic}\nDifficulty: {difficulty.upper()}\n"
        f"Do not return explanations outside JSON.\n"
        f"CONTENT:\n{content_snip[:4500]}\n\nSUMMARY:\n{summary[:1000]}"
    )
    for _ in range(3):
        repaired = call_ollama_json_quiz(recovery_prompt, key="questions")
        rec_qs = repaired.get("questions", []) if isinstance(repaired, dict) else []
        if not rec_qs:
            continue
        q = rec_qs[0]
        if not q.get("topic"):      q["topic"] = topic
        if not q.get("difficulty"): q["difficulty"] = difficulty
        if thinking_cb:
            thinking_cb({
                "agent":            agent_name,
                "step":             "question_finalized",
                "phase":            "GENERATE",
                "question_id":      qid,
                "question_idx":     question_idx,
                "verdict":          "DONE",
                "status":           "PASS",
                "total_attempts":   max_attempts,
                "topic":            topic,
                "difficulty":       difficulty,
                "question_preview": (q.get("question") or q.get("statement") or "")[:160],
                "reason":           f"Recovered and accepted after {max_attempts} attempts returned no question",
            })
        return q

    raise RuntimeError(
        f"Failed to generate any {quiz_type.upper()} question text after {max_attempts} attempts for slot {question_idx + 1}"
    )


def _run_quiz_agentic_loop(ctx, quiz_type, content, summary,
                            user_instr, target_difficulty, thinking_cb=None,
                            max_attempts_per_q=3):
    """Per-question agentic loop:
    Phase 1 PLAN -> Phase 2 GENERATE one + Phase 3 VALIDATE + retry -> Phase 4 DONE.
    Each event has phase/question_id/attempt/verdict so the UI can show a live pipeline."""
    agent_name = getattr(ctx, '_agent_name', 'QuizAgent')
    num        = getattr(ctx, '_quiz_count', 10)

    # ── PHASE 1: PLAN ──────────────────────────────────────────
    if thinking_cb:
        thinking_cb({
            "agent":   "Orchestrator",
            "step":    "phase_start",
            "phase":   "PLAN",
            "verdict": "PASS",
            "reason":  "Phase 1/4 — Planner deciding topic distribution & difficulty…",
        })

    plan = getattr(ctx, '_quiz_plan', None)
    if not plan:
        QuizPlannerAgent().run(ctx)
        plan = getattr(ctx, '_quiz_plan', None) or {}

    topics      = plan.get("topics", [])
    topic_cnt   = plan.get("per_topic_count",      {})
    topic_diff  = plan.get("per_topic_difficulty", {})
    focus_notes = plan.get("focus_notes", "")

    if thinking_cb:
        thinking_cb({
            "agent":   "QuizPlannerAgent",
            "step":    "plan_ready",
            "phase":   "PLAN",
            "verdict": "DONE",
            "reason":  f"Plan ready: {len(topics)} topic(s) over {num} question(s)",
            "data": {
                "topics":               topics,
                "per_topic_count":      topic_cnt,
                "per_topic_difficulty": topic_diff,
                "focus_notes":          focus_notes,
            },
        })

    # build flat sequence of (topic, difficulty) slots
    slots = []
    if topics and topic_cnt:
        for topic in topics:
            count = topic_cnt.get(topic, 0)
            diff  = topic_diff.get(topic, target_difficulty)
            for _ in range(count):
                slots.append((topic, diff))
    if not slots:
        slots = [("General", target_difficulty)] * num

    # respect num exactly: trim or pad
    if len(slots) > num:
        slots = slots[:num]
    while len(slots) < num:
        last = slots[-1] if slots else ("General", target_difficulty)
        slots.append(last)

    # ── PHASE 2 + 3: GENERATE + VALIDATE per question ──────────
    if thinking_cb:
        thinking_cb({
            "agent":   agent_name,
            "step":    "phase_start",
            "phase":   "GENERATE",
            "verdict": "PASS",
            "reason": (
                f"Phase 2/4 — Generating {len(slots)} question(s) one at a time, "
                f"validating each, retrying on fail (max {max_attempts_per_q} attempts each)…"
            ),
            "data":    {"total_slots": len(slots)},
        })

    final_questions = []
    content_snip    = content[:6000]

    for idx, (topic, diff) in enumerate(slots):
        q = _generate_one_question_with_retry(
            quiz_type=quiz_type, topic=topic, difficulty=diff,
            content_snip=content_snip, summary=summary,
            user_instr=user_instr,
            existing_questions=final_questions,
            question_idx=idx,
            max_attempts=max_attempts_per_q,
            thinking_cb=thinking_cb,
            agent_name=agent_name,
        )
        if q:
            final_questions.append(q)

    # ── PHASE 4: DONE ─────────────────────────────────────────
    if thinking_cb:
        thinking_cb({
            "agent":   agent_name,
            "step":    "phase_complete",
            "phase":   "DONE",
            "verdict": "DONE",
            "reason":  f"All phases complete — {len(final_questions)} question(s) ready",
            "data":    {"total": len(final_questions)},
        })

    return final_questions


# ── Manual mode difficulty agents ────────────────────────────────

class _DifficultySubAgent(BaseAgent):
    """Base for Easy/Medium/Hard quiz agents.
    Delegates to the per-question agentic loop at a fixed target difficulty."""
    target_difficulty = "medium"
    max_retries       = 2
    _thinking_log     = []

    def _get_content(self, ctx):
        if ctx.content_type == "video":
            return pack_video_sources(ctx.board_text, ctx.transcript,
                                      ctx.board_entries or [], 7500)
        return (ctx.doc_text or "")[:7500]

    def _build_prompt(self, quiz_type, num, content, summary, user_instructions=""):
        """Backward-compat helper kept for the legacy refine path in /quiz/generate."""
        diff   = self.target_difficulty
        instr  = f"\nAdditional instructions: {user_instructions}\n" if user_instructions else ""
        schema = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
        return (
            f"Generate exactly {num} {quiz_type.upper()} questions at {diff.upper()} difficulty.\n"
            f"{instr}"
            f"STRICT JSON: {schema}\n\n"
            f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
        )

    def _validate_difficulty(self, questions, quiz_type):
        """Backward-compat helper — uses the per-question validator."""
        if not questions:
            return [], []
        kept, thinking = [], []
        for i, q in enumerate(questions):
            ok, reason, checks = _validate_one_question(
                quiz_type, q, "", self.target_difficulty, "",
                existing_questions=questions[:i],
            )
            thinking.append({
                "agent":   self.name,
                "step":    "difficulty_check",
                "verdict": "PASS" if ok else "FAIL",
                "checks":  checks,
                "reason":  reason or "passed",
            })
            if ok:
                kept.append(q)
        if not kept:
            thinking.append({
                "agent": self.name, "step": "override",
                "verdict": "OVERRIDE",
                "reason": "All questions rejected by validator — keeping originals",
            })
            return questions, thinking
        return kept, thinking


class EasyQuizAgent(_DifficultySubAgent):
    name = "EasyQuizAgent"
    target_difficulty = "easy"

    def plan(self, ctx):
        return "Generate easy-level questions (recall, definitions, basics)"

    def should_run(self, ctx):
        return getattr(ctx, '_manual_difficulty', '') == 'easy'


class MediumQuizAgent(_DifficultySubAgent):
    name = "MediumQuizAgent"
    target_difficulty = "medium"

    def plan(self, ctx):
        return "Generate medium-level questions (understanding, application)"

    def should_run(self, ctx):
        return getattr(ctx, '_manual_difficulty', '') == 'medium'


class HardQuizAgent(_DifficultySubAgent):
    name = "HardQuizAgent"
    target_difficulty = "hard"

    def plan(self, ctx):
        return "Generate hard-level questions (analysis, synthesis, edge cases)"

    def should_run(self, ctx):
        return getattr(ctx, '_manual_difficulty', '') == 'hard'


# ── Adaptive Quiz Agent ──────────────────────────────────────────

class AdaptiveQuizAgent(BaseAgent):
    """Reads student performance + content, then runs the per-question agentic
    loop with planner-driven topic + difficulty allocation."""
    name          = "AdaptiveQuizAgent"
    max_retries   = 2
    _thinking_log = []

    def _add_thinking(self, ctx, entry):
        self._thinking_log.append(entry)
        cb = getattr(ctx, "thinking_cb", None)
        if cb:
            cb(entry)

    def plan(self, ctx):
        weak = difficulty_adapter.get_weak_topics()
        if weak:
            topics = ", ".join(f"{t} ({int(s*100)}%)" for t, s in weak.items())
            return f"Adaptive quiz targeting weak topics: {topics}"
        return "No history yet — generating balanced quiz across all topics"

    def execute(self, ctx):
        self._thinking_log = []
        quiz_type  = getattr(ctx, '_quiz_type', 'mcq')
        user_instr = getattr(ctx, '_user_instructions', '')
        perf       = difficulty_adapter.get_all_stats()
        weak       = difficulty_adapter.get_weak_topics()

        self._add_thinking(ctx, {
            "agent":   self.name,
            "step":    "analyze_performance",
            "phase":   "PLAN",
            "verdict": "PASS",
            "data":    {"all_stats": perf, "weak_topics": weak},
            "reason": (
                f"Found {len(weak)} weak topic(s) out of {len(perf)} tracked"
                if perf else "No performance history yet — will generate balanced quiz"
            ),
        })

        content = (
            pack_video_sources(ctx.board_text, ctx.transcript, ctx.board_entries or [], 7500)
            if ctx.content_type == "video"
            else (ctx.doc_text or "")[:7500]
        )
        summary = ctx.summary or state.status.get("summary", "")

        ctx._mode       = "adaptive"
        ctx._agent_name = self.name

        def thinking_cb(entry):
            self._add_thinking(ctx, entry)

        final = _run_quiz_agentic_loop(
            ctx, quiz_type, content, summary, user_instr,
            target_difficulty="medium",
            thinking_cb=thinking_cb,
        )

        state.quiz[quiz_type]      = final
        state._quiz_agent_thinking = self._thinking_log

    def validate(self, ctx):
        quiz_type = getattr(ctx, '_quiz_type', 'mcq')
        qs = state.quiz.get(quiz_type, [])
        if not qs:
            return False, [f"Adaptive agent produced 0 {quiz_type} questions"]
        return True, []



class FlashcardAgent(BaseAgent):
    """Agentic flashcard generator:
    1. Generates cards in batches of 8
    2. Validates each card (front != back, non-empty, not duplicate)
    3. Retries rejected cards up to max_retries times
    4. Stops when 15 unique valid cards are collected or max attempts hit."""
    name        = "FlashcardGen"
    max_retries = 3
    TARGET      = 15

    def plan(self, ctx):
        return f"Agentic flashcard pipeline: generate → validate → retry until {self.TARGET} unique cards"

    def _content(self, ctx):
        if ctx.content_type == "video":
            vis = format_visual_lecture_text(ctx.board_text, ctx.board_entries or [], 4000)
            parts = []
            if vis:
                parts.append(f"=== SLIDES + BOARDS ===\n{vis}")
            parts.append(f"=== SUMMARY ===\n{(ctx.summary or '')[:3500]}")
            return "\n\n".join(parts)[:6500]
        doc = (ctx.doc_text or "")[:4000]
        return (doc + "\n\n=== SUMMARY ===\n" + (ctx.summary or "")[:2500])[:6500]

    def _validate_card(self, card, seen_fronts):
        """Return (ok, reason). Deterministic checks only."""
        front = card.get("front", "").strip()
        back  = card.get("back",  "").strip()
        if not front or not back:
            return False, "empty front or back"
        if front.lower() == back.lower():
            return False, "front == back"
        if len(back) < 8:
            return False, "back too short to be useful"
        if front.lower() in seen_fronts:
            return False, "duplicate front"
        return True, ""

    def execute(self, ctx):
        content     = self._content(ctx)
        seen_fronts = set()
        accepted    = []
        batch_num   = 0
        max_batches = self.max_retries + 2

        while len(accepted) < self.TARGET and batch_num < max_batches:
            batch_num   += 1
            need         = self.TARGET - len(accepted)
            ask          = min(need + 3, 10)    # ask a few extra to absorb rejects
            existing_txt = (
                "\nALREADY ACCEPTED (do not repeat):\n- "
                + "\n- ".join(c["front"] for c in accepted)
                if accepted else ""
            )
            prompt = (
                f"Generate exactly {ask} flashcards for key concepts from this lecture.\n"
                f"Batch {batch_num} — need {need} more accepted cards.\n"
                f"{existing_txt}\n\n"
                "STRICT JSON: {\"cards\":[{\"front\":\"term or question\","
                "\"back\":\"clear explanation (≥10 words)\","
                "\"topic\":\"\",\"difficulty\":\"easy|medium|hard\"}]}\n"
                "RULES: front != back. back must be at least 10 words. "
                "Each card must cover a different concept.\n\n"
                f"CONTENT:\n{content}"
            )
            data  = call_ollama_json_quiz(prompt, key="cards")
            batch = data.get("cards", []) if isinstance(data, dict) else []

            for card in batch:
                ok, reason = self._validate_card(card, seen_fronts)
                if ok:
                    seen_fronts.add(card["front"].strip().lower())
                    accepted.append(card)
                    if len(accepted) >= self.TARGET:
                        break
                else:
                    print(f"  [FlashcardGen] Rejected card: {reason} — '{card.get('front','')[:50]}'")

        state.quiz["flash"] = accepted
        print(f"  [FlashcardGen] ✅ {len(accepted)} unique valid flashcards generated in {batch_num} batch(es)")

    def validate(self, ctx):
        cards = state.quiz.get("flash", [])
        if not cards:
            return False, ["FlashcardAgent produced 0 cards"]
        issues = []
        if len(cards) < 8:
            issues.append(f"only {len(cards)} cards — expected ~{self.TARGET}")
        return True, issues



class SuggestedQuestionsAgent(BaseAgent):
    """Generates suggested chat questions for the tutor."""
    name = "SuggestQuestions"

    def plan(self, ctx):
        return "Generate 5 suggested study questions for the chat tutor"

    def execute(self, ctx):
        if ctx.content_type == "video":
            head = pack_video_sources(
                ctx.board_text, ctx.transcript, ctx.board_entries or [], 5500,
            )
            head = f"{head}\n\n=== SUMMARY ===\n{ctx.summary[:2200]}"
        else:
            doc = (ctx.doc_text or "")[:4000]
            head = f"{doc}\n\n=== SUMMARY ===\n{ctx.summary[:2200]}" if doc else ctx.summary[:4000]
        sq = call_ollama_json_list_quiz(
            "Generate 5 study questions a student might ask about this lecture "
            "(slides, boards, speech). Return a JSON list of strings:\n"
            f"{head[:7800]}",
            fallback_list=[],
        )
        state.suggested_questions = sq if isinstance(sq, list) else []



print("✅ Specialist agents ready.")

# ══════════════════════════════════════════════════════════════
# ORCHESTRATOR AGENT
# ══════════════════════════════════════════════════════════════



class OrchestratorAgent:
    """The brain — decides which agents to run, in what order,
    with what parameters, and handles failures gracefully.

    Agents inside the same step-list entry run in parallel threads.
    Each step waits for the previous step to finish before starting."""

    # Each entry is a list of agent classes that can run concurrently.
    # Entries themselves execute sequentially (step 1 finishes → step 2 starts).
    PIPELINE = {
        "video": [
            [ExtractorAgent, TranscriberAgent],                # parallel: OCR + Whisper
            [SummaryAgent],
            [RAGBuilderAgent, HintsAgent],                     # parallel: index + hints
            [FlashcardAgent, SuggestedQuestionsAgent],         # parallel: cards + suggestions
        ],
        "document": [
            [DocParserAgent],
            [SummaryAgent],
            [RAGBuilderAgent],
            [FlashcardAgent, SuggestedQuestionsAgent],         # parallel: cards + suggestions
        ],
    }

    def estimate_quiz_counts(self, ctx):
        """Adapt question counts to content length. The agentic pipeline runs
        per-question with retry+validation so we keep these tighter than the
        old upfront generator: a typical lecture lands around 18-22 questions."""
        content_len = len(ctx.transcript) + len(ctx.board_text) + len(ctx.doc_text)
        if content_len > 10000:
            return {"mcq": 10, "tf": 6, "fill": 6, "short": 5}
        if content_len > 5000:
            return {"mcq": 8, "tf": 5, "fill": 5, "short": 4}
        return {"mcq": 5, "tf": 4, "fill": 4, "short": 3}

    def build_difficulty_preference(self):
        """Use student performance history to suggest difficulty weighting."""
        perf = state.student_performance
        if not perf:
            return None
        weak = [t for t, scores in perf.items()
                if scores and (sum(scores) / len(scores)) < 0.5]
        if weak:
            return f"focus harder questions on weak topics: {', '.join(weak[:5])}"
        return None

    def _run_agent(self, agent, ctx):
        """Run a single agent, catching crashes so the pipeline continues."""
        try:
            agent.run(ctx)
        except Exception as e:
            self._log(ctx, "ERROR", f"{agent.name} crashed: {e} — continuing")

    def _run_parallel_group(self, agent_classes, ctx):
        """Run a group of agents concurrently and wait for all to finish."""
        agents = [cls() for cls in agent_classes]

        # prepare quiz agent params before launching
        for a in agents:
            if isinstance(a, QuizAgent):
                ctx.quiz_counts     = self.estimate_quiz_counts(ctx)
                ctx.difficulty_pref = self.build_difficulty_preference()
                self._log(ctx, "PLAN",
                          f"Quiz counts: {ctx.quiz_counts}, "
                          f"difficulty: {ctx.difficulty_pref or 'balanced'}")

        if len(agents) == 1:
            self._run_agent(agents[0], ctx)
            return

        # multiple agents in this step — fan out into threads
        self._log(ctx, "PLAN",
                  f"Running in parallel: {', '.join(a.name for a in agents)}")
        with ThreadPoolExecutor(max_workers=len(agents)) as pool:
            futs = {pool.submit(self._run_agent, a, ctx): a for a in agents}
            for fut in as_completed(futs):
                fut.result()   # propagate exceptions if any

    def run(self, content_type, file_path):
        """Execute the full agentic pipeline with parallel steps."""
        ctx = AgentContext(content_type, file_path)
        state.session_id = state._make_session_id(ctx.source_name)
        state.agent_log  = []
        state.quiz_history = []

        pipeline = self.PIPELINE.get(content_type, [])
        total_steps = len(pipeline)

        self._log(ctx, "START",
                  f"Pipeline for {content_type}: {ctx.source_name} "
                  f"({total_steps} steps, parallel where possible)")

        for step_idx, agent_group in enumerate(pipeline):
            names = "+".join(cls.name if hasattr(cls, 'name') else cls.__name__
                            for cls in agent_group)
            pct = int(((step_idx + 0.5) / total_steps) * 100)
            state.update_status(
                state="running", stage=f"{names}…", pct=min(pct, 98),
            )
            self._run_parallel_group(agent_group, ctx)

        # finalize
        timeline = ""
        if content_type == "video":
            timeline = build_lecture_timeline_text(ctx.segment_list, ctx.board_entries)
            if timeline:
                save_text(
                    "lecture_timeline.txt",
                    "=== LECTURE TIMELINE (speech + slides & boards) ===\n\n" + timeline,
                )
        state.update_status(
            state="done", stage="Complete ✓", pct=100,
            summary=ctx.summary,
            transcript=ctx.transcript,
            board_text=ctx.board_text,
            lecture_timeline=timeline,
        )
        state.agent_log = ctx.agent_log
        state.save_session(ctx.source_name)
        self._log(ctx, "DONE",
                  f"Pipeline complete — {len(ctx.agent_log)} agent decisions logged")

    def _log(self, ctx, level, msg):
        entry = {
            "agent": "Orchestrator", "level": level,
            "time": datetime.now().isoformat(), "message": msg,
        }
        ctx.agent_log.append(entry)
        print(f"  [Orchestrator] {msg}")



orchestrator = OrchestratorAgent()



print("✅ Orchestrator ready.")

# ══════════════════════════════════════════════════════════════
# DIFFICULTY ADAPTER  (self-improving quiz difficulty)
# ══════════════════════════════════════════════════════════════



class DifficultyAdapter:
    """Tracks student performance per topic and generates targeted
    questions for weak areas at adjusted difficulty."""

    def record_performance(self, grade_results):
        """Called after every quiz submission.  Stores per-topic score history."""
        for r in grade_results.get("results", []):
            topic = r.get("topic", "General")
            if topic not in state.student_performance:
                state.student_performance[topic] = []
            state.student_performance[topic].append(1.0 if r.get("correct") else 0.0)

    def get_weak_topics(self, threshold=0.5):
        """Return topics where average score is below threshold."""
        weak = {}
        for topic, scores in state.student_performance.items():
            if not scores:
                continue
            avg = sum(scores) / len(scores)
            if avg < threshold:
                weak[topic] = round(avg, 2)
        return weak

    def get_all_stats(self):
        """Return full performance stats for every topic."""
        stats = {}
        for topic, scores in state.student_performance.items():
            if not scores:
                continue
            avg = sum(scores) / len(scores)
            stats[topic] = {
                "average": round(avg, 2),
                "attempts": len(scores),
                "correct": sum(1 for s in scores if s == 1.0),
                "total": len(scores),
                "level": "weak" if avg < 0.5 else "ok" if avg < 0.8 else "strong",
            }
        return stats

    def suggest_difficulty(self, topic):
        """Based on past performance, return recommended difficulty."""
        scores = state.student_performance.get(topic, [])
        if not scores:
            return "medium"
        avg = sum(scores) / len(scores)
        if avg > 0.8:
            return "hard"
        if avg < 0.4:
            return "easy"
        return "medium"

    def generate_targeted(self, quiz_type="mcq", num=5):
        """Generate new questions focused on weak topics at adjusted difficulty."""
        weak = self.get_weak_topics()
        if not weak:
            return []

        topic_list = ", ".join(weak.keys())
        diff_hints = ", ".join(
            f"{t}: {self.suggest_difficulty(t)}" for t in weak
        )
        summary = state.status.get("summary", "")
        if not summary:
            return []

        prompt = (
            f"Generate {num} {quiz_type.upper()} questions focused specifically on "
            f"these WEAK topics the student struggles with: {topic_list}.\n"
            f"Difficulty guidance per topic: {diff_hints}.\n"
            f"The student has been getting these wrong — make questions that help them learn.\n\n"
        )

        if quiz_type == "mcq":
            prompt += (
                "STRICT JSON: {\"questions\":[{\"question\",\"options\":{\"A\",\"B\",\"C\",\"D\"},"
                "\"correct_answer\",\"explanation\",\"topic\",\"difficulty\",\"bloom_level\"}]}\n\n"
            )
        elif quiz_type == "tf":
            prompt += (
                "STRICT JSON: {\"questions\":[{\"statement\",\"answer\",\"explanation\","
                "\"topic\",\"difficulty\"}]}\n\n"
            )
        elif quiz_type == "fill":
            prompt += (
                "STRICT JSON: {\"questions\":[{\"question\",\"answer\",\"hint\","
                "\"topic\",\"difficulty\"}]}\n\n"
            )
        else:
            prompt += (
                "STRICT JSON: {\"questions\":[{\"question\",\"model_answer\",\"key_points\","
                "\"marks\",\"topic\",\"difficulty\"}]}\n\n"
            )

        if state.video_path:
            be = load_board_entries_from_disk()
            bt = state.status.get("board_text", "")
            tr = state.status.get("transcript", "")
            content = (
                pack_video_sources(bt, tr, be, 5500)
                + "\n\n=== SUMMARY ===\n"
                + summary[:4000]
            )
        else:
            content = summary[:3000]
        prompt += f"CONTENT:\n{content[:8000]}"
        data = call_ollama_json_quiz(prompt, key="questions")
        return data.get("questions", []) if isinstance(data, dict) else []



difficulty_adapter = DifficultyAdapter()



print("✅ Difficulty adapter ready.")

# ══════════════════════════════════════════════════════════════
# AUTO-RESTORE  (load last session from Drive on startup)
# ══════════════════════════════════════════════════════════════



def _auto_restore():
    """Check Google Drive for saved sessions. Load the latest one
    so quizzes/summary/flashcards are available immediately."""
    saved = AppState.list_sessions()
    if not saved:
        print("📂 No saved sessions found — start fresh.")
        return

    latest = saved[0]
    state.load_session(latest["session_id"])
    print(f"🔄 Auto-loaded last session: {latest['source_file']} ({latest['session_id']})")
    print(f"   📊 {latest['mcq_count']} MCQ · {latest['tf_count']} T/F · {latest['flash_count']} flashcards")

    # Rebuild the RAG search index from saved segment/board files
    rag_entries = []
    seg_path   = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    board_path = os.path.join(OUTPUT_DIR, "board_entries.json")
    if os.path.exists(seg_path):
        with open(seg_path, "r", encoding="utf-8") as f:
            rag_entries += json.load(f)
    if os.path.exists(board_path):
        with open(board_path, "r", encoding="utf-8") as f:
            for e in json.load(f):
                rag_entries.append({
                    "text": e["text"], "start_str": e["timestamp_str"],
                    "tag": e.get("tag", ""),
                })
    if rag_entries:
        build_rag(rag_entries, ts_key="start_str", src_key="tag")
        print(f"   🔍 RAG index rebuilt with {len(rag_entries)} chunks")



_auto_restore()

# ══════════════════════════════════════════════════════════════
# FLASK API
# ══════════════════════════════════════════════════════════════



app = Flask(__name__)
CORS(app)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024   # 500 MB

# ── Health / Status ───────────────────────────────────────────



@app.route('/')
def index():
    return jsonify(
        app="LectureForge v1.1",
        status="API running",
        model=state.active_model,
        message="Open quizforge_ui.html locally to use the UI.",
    )



@app.route('/status')
def status_route():
    """Pipeline status + live agentic-quiz thinking events.
    The UI uses `quiz_thinking_session` to detect a fresh generation run."""
    payload = dict(state.status)
    payload["model"]                  = state.active_model
    payload["session_id"]             = state.session_id
    payload["quiz_thinking"]          = state._quiz_agent_thinking
    payload["quiz_thinking_session"]  = state._quiz_thinking_session
    payload["quiz_thinking_active"]   = state._quiz_thinking_active
    return jsonify(payload)

# ── Upload ────────────────────────────────────────────────────



@app.route('/upload/video', methods=['POST'])
def upload_video():
    f = request.files.get('video')
    if not f:
        return jsonify(ok=False, error="No file sent.")
    ext = Path(f.filename).suffix.lower().lstrip('.')
    if ext not in ALLOWED_VIDEO:
        return jsonify(ok=False, error=f"Unsupported type: {ext}")
    path = os.path.join(UPLOAD_DIR, secure_filename(f.filename))
    f.save(path)
    state.video_path = path
    state.reset_status()
    threading.Thread(target=orchestrator.run, args=("video", path), daemon=True).start()
    return jsonify(ok=True, filename=f.filename)



@app.route('/upload/document', methods=['POST'])
def upload_document():
    f = request.files.get('document')
    if not f:
        return jsonify(ok=False, error="No file sent.")
    ext = Path(f.filename).suffix.lower().lstrip('.')
    if ext not in ALLOWED_DOCS:
        return jsonify(ok=False, error=f"Unsupported type: {ext}")
    path = os.path.join(UPLOAD_DIR, secure_filename(f.filename))
    f.save(path)
    state.doc_path = path
    state.reset_status()
    threading.Thread(target=orchestrator.run, args=("document", path), daemon=True).start()
    return jsonify(ok=True, filename=f.filename)

# ── Quiz endpoints ────────────────────────────────────────────



@app.route('/quiz/mcq')
def quiz_mcq_route():
    qs = state.quiz.get("mcq", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/true_false')
def quiz_tf_route():
    qs = state.quiz.get("tf", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/fill_blank')
def quiz_fill_route():
    qs = state.quiz.get("fill", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/short_answer')
def quiz_short_route():
    qs = state.quiz.get("short", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/history')
def quiz_history_route():
    """Return past graded quiz submissions saved into this session."""
    return jsonify(history=state.quiz_history, total=len(state.quiz_history))


def _new_quiz_run_id():
    return f"run_{int(time.time() * 1000)}_{os.urandom(4).hex()}"


def _append_generated_quiz_run(source, quiz_type, mode, difficulty, questions):
    """Store a snapshot of a generated quiz so the UI can reload past sets."""
    if not questions:
        return
    with state.lock:
        run = {
            "id":          _new_quiz_run_id(),
            "created_at":  datetime.now().isoformat(),
            "source":      source,          # "async" | "upload"
            "quiz_type":   quiz_type,
            "mode":        mode,            # auto | adaptive | manual
            "difficulty":  difficulty or "",
            "count":       len(questions),
            "questions":   copy.deepcopy(questions),
        }
        state.generated_quiz_runs.insert(0, run)
        state.generated_quiz_runs = state.generated_quiz_runs[:50]


@app.route('/quiz/generated_runs')
def quiz_generated_runs_route():
    with state.lock:
        runs = list(state.generated_quiz_runs)
    light = []
    for r in runs:
        light.append({
            "id":         r.get("id"),
            "created_at": r.get("created_at"),
            "source":     r.get("source"),
            "quiz_type":  r.get("quiz_type"),
            "mode":       r.get("mode"),
            "difficulty": r.get("difficulty"),
            "count":      r.get("count", len(r.get("questions", []))),
        })
    return jsonify(runs=light, total=len(light))


@app.route('/quiz/generated_runs/restore', methods=['POST'])
def quiz_restore_generated_run_route():
    data   = request.get_json(silent=True) or {}
    run_id = (data.get("run_id") or "").strip()
    if not run_id:
        return jsonify(ok=False, error="run_id required"), 400
    with state.lock:
        found = None
        for r in state.generated_quiz_runs:
            if r.get("id") == run_id:
                found = r
                break
        if not found:
            return jsonify(ok=False, error="Run not found"), 404
        qt = found.get("quiz_type")
        qs = found.get("questions") or []
        if qt not in state.quiz:
            return jsonify(ok=False, error="Invalid quiz type in run"), 400
        state.quiz[qt] = copy.deepcopy(qs)
    return jsonify(ok=True, quiz_type=qt, total=len(qs))



def _load_segments_and_board_artifacts():
    segs = []
    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    if os.path.exists(seg_path):
        try:
            with open(seg_path, "r", encoding="utf-8") as f:
                segs = json.load(f)
        except Exception:
            segs = []
    be = load_board_entries_from_disk()
    return segs, be



def _record_quiz_submission(quiz_type, request_answers, result):
    # Ensure each graded question has a source_timestamp for UI clip links.
    segs, be = _load_segments_and_board_artifacts()
    graded_items = result.get("results", [])
    if graded_items:
        fill_missing_timestamps(graded_items, segments=segs, board_entries=be)

    entry = {
        "created_at": datetime.now().isoformat(),
        "quiz_type": quiz_type,
        "mode": "auto",
        "count": result.get("total", len(graded_items)),
        "score": result.get("score", 0),
        "total": result.get("total", len(graded_items)),
        "percentage": result.get("percentage", 0),
        "grade": result.get("grade", ""),
        "answers": request_answers or {},
        "results": graded_items,
    }
    state.quiz_history.append(entry)



@app.route('/quiz/mcq/grade', methods=['POST'])
def grade_mcq_route():
    data   = request.get_json(silent=True) or {}
    result = grade_mcq(state.quiz.get("mcq", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("mcq", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/tf/grade', methods=['POST'])
def grade_tf_route():
    data   = request.get_json(silent=True) or {}
    result = grade_tf(state.quiz.get("tf", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("tf", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/fill/grade', methods=['POST'])
def grade_fill_route():
    data   = request.get_json(silent=True) or {}
    result = grade_fill(state.quiz.get("fill", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("fill", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/short/reveal', methods=['POST'])
def reveal_short_route():
    data = request.get_json(silent=True) or {}
    qs   = state.quiz.get("short", [])
    revealed = {i: qs[i] for i in (data.get("indices") or range(len(qs))) if i < len(qs)}
    return jsonify(revealed=revealed)



@app.route('/quiz/short/grade', methods=['POST'])
def grade_short_route():
    data   = request.get_json(silent=True) or {}
    result = grade_short(state.quiz.get("short", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("short", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/explain', methods=['POST'])
def explain_mistakes_route():
    data  = request.get_json(silent=True) or {}
    wrong = data.get("wrong_questions", [])
    exps  = explain_wrong(wrong)
    return jsonify(explanations=exps)

@app.route('/quiz/generate', methods=['POST'])
def quiz_generate_route():
    """Generate quiz using the full agentic loop: Plan → Generate → Validate → Refine.
    Body: {mode: "adaptive"|"manual", difficulty: "easy"|"medium"|"hard",
           quiz_type: "mcq"|"tf"|"fill"|"short", count: 10,
           instructions: "...", action: "new"|"refine"}"""
    data       = request.get_json(silent=True) or {}
    mode       = data.get("mode", "adaptive")
    difficulty = data.get("difficulty", "medium")
    quiz_type  = data.get("quiz_type", "mcq")
    action     = data.get("action", "new")
    num        = min(int(data.get("count", 10)), 20)
    user_instr = data.get("instructions", "").strip()
    plan_mode  = data.get("plan_mode", "auto")
    manual_topics = data.get("manual_topics", []) or []

    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No content processed yet — upload a file first.",
                       questions=[], thinking=[])

    ctx = type('Ctx', (), {
        'content_type':       'video' if state.video_path else 'document',
        'board_text':         state.status.get("board_text", ""),
        'transcript':         state.status.get("transcript", ""),
        'board_entries':      load_board_entries_from_disk(),
        'doc_text':           state.status.get("summary", ""),
        'summary':            summary,
        'agent_log':          [],
        '_quiz_type':         quiz_type,
        '_quiz_count':        num,
        '_user_instructions': user_instr,
        '_plan_mode':         plan_mode,
        '_manual_topics':     manual_topics,
        '_manual_difficulty': difficulty,
        '_mode':              mode,
        '_quiz_plan':         None,
        '_agent_name':        '',
        'segment_list':       [],
        'file_path':          state.video_path or state.doc_path or "",
        'source_name':        Path(state.video_path or state.doc_path or "unknown").name,
        'thinking_cb':        None,
    })()

    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    if os.path.exists(seg_path):
        try:
            with open(seg_path, "r", encoding="utf-8") as f:
                ctx.segment_list = json.load(f)
        except Exception:
            ctx.segment_list = []

    if state.doc_path:
        doc_path = os.path.join(OUTPUT_DIR, "doc_text.txt")
        if os.path.exists(doc_path):
            with open(doc_path, "r", encoding="utf-8") as f:
                ctx.doc_text = f.read()

    thinking   = []
    def thinking_cb(entry):
        thinking.append(entry)
    ctx.thinking_cb = thinking_cb

    content = (
        pack_video_sources(ctx.board_text, ctx.transcript, ctx.board_entries or [], 7500)
        if ctx.content_type == "video"
        else (ctx.doc_text or "")[:7500]
    )

    # -- Refine mode: send existing questions through Refiner + Validator --
    if action == "refine" and state.quiz.get(quiz_type):
        selected_topics = [t.strip() for t in (manual_topics or []) if str(t).strip()] if plan_mode == "manual_topics" else []
        existing_all = state.quiz[quiz_type]
        if selected_topics:
            existing = [
                q for q in existing_all
                if (q.get("topic") or "").strip() in selected_topics
            ]
            if not existing:
                return jsonify(ok=False, error="No existing questions match selected manual topics.", questions=[], thinking=[])
        else:
            existing = existing_all
        thinking_cb({
            "agent": "QuizRefinerAgent", "step": "refine_start", "verdict": "PASS",
            "phase": "PLAN",
            "reason": f"Refining {len(existing)} existing {quiz_type.upper()} questions with new instructions...",
        })
        thinking_cb({
            "agent": "QuizRefinerAgent", "step": "phase_start", "phase": "GENERATE", "verdict": "PASS",
            "reason": f"Generating refined set for {len(existing)} selected question(s)…",
            "data": {"total_slots": len(existing)},
        })
        weak = difficulty_adapter.get_weak_topics()
        weak_str = ", ".join(weak.keys()) if weak else "none"
        topics_line = ", ".join(selected_topics) if selected_topics else "auto"
        schema = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
        prompt = (
            f"Revise this {quiz_type.upper()} quiz to better follow the student's instructions.\n"
            f"Mode: {mode}  Difficulty: {difficulty.upper()}\n"
            f"Topic scope: {topics_line}\n"
            f"Weak topics to emphasize: {weak_str}\n"
            f"Student instructions: {user_instr or 'None'}\n\n"
            f"EXISTING QUIZ:\n{json.dumps(existing, ensure_ascii=False)}\n\n"
            f"CONTENT:\n{content[:6000]}\n\nSUMMARY:\n{summary[:1500]}\n\n"
            f"STRICT JSON: {schema}"
        )
        thinking_cb({"agent": "QuizRefinerAgent", "step": "llm_prompt", "phase": "GENERATE", "verdict": "PASS",
                     "reason": "Prompt sent to LLM for refine generation",
                     "data": {"prompt_preview": prompt[:1200]}})
        thinking_cb({"agent": "QuizRefinerAgent", "step": "generating", "verdict": "PASS",
                     "reason": "Sending revision prompt to LLM..."})
        data_out = call_ollama_json_quiz(prompt, key="questions")
        qs = data_out.get("questions", []) if isinstance(data_out, dict) else []
        # validate the refined output
        passed, failed = QuizValidatorAgent().validate_batch(
            qs, quiz_type, content[:800], difficulty, user_instr, thinking_cb=thinking_cb
        )
        if failed:
            fixed = QuizRefinerAgent().refine_batch(
                failed, quiz_type, content[:800], difficulty, user_instr, thinking_cb=thinking_cb
            )
            qs = passed + fixed
        else:
            qs = passed
        fill_missing_timestamps(qs, segments=ctx.segment_list, board_entries=ctx.board_entries)
        state.quiz[quiz_type] = qs
        thinking_cb({"agent": "QuizRefinerAgent", "step": "final", "verdict": "DONE",
                     "reason": f"Refined quiz ready -- {len(qs)} questions"})

    # -- Adaptive mode --
    elif mode == "adaptive":
        ctx._mode = "adaptive"
        agent = AdaptiveQuizAgent()
        agent.run(ctx)
        thinking.extend(agent._thinking_log)
        fill_missing_timestamps(
            state.quiz.get(quiz_type, []),
            segments=ctx.segment_list, board_entries=ctx.board_entries,
        )

    # -- Manual mode --
    elif mode == "manual":
        agent_map = {"easy": EasyQuizAgent, "medium": MediumQuizAgent, "hard": HardQuizAgent}
        agent = agent_map.get(difficulty, MediumQuizAgent)()
        ctx._mode       = "manual"
        ctx._agent_name = agent.name
        thinking_cb({"agent": agent.name, "step": "start", "verdict": "PASS",
                     "reason": f"Manual {difficulty.upper()} mode -- {num} {quiz_type.upper()} questions"})
        if user_instr:
            thinking_cb({"agent": agent.name, "step": "user_instructions", "verdict": "PASS",
                         "reason": f"Student instructions: {user_instr}"})

        max_attempts = 3
        final = _run_quiz_agentic_loop(
            ctx, quiz_type, content, summary, user_instr,
            target_difficulty=difficulty, thinking_cb=thinking_cb,
            max_attempts_per_q=max_attempts,
        )
        for q in final:
            q["difficulty"] = difficulty
        fill_missing_timestamps(final, segments=ctx.segment_list, board_entries=ctx.board_entries)
        state.quiz[quiz_type] = final

    else:
        return jsonify(ok=False, error=f"Unknown mode: {mode}", questions=[], thinking=[])

    state._quiz_agent_thinking = thinking
    final_qs = state.quiz.get(quiz_type, [])
    return jsonify(ok=True, questions=final_qs, total=len(final_qs),
                   thinking=thinking, mode=mode,
                   difficulty=difficulty if mode == "manual" else "adaptive")

# ── Study endpoints (flashcards, hints, chat) ─────────────────



@app.route('/flashcards')
def flashcards_route():
    cards = state.quiz.get("flash", [])
    return jsonify(cards=cards, total=len(cards))


@app.route('/flashcards/generate', methods=['POST'])
def flashcards_generate_route():
    """Generate flashcards on-demand from the Flashcards section."""
    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No processed content yet. Upload and process first."), 400

    if state.video_path:
        cards = generate_flashcards(
            summary,
            board_text=state.status.get("board_text", ""),
            board_entries=load_board_entries_from_disk(),
        )
    else:
        cards = generate_flashcards(summary, doc_text=summary)

    # Deduplicate by front text
    seen, unique = set(), []
    for c in cards or []:
        front = (c.get("front") or "").strip().lower()
        if front and front not in seen:
            seen.add(front)
            unique.append(c)
    state.quiz["flash"] = unique
    return jsonify(ok=True, cards=unique, total=len(unique))



@app.route('/exam_hints')
def exam_hints_route():
    return jsonify(**state.quiz.get("hints", {}))


@app.route('/quiz/topics')
def quiz_topics_route():
    """Return topic options for manual topic plan mode."""
    topics = []

    # 1) AI-ranked exam topics
    hints = state.quiz.get("hints", {}) or {}
    ai = hints.get("ai_analysis", {}) if isinstance(hints, dict) else {}
    ranked = ai.get("ai_important_topics", []) if isinstance(ai, dict) else []
    for item in ranked[:20]:
        if isinstance(item, dict):
            t = (item.get("topic") or "").strip()
            if t:
                topics.append(t)

    # 2) performance-tracked topics
    for t in list(difficulty_adapter.get_all_stats().keys())[:20]:
        if t and t not in topics:
            topics.append(t)

    # 3) weak topics
    for t in list(difficulty_adapter.get_weak_topics().keys())[:20]:
        if t and t not in topics:
            topics.append(t)

    # 4) sensible fallback
    if not topics:
        topics = ["Core Concepts", "Key Principles", "Applications", "Problem Solving"]

    return jsonify(topics=topics[:30], total=len(topics[:30]))



@app.route('/suggested_questions')
def suggested_questions_route():
    return jsonify(questions=state.suggested_questions)



GREETING_WORDS = {"hi","hello","hey","sup","yo","thanks","thank you","bye","ok","okay","good morning","good evening","hola"}



# ── Socratic tutor prompts ─────────────────────────────────────

SOCRATIC_SYSTEM = """You are a Socratic AI lecture tutor. Your role is to help students \
DISCOVER answers through guided questioning, not by giving answers directly.

STRICT RULES:
1. On the FIRST exchange about a topic: NEVER give the answer directly.
   Instead, ask 1-2 short probing questions that activate prior knowledge.
2. When a student attempts an answer: acknowledge what is correct, gently correct \
   what is wrong, then ask a follow-up question that moves them closer to the full answer.
3. Give the COMPLETE answer only when:
   - The student has made 2 or more genuine attempts, OR
   - The student says something like "just tell me", "I give up", "I don't know", \
     "tell me the answer", or "what is it".
4. Every Socratic response must end with a question mark — you are always probing.
5. Keep probing questions short (1-2 sentences) and focus on ONE concept at a time.
6. Reference lecture content naturally: "The lecture mentioned at [timestamp]…"
7. Be warm and encouraging — never make the student feel bad for not knowing.

EXAMPLE:
Student: "What is gradient descent?"
BAD response: "Gradient descent is an optimization algorithm that…"
GOOD response: "Great question! Let's think about it together. Imagine you're \
standing on a hilly landscape and want to reach the lowest point — what strategy \
would you use to get there?"
"""

DIRECT_SYSTEM = "You are a helpful lecture tutor. Answer clearly and concisely. Cite timestamps when possible."

# phrases that signal the student wants the answer directly
_GIVE_UP_PHRASES = (
    "just tell me", "tell me the answer", "i give up", "i don't know",
    "what is it", "what's the answer", "i have no idea", "please just answer",
    "stop asking", "skip the questions", "idk", "no idea", "بলো", "বলো",
)

def _student_gave_up(text):
    t = text.lower().strip()
    return any(ph in t for ph in _GIVE_UP_PHRASES)

def _build_history_block(history):
    """Convert [{role, content}] list into a readable conversation string."""
    if not history:
        return ""
    lines = []
    for turn in history[-6:]:   # last 6 turns = 3 exchanges
        role    = "Student" if turn.get("role") == "user" else "Tutor"
        content = str(turn.get("content", "")).strip()[:400]
        lines.append(f"{role}: {content}")
    return "=== RECENT CONVERSATION ===\n" + "\n".join(lines)

def _count_student_attempts(history):
    """Count how many times the student has already responded (not first-time)."""
    return sum(1 for t in history if t.get("role") == "user")



@app.route('/chat', methods=['POST'])
def chat_route():
    data    = request.get_json(silent=True) or {}
    q       = data.get("question", "").strip()
    history = data.get("history",  [])        # [{role:"user"|"assistant", content:"..."}]
    mode    = data.get("mode",     "socratic") # "socratic" | "direct"

    if not q:
        return jsonify(answer="Please ask a question.", timestamps=[], mode_used=mode, socratic=False)

    # ── greeting shortcut ──────────────────────────────────────
    if q.lower().strip("!?., ") in GREETING_WORDS or (len(q.split()) <= 2 and not any(c.isdigit() for c in q)):
        answer = call_ollama(
            f"The student said: \"{q}\"\nRespond with a brief friendly greeting. "
            "You are their Socratic lecture tutor. Keep it to 1-2 sentences.",
            system="You are a friendly lecture tutor.",
        )
        return jsonify(answer=answer, timestamps=[], mode_used="direct", socratic=False)

    # ── RAG retrieval + grounding ──────────────────────────────
    ctx, tss = rag.query(q, top_k=5)
    summary_snip = (state.status.get("summary") or "")[:2000]
    extra = []
    if summary_snip:
        extra.append(f"=== LECTURE SUMMARY ===\n{summary_snip}")
    if state.video_path:
        be  = load_board_entries_from_disk()
        bt  = state.status.get("board_text", "")
        vis = format_visual_lecture_text(bt, be, 2800)
        if vis:
            extra.append(f"=== SLIDES + BOARDS (OCR) ===\n{vis}")
    grounding = "\n\n".join(extra)

    # ── clip duration estimation ───────────────────────────────
    for i, ts in enumerate(tss):
        chunk_idx = None
        try:
            chunk_idx = rag.timestamps.index(ts)
        except (ValueError, AttributeError):
            pass
        words = len(rag.chunks[chunk_idx].split()) if (chunk_idx is not None and chunk_idx < len(rag.chunks)) else 50
        tss[i] = {**ts, "clip_dur": max(10, min(60, int(words / 2.5)))}

    # ── decide whether to be Socratic or direct ────────────────
    attempts       = _count_student_attempts(history)
    gave_up        = _student_gave_up(q)
    use_socratic   = (mode == "socratic") and (not gave_up) and (attempts < 2)
    history_block  = _build_history_block(history)

    if use_socratic:
        prompt = (
            "Use the lecture material below to guide a Socratic dialogue.\n"
            "Do NOT answer directly — ask probing questions that help the student think.\n\n"
            f"{grounding}\n\n"
            f"=== RELEVANT SNIPPETS ===\n{ctx}\n\n"
            f"{history_block}\n\n"
            f"Student's latest message: {q}\n\n"
            "Your response (Socratic — end with a question):"
        )
        system = SOCRATIC_SYSTEM
    else:
        # direct answer: student gave up, or already tried 2+ times, or mode is direct
        reveal_note = ""
        if gave_up:
            reveal_note = "The student has indicated they need the answer directly. Give a clear, complete explanation.\n\n"
        elif attempts >= 2:
            reveal_note = f"The student has attempted this {attempts} times. Now give a full, clear explanation.\n\n"

        prompt = (
            f"{reveal_note}"
            "Answer using the lecture material below. Be clear and cite timestamps.\n\n"
            f"{grounding}\n\n=== RELEVANT SNIPPETS ===\n{ctx}\n\n"
            f"{history_block}\n\n"
            f"QUESTION: {q}"
        )
        system = DIRECT_SYSTEM

    answer = call_ollama(prompt, system=system)
    return jsonify(
        answer    = answer,
        timestamps= tss,
        mode_used = "socratic" if use_socratic else "direct",
        socratic  = use_socratic,
    )

# ── Lesson Planner ────────────────────────────────────────────



@app.route('/lesson_plan', methods=['POST'])
def lesson_plan_route():
    data  = request.get_json(silent=True) or {}
    topic = data.get("topic", "")
    level = data.get("level", "undergraduate")
    dur   = data.get("duration", 45)
    if not topic:
        return jsonify(error="No topic provided.")
    prompt = (
        f"Create a complete {dur}-minute lesson plan about '{topic}' for {level} students.\n"
        "Include: Learning Objectives, Prerequisites, Introduction (5 min), "
        "Main Content with timed sub-topics, Activities/Exercises, "
        "Summary & Key Takeaways, Assessment Questions, and Homework.\n"
        "Format in clear Markdown with headers."
    )
    plan = call_ollama(prompt, system="You are an expert curriculum designer.")
    return jsonify(plan=plan, topic=topic, level=level, duration=dur)

# ── Sessions (list / load / delete past processed files) ──────



@app.route('/sessions')
def list_sessions_route():
    return jsonify(sessions=AppState.list_sessions())



@app.route('/sessions/load', methods=['POST'])
def load_session_route():
    data = request.get_json(silent=True) or {}
    sid  = data.get("session_id", "")
    if not sid:
        return jsonify(ok=False, error="No session_id provided.")
    ok = state.load_session(sid)
    if not ok:
        return jsonify(ok=False, error=f"Session '{sid}' not found.")

    # Rebuild RAG from restored transcript/board data
    rag_entries = []
    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    board_path = os.path.join(OUTPUT_DIR, "board_entries.json")
    if os.path.exists(seg_path):
        with open(seg_path, "r", encoding="utf-8") as f:
            rag_entries += json.load(f)
    if os.path.exists(board_path):
        with open(board_path, "r", encoding="utf-8") as f:
            for e in json.load(f):
                rag_entries.append({
                    "text": e["text"], "start_str": e["timestamp_str"],
                    "tag": e.get("tag", ""),
                })
    if rag_entries:
        build_rag(rag_entries, ts_key="start_str", src_key="tag")

    return jsonify(ok=True, session_id=sid, source_file=data.get("source_file", ""))



@app.route('/sessions/delete', methods=['POST'])
def delete_session_route():
    data = request.get_json(silent=True) or {}
    sid  = data.get("session_id", "")
    if not sid:
        return jsonify(ok=False, error="No session_id provided.")
    folder = os.path.join(SESSIONS_DIR, sid)
    if os.path.isdir(folder):
        shutil.rmtree(folder)
        return jsonify(ok=True, deleted=sid)
    return jsonify(ok=False, error="Session not found."), 404



@app.route('/sessions/current')
def current_session_route():
    return jsonify(
        session_id=state.session_id,
        source_file=Path(state.video_path or state.doc_path or "").name,
        has_data=state.status.get("state") == "done",
    )

# ── Agentic endpoints ─────────────────────────────────────────



@app.route('/agent_log')
def agent_log_route():
    return jsonify(log=state.agent_log, total=len(state.agent_log))



@app.route('/quiz/regenerate', methods=['POST'])
def quiz_regenerate_route():
    data      = request.get_json(silent=True) or {}
    quiz_type = data.get("type", "mcq")
    num       = min(int(data.get("count", 5)), 20)
    qs        = difficulty_adapter.generate_targeted(quiz_type=quiz_type, num=num)
    if not qs:
        weak = difficulty_adapter.get_weak_topics()
        if not weak:
            return jsonify(questions=[], message="No weak topics detected yet — take a quiz first.")
        return jsonify(questions=[], message="No summary available to generate from.")
    return jsonify(questions=qs, total=len(qs), weak_topics=difficulty_adapter.get_weak_topics())



@app.route('/quiz/regenerate_one', methods=['POST'])
def quiz_regenerate_one_route():
    """Regenerate a single question in-place using the per-question agentic loop.
    Body: {quiz_type, index, instructions?, difficulty?}.
    Streams the same Generate -> Validate -> Retry events the bulk loop emits."""
    data        = request.get_json(silent=True) or {}
    quiz_type   = data.get("quiz_type", "mcq")
    try:
        idx = int(data.get("index", 0))
    except (TypeError, ValueError):
        return jsonify(ok=False, error="Invalid index", thinking=[])
    user_instr  = (data.get("instructions") or "").strip()
    requested_diff = data.get("difficulty")

    qs = state.quiz.get(quiz_type, [])
    if not qs or idx < 0 or idx >= len(qs):
        return jsonify(ok=False, error="Invalid question index", thinking=[])

    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No content available", thinking=[])

    if state.video_path:
        be = load_board_entries_from_disk()
        bt = state.status.get("board_text", "")
        tr = state.status.get("transcript", "")
        content = pack_video_sources(bt, tr, be, 7500)
    else:
        content = summary[:7500]

    orig  = qs[idx]
    topic = orig.get("topic") or "General"
    target_difficulty = requested_diff or orig.get("difficulty") or "medium"

    thinking = []
    def push(entry): thinking.append(entry)

    push({
        "agent":       "Orchestrator",
        "step":        "regen_start",
        "phase":       "GENERATE",
        "verdict":     "PASS",
        "question_id": f"q{idx+1}",
        "reason":      f"Regenerating Q{idx+1} ({topic}, {target_difficulty}) on demand",
    })

    # exclude the question being replaced from the dedup set
    existing = [q for i, q in enumerate(qs) if i != idx]

    new_q = _generate_one_question_with_retry(
        quiz_type=quiz_type, topic=topic, difficulty=target_difficulty,
        content_snip=content, summary=summary,
        user_instr=user_instr,
        existing_questions=existing,
        question_idx=idx,
        max_attempts=3,
        thinking_cb=push,
        agent_name="QuizRegenerator",
    )

    if not new_q:
        return jsonify(ok=False, error="Failed to generate replacement",
                       thinking=thinking)

    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    segs = []
    if os.path.exists(seg_path):
        try:
            with open(seg_path, "r", encoding="utf-8") as f:
                segs = json.load(f)
        except Exception:
            segs = []
    be = load_board_entries_from_disk()
    fill_missing_timestamps([new_q], segments=segs, board_entries=be)

    qs[idx] = new_q
    state.quiz[quiz_type] = qs
    state._quiz_agent_thinking = thinking
    return jsonify(ok=True, question=new_q, index=idx,
                   thinking=thinking, total=len(qs))



@app.route('/quiz/thinking')
def quiz_thinking_route():
    """Return the thinking log from the last quiz generation."""
    return jsonify(thinking=state._quiz_agent_thinking,
                   total=len(state._quiz_agent_thinking))



def _new_job_id():
    short_hash = hashlib.md5(f"{time.time()}_{os.urandom(8)}".encode()).hexdigest()[:10]
    return f"job_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{short_hash}"



def _job_push_thinking(job_id, entry):
    with state.quiz_generation_lock:
        job = state.quiz_generation_jobs.get(job_id)
        if not job:
            return
        job["thinking"].append(entry)



@app.route('/quiz/generate_async', methods=['POST'])
def quiz_generate_async_route():
    """Kick off quiz generation in a background thread.
    Frontend polls /quiz/generate_job?job_id=... for live thinking entries.
    Uses the full Plan → Generate → Validate → Refine agentic loop."""
    data       = request.get_json(silent=True) or {}
    mode       = data.get("mode", "adaptive")
    difficulty = data.get("difficulty", "medium")
    quiz_type  = data.get("quiz_type", "mcq")
    action     = data.get("action", "new")
    num        = min(int(data.get("count", 10)), 20)
    user_instr = data.get("instructions", "").strip()
    plan_mode  = data.get("plan_mode", "auto")
    manual_topics = data.get("manual_topics", []) or []

    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No content processed yet — upload a file first.", job_id=None)

    job_id = _new_job_id()
    with state.quiz_generation_lock:
        state.quiz_generation_jobs[job_id] = {
            "status": "running", "thinking": [], "result": None, "error": "",
        }

    # Bump the global thinking session so the UI also sees this run via /status
    state._quiz_thinking_session += 1
    state._quiz_thinking_active   = True
    state._quiz_agent_thinking    = []

    def worker():
        try:
            def push(entry):
                _job_push_thinking(job_id, entry)
                state._quiz_agent_thinking.append(entry)

            push({"agent": "Orchestrator", "step": "start", "verdict": "PASS",
                  "reason": f"{mode} mode started for {num} {quiz_type.upper()} questions — {difficulty.upper()} difficulty"})

            ctx = type('Ctx', (), {
                'content_type':       'video' if state.video_path else 'document',
                'board_text':         state.status.get("board_text", ""),
                'transcript':         state.status.get("transcript", ""),
                'board_entries':      load_board_entries_from_disk(),
                'doc_text':           state.status.get("summary", ""),
                'summary':            summary,
                'agent_log':          [],
                '_quiz_type':         quiz_type,
                '_quiz_count':        num,
                '_user_instructions': user_instr,
                '_plan_mode':         plan_mode,
                '_manual_topics':     manual_topics,
                '_manual_difficulty': difficulty,
                '_mode':              mode,
                '_quiz_plan':         None,
                '_agent_name':        '',
                'segment_list':       [],
                'file_path':          state.video_path or state.doc_path or "",
                'source_name':        Path(state.video_path or state.doc_path or "unknown").name,
                'thinking_cb':        push,
            })()

            seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
            if os.path.exists(seg_path):
                try:
                    with open(seg_path, "r", encoding="utf-8") as f:
                        ctx.segment_list = json.load(f)
                except Exception:
                    ctx.segment_list = []

            if state.doc_path:
                doc_entries_path = os.path.join(OUTPUT_DIR, "doc_text.txt")
                if os.path.exists(doc_entries_path):
                    with open(doc_entries_path, "r", encoding="utf-8") as f:
                        ctx.doc_text = f.read()

            content = (
                pack_video_sources(ctx.board_text, ctx.transcript,
                                   ctx.board_entries or [], 7500)
                if ctx.content_type == "video"
                else (ctx.doc_text or "")[:7500]
            )

            # Refine mode
            if action == "refine" and state.quiz.get(quiz_type):
                selected_topics = [t.strip() for t in (manual_topics or []) if str(t).strip()] if plan_mode == "manual_topics" else []
                existing_all = state.quiz[quiz_type]
                if selected_topics:
                    existing = [
                        q for q in existing_all
                        if (q.get("topic") or "").strip() in selected_topics
                    ]
                    if not existing:
                        raise ValueError("No existing questions match selected manual topics.")
                else:
                    existing = existing_all
                push({"agent": "QuizRefinerAgent", "step": "refine_start", "verdict": "PASS",
                      "phase": "PLAN",
                      "reason": f"Refining {len(existing)} questions with new instructions…"})
                push({"agent": "QuizRefinerAgent", "step": "phase_start", "phase": "GENERATE", "verdict": "PASS",
                      "reason": f"Generating refined set for {len(existing)} selected question(s)…",
                      "data": {"total_slots": len(existing)}})
                weak    = difficulty_adapter.get_weak_topics()
                schema  = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
                topics_line = ", ".join(selected_topics) if selected_topics else "auto"
                prompt  = (
                    f"Revise this {quiz_type.upper()} quiz per student instructions.\n"
                    f"Mode: {mode}  Difficulty: {difficulty.upper()}\n"
                    f"Topic scope: {topics_line}\n"
                    f"Weak topics: {', '.join(weak.keys()) if weak else 'none'}\n"
                    f"Instructions: {user_instr or 'None'}\n\n"
                    f"EXISTING QUIZ:\n{json.dumps(existing, ensure_ascii=False)}\n\n"
                    f"CONTENT:\n{content[:6000]}\n\nSUMMARY:\n{summary[:1500]}\n\n"
                    f"STRICT JSON: {schema}"
                )
                push({"agent": "QuizRefinerAgent", "step": "llm_prompt", "phase": "GENERATE", "verdict": "PASS",
                      "reason": "Prompt sent to LLM for refine generation",
                      "data": {"prompt_preview": prompt[:1200]}})
                push({"agent": "QuizRefinerAgent", "step": "generating", "verdict": "PASS",
                      "reason": "Sending revision prompt to LLM…"})
                data_out = call_ollama_json_quiz(prompt, key="questions")
                qs = data_out.get("questions", []) if isinstance(data_out, dict) else []
                passed, failed = QuizValidatorAgent().validate_batch(
                    qs, quiz_type, content[:800], difficulty, user_instr, thinking_cb=push
                )
                if failed:
                    fixed = QuizRefinerAgent().refine_batch(
                        failed, quiz_type, content[:800], difficulty, user_instr, thinking_cb=push
                    )
                    qs = passed + fixed
                else:
                    qs = passed
                fill_missing_timestamps(qs, segments=ctx.segment_list, board_entries=ctx.board_entries)
                state.quiz[quiz_type] = qs
                push({"agent": "QuizRefinerAgent", "step": "final", "verdict": "DONE",
                      "reason": f"Refined quiz ready — {len(qs)} questions"})

            elif mode == "adaptive":
                ctx._mode = "adaptive"
                agent = AdaptiveQuizAgent()
                agent.run(ctx)   # internally calls _run_quiz_agentic_loop via thinking_cb=push
                fill_missing_timestamps(
                    state.quiz.get(quiz_type, []),
                    segments=ctx.segment_list, board_entries=ctx.board_entries,
                )

            elif mode == "manual":
                agent_map = {"easy": EasyQuizAgent, "medium": MediumQuizAgent,
                             "hard": HardQuizAgent}
                agent = agent_map.get(difficulty, MediumQuizAgent)()
                ctx._mode       = "manual"
                ctx._agent_name = agent.name
                push({"agent": agent.name, "step": "start", "verdict": "PASS",
                      "reason": f"Manual {difficulty.upper()} mode — {num} {quiz_type.upper()} questions"})
                if user_instr:
                    push({"agent": agent.name, "step": "user_instructions", "verdict": "PASS",
                          "reason": f"Student instructions: {user_instr}"})
                max_attempts = 3
                final = _run_quiz_agentic_loop(
                    ctx, quiz_type, content, summary, user_instr,
                    target_difficulty=difficulty, thinking_cb=push,
                    max_attempts_per_q=max_attempts,
                )
                for q in final:
                    q["difficulty"] = difficulty
                fill_missing_timestamps(final, segments=ctx.segment_list, board_entries=ctx.board_entries)
                state.quiz[quiz_type] = final

            else:
                raise ValueError(f"Unknown mode: {mode}")

            final_qs = state.quiz.get(quiz_type, [])
            diff_label = difficulty if mode == "manual" else "adaptive"
            _append_generated_quiz_run("async", quiz_type, mode, diff_label, final_qs)
            with state.quiz_generation_lock:
                job = state.quiz_generation_jobs.get(job_id)
                if job:
                    job["status"] = "done"
                    job["result"] = {
                        "mode":       mode,
                        "difficulty": difficulty if mode == "manual" else "adaptive",
                        "total":      len(final_qs),
                        "questions":  final_qs,
                    }
        except Exception as e:
            with state.quiz_generation_lock:
                job = state.quiz_generation_jobs.get(job_id)
                if job:
                    job["status"] = "error"
                    job["error"]  = str(e)
        finally:
            state._quiz_thinking_active = False

    threading.Thread(target=worker, daemon=True).start()
    return jsonify(ok=True, job_id=job_id)



@app.route('/quiz/generate_job')
def quiz_generate_job_route():
    job_id = request.args.get("job_id", "")
    if not job_id:
        return jsonify(ok=False, error="Missing job_id"), 400
    with state.quiz_generation_lock:
        job = state.quiz_generation_jobs.get(job_id)
        if not job:
            return jsonify(ok=False, error="Job not found"), 404
        return jsonify(
            ok=True,
            status=job.get("status", ""),
            thinking=job.get("thinking", []),
            result=job.get("result", None),
            error=job.get("error", ""),
        )



@app.route('/performance')
def performance_route():
    stats = difficulty_adapter.get_all_stats()
    weak  = difficulty_adapter.get_weak_topics()
    return jsonify(stats=stats, weak_topics=weak, total_topics=len(stats))

# ── Video streaming (for in-UI playback at source timestamps) ─



def _stream_video_path(file_path, mimetype="video/mp4"):
    """Serve MP4 with HTTP Range so browsers can seek without downloading the whole file."""
    if not os.path.exists(file_path):
        return jsonify(error="File not found."), 404
    file_size = os.path.getsize(file_path)
    range_header = request.headers.get("Range")
    if not range_header:
        def gen_full():
            with open(file_path, "rb") as f:
                while True:
                    chunk = f.read(1024 * 1024)
                    if not chunk:
                        break
                    yield chunk

        return Response(
            gen_full(),
            status=200,
            mimetype=mimetype,
            headers={
                "Content-Length": str(file_size),
                "Accept-Ranges": "bytes",
                "Cache-Control": "no-cache",
            },
        )
    try:
        byte_range = range_header.replace("bytes=", "").split("-")
        start = int(byte_range[0]) if byte_range[0] else 0
        end = int(byte_range[1]) if byte_range[1] else file_size - 1
    except (ValueError, IndexError):
        return Response(status=416)
    end = min(max(start, end), file_size - 1)
    start = min(max(0, start), file_size - 1)
    if start > end:
        return Response(status=416)
    length = end - start + 1

    def gen_chunk():
        with open(file_path, "rb") as f:
            f.seek(start)
            remaining = length
            while remaining > 0:
                chunk = f.read(min(1024 * 1024, remaining))
                if not chunk:
                    break
                remaining -= len(chunk)
                yield chunk

    return Response(
        gen_chunk(),
        status=206,
        mimetype=mimetype,
        headers={
            "Content-Range": f"bytes {start}-{end}/{file_size}",
            "Accept-Ranges": "bytes",
            "Content-Length": str(length),
            "Cache-Control": "no-cache",
        },
    )



@app.route('/stream/video')
def stream_video_route():
    path = state.video_path
    if not path or not os.path.exists(path):
        return jsonify(error="No video available."), 404
    return _stream_video_path(path)



@app.route('/clip/video')
def clip_video_route():
    """Extract a short clip around a timestamp. ?t=MM:SS&dur=15"""
    path = state.video_path
    if not path or not os.path.exists(path):
        return jsonify(error="No video available."), 404

    ts_str = request.args.get("t", "00:00")
    dur    = min(int(request.args.get("dur", 20)), 90)

    # parse MM:SS or HH:MM:SS to seconds
    parts = ts_str.split(":")
    if len(parts) == 3:
        start = int(parts[0]) * 3600 + int(parts[1]) * 60 + int(parts[2])
    elif len(parts) == 2:
        start = int(parts[0]) * 60 + int(parts[1])
    else:
        start = int(float(ts_str))

    # Cut should start at the exact timestamp shown in the UI.
    start = max(0, start)
    # Suffix bumps cache when encode settings change (browser-safe yuv420p + main).
    clip_name = f"clip_{start}_{dur}_web.mp4"
    clip_path = os.path.join(CLIPS_DIR, clip_name)

    if not os.path.exists(clip_path):
        cmd = [
            "ffmpeg", "-y", "-ss", str(start), "-i", path,
            "-t", str(dur),
            "-c:v", "libx264", "-preset", "veryfast",
            "-profile:v", "main", "-pix_fmt", "yuv420p",
            "-c:a", "aac", "-b:a", "128k", "-ar", "44100", "-ac", "2",
            "-movflags", "+faststart",
            clip_path,
        ]
        try:
            subprocess.run(cmd, capture_output=True, timeout=120)
        except Exception as e:
            return jsonify(error=f"Clip extraction failed: {e}"), 500

    if not os.path.exists(clip_path):
        return jsonify(error="Clip extraction produced no output."), 500
    return _stream_video_path(clip_path)

# ── Downloads ─────────────────────────────────────────────────



DOWNLOAD_MAP = {
    "summary":    f"{OUTPUT_DIR}/summary.txt",
    "transcript": f"{OUTPUT_DIR}/transcript.txt",
    "board":      f"{OUTPUT_DIR}/board_text.txt",
    "timeline":   f"{OUTPUT_DIR}/lecture_timeline.txt",
}



@app.route('/download/<file_type>')
def download_route(file_type):
    path = DOWNLOAD_MAP.get(file_type)
    if path and os.path.exists(path):
        return send_file(path, as_attachment=True)
    return jsonify(error="File not found."), 404

# ── Print-friendly quiz export ────────────────────────────────



@app.route('/export/quiz')
def export_quiz_route():
    """Returns a self-contained printable HTML page of all quizzes."""
    q = state.quiz
    rows = []
    for i, m in enumerate(q.get("mcq", [])):
        opts = "".join(f"<li><b>{k}.</b> {v}</li>" for k, v in m.get("options", {}).items())
        rows.append(f"<div class='eq'><p><b>Q{i+1}.</b> {m.get('question','')}</p><ul>{opts}</ul></div>")
    for i, t in enumerate(q.get("tf", [])):
        rows.append(f"<div class='eq'><p><b>T/F {i+1}.</b> {t.get('statement','')}</p></div>")
    for i, f in enumerate(q.get("fill", [])):
        rows.append(f"<div class='eq'><p><b>Fill {i+1}.</b> {f.get('question','')}</p></div>")
    for i, s in enumerate(q.get("short", [])):
        rows.append(f"<div class='eq'><p><b>Short {i+1}.</b> ({s.get('marks',4)} marks) {s.get('question','')}</p>"
                     f"<div class='ans-space'></div></div>")
    body = "\n".join(rows) or "<p>No quiz data available.</p>"

    # answer key
    keys = []
    for i, m in enumerate(q.get("mcq", [])):
        keys.append(f"Q{i+1}: {m.get('correct_answer','')}")
    for i, t in enumerate(q.get("tf", [])):
        keys.append(f"T/F {i+1}: {t.get('answer','')}")
    for i, f in enumerate(q.get("fill", [])):
        keys.append(f"Fill {i+1}: {f.get('answer','')}")
    key_html = " &nbsp;|&nbsp; ".join(keys)

    html = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<title>LectureForge — Printable Quiz</title>
<style>
body{{font-family:Georgia,serif;max-width:750px;margin:40px auto;color:#1a1a2e;line-height:1.7}}
h1{{text-align:center;border-bottom:2px solid #333;padding-bottom:8px}}
.eq{{margin-bottom:18px;page-break-inside:avoid}}
.eq p{{margin:0 0 4px}} .eq ul{{margin:4px 0;padding-left:24px}}
.ans-space{{border-bottom:1px dashed #999;height:60px}}
.key{{margin-top:40px;padding-top:16px;border-top:2px solid #333;font-size:12px;color:#555}}
@media print{{.no-print{{display:none}}}}
</style></head><body>
<button class="no-print" onclick="window.print()" style="float:right;padding:8px 16px;cursor:pointer">🖨 Print / Save PDF</button>
<h1>LectureForge Quiz</h1>
{body}
<div class="key"><b>Answer Key:</b><br>{key_html}</div>
</body></html>"""
    return html, 200, {'Content-Type': 'text/html; charset=utf-8'}



print("✅ Flask API defined.")

# ══════════════════════════════════════════════════════════════
# LAUNCH  (Flask + ngrok)
# ══════════════════════════════════════════════════════════════



from pyngrok import ngrok, conf



conf.get_default().auth_token = NGROK_AUTH_TOKEN
ngrok.kill()



def _run_flask():
    app.run(host='0.0.0.0', port=5000, debug=False, use_reloader=False)



threading.Thread(target=_run_flask, daemon=True).start()
time.sleep(2)



tunnel     = ngrok.connect(5000, "http")
public_url = tunnel.public_url



print(f"""
{'='*60}
  🎓 LectureForge v1.1 LIVE  (backend-only mode)
  ➜  {public_url}
{'='*60}

  1. Copy the URL above
  2. Open quizforge_ui.html in your browser
  3. Paste the URL → Connect
  4. Upload your lecture → Process Content

  Model : {state.active_model}
  CORS  : enabled (local UI supported)
""")



import subprocess, time

# Start Ollama in background
process = subprocess.Popen(
    ["ollama", "serve"],
    stdout=subprocess.PIPE,
    stderr=subprocess.PIPE,
)

time.sleep(5)

print("🚀 Attempted to start Ollama")



import requests

def check_ollama():
    try:
        r = requests.get("http://127.0.0.1:11434/api/tags", timeout=3)
        if r.status_code == 200:
            print("✅ Ollama is RUNNING")
            print("📦 Models:", r.json())
        else:
            print("⚠️ Ollama responded but something is off:", r.status_code)
    except Exception as e:
        print("❌ Ollama is NOT running")
        print("Error:", e)

check_ollama()



import time

while True:
    print("Still running...")
    time.sleep(60)  # runs every 60 seconds



