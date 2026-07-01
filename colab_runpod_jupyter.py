# ╔══════════════════════════════════════════════════════════════╗
# ║      LectureForge v1.1  —  Backend (Jupyter / RunPod)          ║
# ║  Run as notebook CELLS on a CUDA PyTorch RunPod pod (GPU)     ║
# ║  UI → open LectureAssis.html locally on your PC               ║
# ╚══════════════════════════════════════════════════════════════╝
#
# Cells are delimited by  "# %%"  — paste each block into its own notebook cell.
#   Cell 1 = install (run once)
#   Cell 2 = backend (starts Ollama, pulls models, serves Flask in a thread)
#   Cell 3 = optional: run the bulk evaluation directly from the notebook
# Expose HTTP port 5000 in the pod, then open the proxy URL in LectureAssis.html.

# %% [Cell 1 — install once]
# Use a CUDA PyTorch pod so torch is already present — do NOT reinstall torch.
!pip install -q easyocr faster-whisper sentence-transformers faiss-cpu pillow numpy werkzeug flask flask-cors python-docx python-pptx PyMuPDF opencv-python-headless matplotlib requests
!apt-get update -qq && apt-get install -y -qq ffmpeg curl zstd
!curl -fsSL https://ollama.com/install.sh | sh
# Optional — equation→LaTeX extraction for math PDFs (needs GPU). Enable by also
# setting the env var LF_MATH_PDF=1 in Cell 2's CONFIG:
# !pip install -q nougat-ocr
print("✅ Cell 1 done — now run Cell 2")


# %% [Cell 2 — backend: run after Cell 1]
# ══════════════════════════════════════════════════════════════
# CONFIG
# ══════════════════════════════════════════════════════════════



OLLAMA_MODEL     = "gemma3:12b"              # primary model — summaries, chat, grading
OLLAMA_FALLBACK  = "gemma3:4b"               # fallback if primary fails
QUIZ_MODEL       = "gemma3:4b"               # fast model for quiz generation
QUIZ_CTX         = 16384                     # smaller context = faster quiz inference
OLLAMA_BASE      = "http://localhost:11434"

# RunPod persistence: use the pod's volume (default mount is /workspace).
# Override with the LECTUREFORGE_DIR env var if your volume is mounted elsewhere.
import os as _os
BASE_DIR         = _os.environ.get("LECTUREFORGE_DIR", "/workspace/LectureForge")
UPLOAD_DIR       = f"{BASE_DIR}/uploads"
OUTPUT_DIR       = f"{BASE_DIR}/outputs"
CLIPS_DIR        = f"{BASE_DIR}/outputs/clips"

# Local persistence (survives across runs as long as the pod volume persists)
DRIVE_DIR        = BASE_DIR
SESSIONS_DIR     = f"{BASE_DIR}/sessions"
OCR_CACHE_DIR    = f"{BASE_DIR}/ocr_cache"

# Port the Flask API listens on — expose this as an HTTP port in the RunPod template
API_PORT         = int(_os.environ.get("PORT", "5000"))

# EasyOCR languages — add more codes from https://www.jaided.ai/easyocr/
OCR_LANGUAGES    = ["en", "bn"]          # English + Bengali

# ══════════════════════════════════════════════════════════════
# IMPORTS
# ══════════════════════════════════════════════════════════════



# (Dependencies are installed by Cell 1 above. To enable equation→LaTeX PDF
#  extraction set the env var before launch, e.g. in a cell:  %env LF_MATH_PDF=1 )



import os, json, re, time, copy, difflib, threading, subprocess, shutil, hashlib
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed



import cv2
import numpy as np
import easyocr
import faiss
import requests as http_requests
from faster_whisper import WhisperModel
from sentence_transformers import SentenceTransformer
from flask import Flask, request, jsonify, send_file, Response
from flask_cors import CORS
from werkzeug.utils import secure_filename

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt



# RunPod: no Google Drive — persistence is the local pod volume under BASE_DIR.



for d in [UPLOAD_DIR, OUTPUT_DIR, CLIPS_DIR, DRIVE_DIR, SESSIONS_DIR, OCR_CACHE_DIR]:
    os.makedirs(d, exist_ok=True)



ALLOWED_VIDEO = {'mp4', 'avi', 'mov', 'mkv', 'webm'}
ALLOWED_DOCS  = {'pdf', 'docx', 'pptx', 'ppt'}



print("✅ Imports & config ready.")
print(f"✅ Persistence ready (local volume) — sessions saved to {SESSIONS_DIR}")

# ══════════════════════════════════════════════════════════════
# APP STATE — single place for all runtime data
# ══════════════════════════════════════════════════════════════



class AppState:
    """Keeps every piece of mutable runtime data in one object
    instead of scattered globals."""

    def __init__(self):
        self.lock         = threading.Lock()
        self.active_model = OLLAMA_MODEL
        self.video_path   = ""
        self.doc_path     = ""
        self.suggested_questions = []
        self.session_id   = ""
        self.agent_log    = []
        self.student_performance = {}    # topic -> list of 0.0/1.0 scores
        self._quiz_agent_thinking = []   # thinking log for manual/adaptive quiz agents
        self._quiz_thinking_session = 0  # bumps each time a fresh agentic generation starts
        self._quiz_thinking_active  = False  # true while the agentic pipeline is running
        self.quiz_history = []          # list of graded quiz submissions saved per session
        self.generated_quiz_runs = []   # snapshots of each generation (async + upload) for UI history
        self.quiz_generation_jobs = {} # quiz generation jobs for live thinking
        self.quiz_generation_lock = threading.Lock()

        self.status = {
            "state": "idle", "stage": "Idle", "pct": 0, "error": "",
            "summary": "", "transcript": "", "board_text": "", "lecture_timeline": "",
        }

        self.quiz = {
            "mcq": [], "tf": [], "fill": [], "short": [],
            "flash": [], "hints": {},
        }

    def update_status(self, **kw):
        with self.lock:
            self.status.update(kw)

    def reset_status(self):
        with self.lock:
            self.status.update(
                state="running", stage="Upload received…", pct=2, error="",
                summary="", transcript="", board_text="", lecture_timeline="",
            )
            self._quiz_agent_thinking   = []
            self._quiz_thinking_active  = False
            self._quiz_thinking_session += 1
            self.generated_quiz_runs = []  # new material — clear past-set history

    # ── Persistence ───────────────────────────────────────────

    def _make_session_id(self, filename):
        name = Path(filename).stem
        short_hash = hashlib.md5(f"{name}{time.time()}".encode()).hexdigest()[:6]
        date_str = datetime.now().strftime("%Y%m%d_%H%M")
        return f"{date_str}_{name}_{short_hash}"

    def save_session(self, source_filename):
        """Dump full session to Google Drive, including performance and source file."""
        with self.lock:
            if not self.session_id:
                self.session_id = self._make_session_id(source_filename)

            folder = os.path.join(SESSIONS_DIR, self.session_id)
            os.makedirs(folder, exist_ok=True)

            payload = {
                "session_id":            self.session_id,
                "source_file":           source_filename,
                "saved_at":              datetime.now().isoformat(),
                "summary":               self.status.get("summary", ""),
                "transcript":            self.status.get("transcript", ""),
                "board_text":            self.status.get("board_text", ""),
                "lecture_timeline":      self.status.get("lecture_timeline", ""),
                "suggested_questions":   self.suggested_questions,
                "quiz":                  self.quiz,
                "quiz_history":         self.quiz_history,
                "generated_quiz_runs":   self.generated_quiz_runs,
                "student_performance":   self.student_performance,
            }

            session_path = os.path.join(folder, "session.json")
            with open(session_path, "w", encoding="utf-8") as f:
                json.dump(payload, f, indent=2, ensure_ascii=False)

            # copy text artifacts
            for fname in ["summary.txt", "transcript.txt", "board_text.txt",
                           "lecture_timeline.txt", "board_entries.json", "transcript_segments.json"]:
                src = os.path.join(OUTPUT_DIR, fname)
                if os.path.exists(src):
                    shutil.copy2(src, os.path.join(folder, fname))

            # copy the original uploaded file so /stream/video works after restore
            source_path = self.video_path or self.doc_path
            if source_path and os.path.exists(source_path):
                dest = os.path.join(folder, Path(source_path).name)
                if not os.path.exists(dest):
                    shutil.copy2(source_path, dest)

        print(f"💾 Session saved → {folder}")

    def load_session(self, session_id):
        """Restore a previously saved session from Drive."""
        folder = os.path.join(SESSIONS_DIR, session_id)
        session_path = os.path.join(folder, "session.json")
        if not os.path.exists(session_path):
            print(f"  Session {session_id} not found.")
            return False

        try:
            with open(session_path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except (json.JSONDecodeError, OSError) as e:
            print(f"  Corrupt session file: {e}")
            return False

        with self.lock:
            self.session_id          = session_id
            self.agent_log           = []
            self.suggested_questions = data.get("suggested_questions", [])
            self.quiz                = data.get("quiz", self.quiz)
            self.quiz_history       = data.get("quiz_history", [])
            self.generated_quiz_runs = data.get("generated_quiz_runs", [])
            self.status.update(
                state="done", stage="Restored from Drive ✓", pct=100, error="",
                summary=data.get("summary", ""),
                transcript=data.get("transcript", ""),
                board_text=data.get("board_text", ""),
                lecture_timeline=data.get("lecture_timeline", ""),
            )

            # merge performance history (accumulate, don't replace)
            for topic, scores in data.get("student_performance", {}).items():
                if topic not in self.student_performance:
                    self.student_performance[topic] = []
                self.student_performance[topic].extend(scores)

            # restore video/doc path from the saved copy
            self.video_path = ""
            self.doc_path   = ""
            source_file = data.get("source_file", "")
            if source_file:
                saved_copy = os.path.join(folder, source_file)
                if os.path.exists(saved_copy):
                    ext = Path(source_file).suffix.lower().lstrip('.')
                    if ext in ALLOWED_VIDEO:
                        self.video_path = saved_copy
                    elif ext in ALLOWED_DOCS:
                        self.doc_path = saved_copy

        # copy text files back to OUTPUT_DIR so /download endpoints work
        for fname in ["summary.txt", "transcript.txt", "board_text.txt",
                       "lecture_timeline.txt", "board_entries.json", "transcript_segments.json"]:
            src = os.path.join(folder, fname)
            if os.path.exists(src):
                shutil.copy2(src, os.path.join(OUTPUT_DIR, fname))

        with self.lock:
            if not self.status.get("lecture_timeline") and self.video_path:
                tl = rebuild_lecture_timeline_from_disk()
                if tl:
                    self.status["lecture_timeline"] = tl

        print(f"✅ Session restored: {session_id}")
        return True

    @staticmethod
    def list_sessions():
        sessions = []
        if not os.path.isdir(SESSIONS_DIR):
            return sessions
        for name in sorted(os.listdir(SESSIONS_DIR), reverse=True):
            sp = os.path.join(SESSIONS_DIR, name, "session.json")
            if os.path.exists(sp):
                try:
                    with open(sp, "r", encoding="utf-8") as f:
                        meta = json.load(f)
                    sessions.append({
                        "session_id":  name,
                        "source_file": meta.get("source_file", ""),
                        "saved_at":    meta.get("saved_at", ""),
                        "has_summary": bool(meta.get("summary")),
                        "mcq_count":   len(meta.get("quiz", {}).get("mcq", [])),
                        "tf_count":    len(meta.get("quiz", {}).get("tf", [])),
                        "flash_count": len(meta.get("quiz", {}).get("flash", [])),
                    })
                except Exception:
                    pass
        return sessions



state = AppState()

# ══════════════════════════════════════════════════════════════
# KEEP-ALIVE  (prevents Colab disconnecting after 90 min idle)
# ══════════════════════════════════════════════════════════════



def _keep_alive():
    while True:
        try:
            with open(os.path.join(OUTPUT_DIR, '.heartbeat'), 'w') as f:
                f.write(str(time.time()))
        except Exception:
            pass
        time.sleep(45)



threading.Thread(target=_keep_alive, daemon=True).start()
print("✅ Keep-alive thread started.")

# ══════════════════════════════════════════════════════════════
# LAZY MODEL LOADERS
# ══════════════════════════════════════════════════════════════



_ocr_reader = _whisper_model = _embed_model = None



def _gpu_available():
    """Check once whether CUDA is usable."""
    try:
        import torch
        return torch.cuda.is_available()
    except ImportError:
        return False



def get_ocr():
    global _ocr_reader
    if _ocr_reader is None:
        gpu = _gpu_available()
        print(f"  Loading EasyOCR ({OCR_LANGUAGES}) on {'GPU' if gpu else 'CPU'}…")
        _ocr_reader = easyocr.Reader(OCR_LANGUAGES, gpu=gpu)
    return _ocr_reader



def get_whisper():
    global _whisper_model
    if _whisper_model is None:
        # always use small — fits next to Ollama on T4; medium/large OOM easily
        for device, ctype in [("cuda", "float16"), ("cpu", "int8")]:
            if device == "cuda" and not _gpu_available():
                continue
            try:
                print(f"  Loading Whisper 'small' on {device.upper()} ({ctype})…")
                _whisper_model = WhisperModel("small", device=device, compute_type=ctype)
                print(f"  ✅ Whisper 'small' on {device.upper()}")
                return _whisper_model
            except Exception as e:
                print(f"  ⚠ Whisper small on {device} failed: {e}")
                _whisper_model = None
        raise RuntimeError("Could not load Whisper small on GPU or CPU")
    return _whisper_model



def get_embedder():
    global _embed_model
    if _embed_model is None:
        print("  Loading SentenceTransformer all-MiniLM-L6-v2…")
        _embed_model = SentenceTransformer('all-MiniLM-L6-v2')
    return _embed_model



print("✅ Lazy loaders defined.")

# ══════════════════════════════════════════════════════════════
# OLLAMA SETUP  (start server + pull model)
# ══════════════════════════════════════════════════════════════



def start_ollama():
    env = {
        **os.environ,
        "OLLAMA_KEEP_ALIVE":    "-1",
        "CUDA_VISIBLE_DEVICES": "0",
        "OLLAMA_GPU_OVERHEAD":  "500000000",
    }
    subprocess.Popen(
        ["ollama", "serve"],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, env=env,
    )
    time.sleep(5)
    print("✅ Ollama server started (GPU mode).")



def pull_model(name):
    print(f"  Pulling {name}… (~7 GB first time)")
    r = subprocess.run(["ollama", "pull", name], capture_output=False)
    if r.returncode != 0:
        return False
    # Pre-warm with 32 K context window
    subprocess.run(
        ["ollama", "run", name, "--", "hi"],
        capture_output=True, text=True,
        env={**os.environ, "OLLAMA_NUM_CTX": "32768"},
    )
    print(f"✅ {name} ready | context: 32K tokens")
    return True



def setup_ollama():
    start_ollama()
    if pull_model(OLLAMA_MODEL):
        state.active_model = OLLAMA_MODEL
    elif pull_model(OLLAMA_FALLBACK):
        state.active_model = OLLAMA_FALLBACK
    else:
        print("❌ No model available — check internet connection.")
    # pull the fast quiz model if it's a different model
    if QUIZ_MODEL and QUIZ_MODEL != state.active_model:
        print(f"  Pulling fast quiz model ({QUIZ_MODEL})…")
        pull_model(QUIZ_MODEL)
    print(f"✅ Active model: {state.active_model} | Quiz model: {QUIZ_MODEL or state.active_model}")



setup_ollama()

# ══════════════════════════════════════════════════════════════
# LLM WRAPPER
# ══════════════════════════════════════════════════════════════



def call_ollama(prompt, system="You are an expert educational AI assistant.",
                json_mode=False, max_retries=3, timeout=300,
                model=None, num_ctx=None):
    """Send a prompt to the local Ollama instance and return the response text."""
    payload = {
        "model":   model or state.active_model,
        "prompt":  prompt,
        "system":  system,
        "stream":  False,
        "options": {"num_ctx": num_ctx or 32768, "temperature": 0.2},
    }
    if json_mode:
        payload["format"] = "json"

    for attempt in range(1, max_retries + 1):
        try:
            resp = http_requests.post(
                f"{OLLAMA_BASE}/api/generate", json=payload, timeout=timeout,
            )
            resp.raise_for_status()
            text = resp.json().get("response", "").strip()
            if json_mode:
                text = re.sub(r'^```(?:json)?\s*', '', text)
                text = re.sub(r'\s*```$', '', text)
            return text
        except Exception as e:
            if attempt == max_retries:
                raise
            print(f"  LLM attempt {attempt} failed: {e} — retrying…")
            time.sleep(2)



def call_ollama_json(prompt, system="You are an expert educational AI assistant.",
                     fallback=None, model=None, num_ctx=None):
    """call_ollama with JSON mode. Returns parsed dict/list, or *fallback* on error."""
    try:
        raw = call_ollama(prompt, system=system, json_mode=True,
                          model=model, num_ctx=num_ctx)
        return json.loads(raw)
    except Exception as e:
        print(f"  JSON parse error: {e}")
        return fallback if fallback is not None else {}



def _extract_items(data, key):
    """Pull the list of items out of a model JSON response, tolerating shape
    drift. Models often ignore the exact wrapper and return a bare array
    `[...]`, or use a different key (e.g. 'flashcards' instead of 'cards') —
    which previously caused silently-empty results. Handles:
      {key: [...]}  |  [...]  |  {some_other_key: [...]}"""
    if isinstance(data, list):
        return data
    if isinstance(data, dict):
        v = data.get(key)
        if isinstance(v, list):
            return v
        for val in data.values():          # first list-valued field
            if isinstance(val, list):
                return val
    return []


def call_ollama_json_quiz(prompt, key="questions", force_primary=False):
    """Normalises to {key: [...]} and tolerates bare-array / wrong-key replies.
    Default: try fast quiz model first, fall back to primary on empty.
    force_primary=True: skip the fast model entirely and use the primary model —
    for BULK array tasks (e.g. flashcards) the small 4B model reliably returns
    broken/empty JSON, so trying it first only wastes a slow round-trip."""
    fb = {key: []}
    if not force_primary:
        data  = call_ollama_json(prompt, fallback=fb, model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
        items = _extract_items(data, key)
        if items:
            return {key: items}
        print(f"  ⚠ Empty or bad '{key}' from {QUIZ_MODEL} — using {state.active_model}…")
    data2 = call_ollama_json(prompt, fallback=fb)        # primary (larger) model
    return {key: _extract_items(data2, key)}



def call_ollama_json_list_quiz(prompt, fallback_list, force_primary=False):
    """JSON list (e.g. suggested questions). force_primary=True skips the fast
    model and goes straight to the primary one (4B is unreliable at JSON lists)."""
    if not force_primary:
        data  = call_ollama_json(prompt, fallback=fallback_list,
                                 model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
        items = _extract_items(data, "items")
        if items:
            return items
        print(f"  ⚠ Empty list from {QUIZ_MODEL} — using {state.active_model}…")
    data2 = call_ollama_json(prompt, fallback=fallback_list)
    return _extract_items(data2, "items") or (fallback_list if isinstance(fallback_list, list) else [])



print("✅ LLM wrapper ready.")

# ══════════════════════════════════════════════════════════════
# SHARED HELPERS
# ══════════════════════════════════════════════════════════════



def fmt_ts(seconds):
    """Convert seconds → 'MM:SS' string."""
    m, s = divmod(int(seconds), 60)
    return f"{m:02d}:{s:02d}"



def nearest_timestamp(segment_list, index):
    """Pick the segment timestamp closest to a question's ordinal index."""
    if not segment_list:
        return ""
    seg_idx = min(
        int(index / max(len(segment_list), 1) * len(segment_list)),
        len(segment_list) - 1,
    )
    return segment_list[seg_idx].get("start_str", "")



def fill_missing_timestamps(questions, segments=None, board_entries=None):
    """Back-fill source_timestamp on questions that don't have one."""
    for i, q in enumerate(questions):
        if q.get("source_timestamp"):
            continue
        if segments:
            q["source_timestamp"] = nearest_timestamp(segments, i)
        elif board_entries:
            q["source_timestamp"] = board_entries[min(i, len(board_entries) - 1)]["timestamp_str"]



def compute_grade(score, total):
    """Return (percentage, letter_grade) from a raw score."""
    pct = round(score / max(total, 1) * 100)
    if   pct >= 90: grade = "A"
    elif pct >= 80: grade = "B"
    elif pct >= 70: grade = "C"
    elif pct >= 60: grade = "D"
    else:           grade = "F"
    return pct, grade



def save_text(filename, content):
    """Write a text file into OUTPUT_DIR."""
    with open(f"{OUTPUT_DIR}/{filename}", "w", encoding="utf-8") as f:
        f.write(content)



def save_json(filename, data):
    """Write a JSON file into OUTPUT_DIR."""
    with open(f"{OUTPUT_DIR}/{filename}", "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

# ══════════════════════════════════════════════════════════════
# VIDEO EXTRACTION  (OCR + Whisper)
# ══════════════════════════════════════════════════════════════



def extract_frames(video_path, interval_sec=5):
    cap   = cv2.VideoCapture(video_path)
    fps   = cap.get(cv2.CAP_PROP_FPS) or 25.0
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    dur   = total / max(fps, 1)
    # shorter lectures → denser sampling so slides are not missed
    if dur > 3600:
        interval_sec = max(interval_sec, 10)
    step  = max(1, int(fps * interval_sec))
    frames, n = [], 0
    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        if n % step == 0:
            frames.append((n / fps, frame))
        n += 1
    cap.release()
    print(f"  Sampled {len(frames)} frames | ~{total/fps:.0f}s duration")
    return frames



def classify_frame(frame):
    """Returns (type, preprocessed_frame).  Types: slide / blackboard / whiteboard / scene."""
    gray   = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
    mean_v = np.mean(gray)
    std_v  = np.std(gray)

    def sharpen(f):
        k = np.array([[0, -1, 0], [-1, 5, -1], [0, -1, 0]])
        return cv2.filter2D(f, -1, k)

    def resize960(f):
        h, w = f.shape[:2]
        return cv2.resize(
            f, (960, int(h * 960 / w)),
            interpolation=cv2.INTER_CUBIC if w < 960 else cv2.INTER_AREA,
        )

    def clahe_color(f, clip=2.8):
        lab = cv2.cvtColor(f, cv2.COLOR_BGR2LAB)
        l, a, b = cv2.split(lab)
        cl = cv2.createCLAHE(clipLimit=clip, tileGridSize=(8, 8)).apply(l)
        return cv2.cvtColor(cv2.merge([cl, a, b]), cv2.COLOR_LAB2BGR)

    # Black / dark-theme slides & OLED decks — must run BEFORE chalk invert (invert ruins navy UI)
    if mean_v < 130 and std_v >= 16:
        enh = clahe_color(frame, clip=3.6)
        k = np.array([[0, -0.5, 0], [-0.5, 3, -0.5], [0, -0.5, 0]])
        enh = cv2.filter2D(enh, -1, k)
        return 'slide', resize960(enh)

    # Uniform dark chalk wall (little edge detail — not a busy dark slide)
    if mean_v < 80 and std_v < 16:
        inv = cv2.bitwise_not(frame)
        lab = cv2.cvtColor(inv, cv2.COLOR_BGR2LAB)
        l, a, b = cv2.split(lab)
        cl = cv2.createCLAHE(clipLimit=3.5, tileGridSize=(8, 8)).apply(l)
        enh = cv2.cvtColor(cv2.merge([cl, a, b]), cv2.COLOR_LAB2BGR)
        return 'blackboard', resize960(enh)

    # Mid-brightness projected slides (beige template, lit room)
    if std_v >= 20 and 78 <= mean_v <= 195:
        enh = clahe_color(frame, clip=3.0)
        k = np.array([[0, -0.5, 0], [-0.5, 3, -0.5], [0, -0.5, 0]])
        enh = cv2.filter2D(enh, -1, k)
        return 'slide', resize960(enh)

    # Bright classic slides (white / light gray backgrounds)
    if mean_v > 152 and std_v > 16:
        g = cv2.cvtColor(sharpen(frame), cv2.COLOR_BGR2GRAY)
        _, p = cv2.threshold(g, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        return 'slide', resize960(cv2.cvtColor(p, cv2.COLOR_GRAY2BGR))

    # Physical whiteboards / light walls with writing
    if 65 <= mean_v <= 200:
        enh = clahe_color(frame, clip=2.5)
        k = np.array([[0, -0.5, 0], [-0.5, 3, -0.5], [0, -0.5, 0]])
        enh = cv2.filter2D(enh, -1, k)
        return 'whiteboard', resize960(enh)

    return 'scene', frame



def frames_are_similar(f1, f2, threshold=0.97):
    g1 = cv2.resize(cv2.cvtColor(f1, cv2.COLOR_BGR2GRAY), (160, 90))
    g2 = cv2.resize(cv2.cvtColor(f2, cv2.COLOR_BGR2GRAY), (160, 90))
    # Fast mean-absolute-diff pre-filter — avoids expensive matchTemplate for
    # frames that are obviously identical (mad < 2) or obviously different (mad > 35)
    mad = np.mean(np.abs(g1.astype(np.float32) - g2.astype(np.float32)))
    if mad < 2.0:
        return True
    if mad > 35.0:
        return False
    r = cv2.matchTemplate(
        g1.astype(np.float32), g2.astype(np.float32), cv2.TM_CCORR_NORMED,
    )
    return r[0][0] > threshold



def _ocr_frame(reader, proc, ftype):
    """Run EasyOCR with thresholds tuned per frame type."""
    try:
        kw = dict(detail=0, paragraph=True, min_size=8,
                  text_threshold=0.45, low_text=0.25, width_ths=0.9)
        if ftype == 'slide':
            # dark-theme slides: light text on dark bg needs lower thresholds
            kw.update(min_size=8, text_threshold=0.38, low_text=0.22, width_ths=0.85)
        elif ftype == 'scene':
            kw.update(min_size=6, text_threshold=0.35, low_text=0.2, width_ths=0.7)
        else:
            kw.update(min_size=12, text_threshold=0.52, low_text=0.28, width_ths=0.85)
        parts = []
        for t in reader.readtext(proc, **kw):
            s = (t[1] if isinstance(t, (list, tuple)) and len(t) > 1 else str(t)).strip()
            if len(s) > 2:
                parts.append(s)
        return " | ".join(parts)
    except Exception as ex:
        print(f"\n  OCR error: {ex}")
        return ""



BOARD_TEXT_EMPTY = "No significant board/slide text detected."



# ── OCR Drive cache helpers ──────────────────────────────────────

def _video_ocr_hash(video_path, chunk_size=65536):
    """MD5 of first 4 MB of the video — fast proxy for file identity."""
    h = hashlib.md5()
    read = 0
    with open(video_path, "rb") as f:
        while read < 4 * 1024 * 1024:
            buf = f.read(min(chunk_size, 4 * 1024 * 1024 - read))
            if not buf:
                break
            h.update(buf)
            read += len(buf)
    # also fold in file size so renamed copies are treated differently
    h.update(str(os.path.getsize(video_path)).encode())
    return h.hexdigest()

def _ocr_cache_path(video_path):
    return os.path.join(OCR_CACHE_DIR, f"{_video_ocr_hash(video_path)}.json")

def _load_ocr_cache(video_path):
    """Return (board_text, entries) from Drive cache, or None if missing/corrupt."""
    try:
        path = _ocr_cache_path(video_path)
        if not os.path.isfile(path):
            return None
        with open(path, "r", encoding="utf-8") as f:
            cached = json.load(f)
        entries    = cached.get("entries", [])
        board_text = cached.get("board_text", BOARD_TEXT_EMPTY)
        print(f"  ✅ OCR cache hit — {len(entries)} entries loaded from Drive (skipping re-scan)")
        return board_text, entries
    except Exception as e:
        print(f"  ⚠ OCR cache read failed: {e} — will re-run OCR")
        return None

def _save_ocr_cache(video_path, board_text, entries):
    """Persist OCR result to Drive for future runs."""
    try:
        path = _ocr_cache_path(video_path)
        with open(path, "w", encoding="utf-8") as f:
            json.dump({"board_text": board_text, "entries": entries},
                      f, ensure_ascii=False)
        print(f"  💾 OCR result cached → {path}")
    except Exception as e:
        print(f"  ⚠ OCR cache save failed (non-fatal): {e}")



def format_visual_lecture_text(board_text, board_entries, max_chars=14000):
    """Flatten OCR: slides (light/dark), whiteboards, chalkboards, on-screen text — one string for LLM/RAG."""
    lines = []
    for e in (board_entries or []):
        tag = (e.get("tag") or e.get("frame_type") or "Visual").strip()
        ts  = e.get("timestamp_str", "")
        tx  = (e.get("text") or "").strip()
        if len(tx) < 2:
            continue
        lines.append(f"[{ts}] [{tag}] {tx}")
    blob = "\n".join(lines)
    bt = (board_text or "").strip()
    if bt and bt != BOARD_TEXT_EMPTY and bt not in blob:
        blob = (bt + "\n" + blob).strip() if blob else bt
    elif not blob and bt and bt != BOARD_TEXT_EMPTY:
        blob = bt
    return blob[:max_chars]



def pack_video_sources(board_text, transcript, board_entries, max_total=9000):
    """Single block: all visual OCR + speech — used by summary, quiz, flashcards."""
    vis = format_visual_lecture_text(board_text, board_entries, max(2000, max_total - 2000))
    tr  = (transcript or "").strip()
    tr  = tr[: max(800, max_total - len(vis) - 80)]
    parts = []
    if vis:
        parts.append(f"=== SLIDES + BOARDS (OCR: white/black themes, on-screen) ===\n{vis}")
    if tr:
        parts.append(f"=== SPEECH (TRANSCRIPT) ===\n{tr}")
    return "\n\n".join(parts)[:max_total] if parts else tr



def load_board_entries_from_disk():
    path = os.path.join(OUTPUT_DIR, "board_entries.json")
    if not os.path.isfile(path):
        return []
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError):
        return []



def build_lecture_timeline_text(segment_list, board_entries):
    """Speech + slide/board OCR in time order (one view, like a full transcript)."""
    events = []
    for s in segment_list or []:
        t = s.get("start_sec")
        if t is None:
            continue
        ts = s.get("start_str", "")
        txt = (s.get("text") or "").strip()
        if not txt:
            continue
        events.append((float(t), f"[{ts}] (Speech) {txt}"))
    for e in board_entries or []:
        t = e.get("timestamp_sec")
        if t is None:
            continue
        tag = e.get("tag", "Visual")
        ts = e.get("timestamp_str", "")
        txt = (e.get("text") or "").strip()
        if not txt:
            continue
        events.append((float(t), f"[{ts}] [{tag}] {txt}"))
    events.sort(key=lambda x: x[0])
    return "\n\n".join(line for _, line in events) if events else ""



def rebuild_lecture_timeline_from_disk():
    sp = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    bp = os.path.join(OUTPUT_DIR, "board_entries.json")
    segs, boards = [], []
    try:
        if os.path.isfile(sp):
            with open(sp, "r", encoding="utf-8") as f:
                segs = json.load(f)
    except (json.JSONDecodeError, OSError):
        pass
    try:
        if os.path.isfile(bp):
            with open(bp, "r", encoding="utf-8") as f:
                boards = json.load(f)
    except (json.JSONDecodeError, OSError):
        pass
    return build_lecture_timeline_text(segs, boards)



def extract_board_text(video_path):
    """Run EasyOCR on sampled frames.  Returns (board_text_str, entries_list).
    Results are cached on Google Drive by video hash — re-uploads of the same
    file skip the expensive OCR pass entirely."""
    cached = _load_ocr_cache(video_path)
    if cached is not None:
        board_text, entries = cached
        save_text("board_text.txt",
                  f"=== BOARD & SLIDE TEXT ===\n(from cache)\n\n{board_text}")
        save_json("board_entries.json", entries)
        return board_text, entries

    reader    = get_ocr()
    frames    = extract_frames(video_path)
    entries   = []
    type_log  = {t: 0 for t in ('blackboard', 'whiteboard', 'slide', 'scene')}
    prev_text = ""
    prev_frm  = None
    skipped   = 0

    LABEL_MAP = {'blackboard': 'Board', 'whiteboard': 'Board', 'slide': 'Slide', 'scene': 'On-screen'}

    def process_frame(ts, frame):
        nonlocal prev_text, entries, type_log
        ftype, proc = classify_frame(frame)
        type_log[ftype] = type_log.get(ftype, 0) + 1

        h, w = frame.shape[:2]
        rw = max(w, 1)
        resized = cv2.resize(frame, (960, int(h * 960 / rw)))
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        eq = cv2.cvtColor(cv2.equalizeHist(gray), cv2.COLOR_GRAY2BGR)
        h2, w2 = eq.shape[:2]
        eq_r = cv2.resize(eq, (960, int(h2 * 960 / max(w2, 1))))

        candidates = [(ftype, proc)]
        if ftype == 'scene':
            candidates.append(('scene', resized))
            candidates.append(('scene', eq_r))
        else:
            candidates.append(('slide', resized))
            candidates.append(('slide', eq_r))

        text = ""
        for c_ft, c_im in candidates:
            t = _ocr_frame(reader, c_im, c_ft)
            if t:
                text = t
                break
        if not text:
            return
        if difflib.SequenceMatcher(None, prev_text, text).ratio() < 0.78:
            ts_str = fmt_ts(ts)
            entries.append({
                "timestamp_sec": round(ts, 1),
                "timestamp_str": ts_str,
                "frame_type": ftype,
                "tag": LABEL_MAP.get(ftype, 'On-screen'),
                "text": text,
            })
            prev_text = text

    for i, (ts, frame) in enumerate(frames):
        print(f"  OCR frame {i+1}/{len(frames)} @{ts:.0f}s  skipped={skipped}", end='\r')

        if prev_frm is not None and frames_are_similar(prev_frm, frame, threshold=0.97):
            skipped += 1
            prev_frm = frame
            continue
        prev_frm = frame
        process_frame(ts, frame)

    # second pass: denser sampling if almost nothing found (talking-head + small text)
    if len(entries) < 2 and frames:
        print("\n  OCR sparse — second pass every 3s including scene-like frames…")
        frames2 = extract_frames(video_path, interval_sec=3)
        prev_frm = None
        for ts, frame in frames2:
            if prev_frm is not None and frames_are_similar(prev_frm, frame, threshold=0.97):
                prev_frm = frame
                continue
            prev_frm = frame
            process_frame(ts, frame)

    board_text = "\n".join(
        f"[{e['timestamp_str']}] {e['tag']}  {e['text']}" for e in entries
    ) or BOARD_TEXT_EMPTY

    save_text("board_text.txt",
              f"=== BOARD & SLIDE TEXT ===\nFrame types: {type_log}\n\n{board_text}")
    save_json("board_entries.json", entries)

    print(f"\n  OCR done — {len(entries)} unique states | {type_log}")
    _save_ocr_cache(video_path, board_text, entries)
    return board_text, entries



def transcribe_video(video_path):
    """Transcribe audio with Faster-Whisper.  Returns (transcript_str, segment_list)."""
    audio_tmp = f"{OUTPUT_DIR}/_audio.wav"
    print("  Extracting audio…")
    subprocess.run([
        'ffmpeg', '-y', '-i', video_path,
        '-vn', '-acodec', 'pcm_s16le', '-ar', '16000', '-ac', '1', audio_tmp,
    ], capture_output=True, check=False)

    if not os.path.exists(audio_tmp):
        return "Audio extraction failed.", []

    model = get_whisper()
    print("  Transcribing…")
    gen, info = model.transcribe(
        audio_tmp, beam_size=5, language=None,
        vad_filter=True,
        vad_parameters={"min_silence_duration_ms": 500},
        word_timestamps=False,
    )
    print(f"  Detected: {info.language} (p={info.language_probability:.2f})")

    segs, lines = [], []
    for seg in gen:
        ts = fmt_ts(seg.start)
        segs.append({
            "start_sec": round(seg.start, 1),
            "end_sec":   round(seg.end,   1),
            "start_str": ts,
            "end_str":   fmt_ts(seg.end),
            "text":      seg.text.strip(),
        })
        lines.append(f"[{ts}]  {seg.text.strip()}")

    transcript = "\n".join(lines) or "No speech detected."
    try:
        os.remove(audio_tmp)
    except OSError:
        pass

    save_text("transcript.txt", "=== LECTURE TRANSCRIPT ===\n\n" + transcript)
    save_json("transcript_segments.json", segs)

    print("  Transcript saved.")
    return transcript, segs



print("✅ Video extraction ready.")

# ══════════════════════════════════════════════════════════════
# DOCUMENT EXTRACTION  (PDF / DOCX / PPTX)
# ══════════════════════════════════════════════════════════════



def _nougat_extract(path):
    """Math-aware PDF extraction via Nougat — outputs Markdown with LaTeX for
    equations (far better than plain text for calculus/physics PDFs). Requires
    `pip install nougat-ocr` and a GPU. Returns the text, or '' on any failure
    so the caller can fall back to PyMuPDF."""
    try:
        import tempfile, glob
        outdir = tempfile.mkdtemp(prefix="nougat_")
        # CLI is the most stable interface; uses the GPU automatically.
        subprocess.run(["nougat", path, "-o", outdir],
                       capture_output=True, timeout=2400)
        mmds = glob.glob(os.path.join(outdir, "*.mmd"))
        if not mmds:
            return ""
        with open(mmds[0], "r", encoding="utf-8") as f:
            return f.read().strip()
    except Exception as e:
        print(f"  [extract_pdf] Nougat unavailable/failed ({e}); using PyMuPDF")
        return ""


def extract_pdf(path):
    # Math-aware path (Nougat → LaTeX) when enabled with LF_MATH_PDF=1; this is
    # the only way equations survive. Falls back to fast PyMuPDF text otherwise
    # or if Nougat fails.
    if os.environ.get("LF_MATH_PDF", "0") == "1":
        mmd = _nougat_extract(path)
        if mmd and len(mmd) > 50:
            print(f"  [extract_pdf] Nougat math extraction OK ({len(mmd)} chars, LaTeX preserved)")
            return [{"page": 1, "text": mmd, "source": "pdf-nougat"}]

    import fitz
    doc = fitz.open(path)
    entries = []
    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            entries.append({"page": i + 1, "text": text, "source": "pdf"})
    doc.close()
    return entries



def extract_docx(path):
    from docx import Document
    doc = Document(path)
    entries, buf, sec = [], [], 1
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        if 'heading' in para.style.name.lower() and buf:
            entries.append({"section": sec, "text": " ".join(buf), "source": "docx"})
            buf = []
            sec += 1
        buf.append(t)
    if buf:
        entries.append({"section": sec, "text": " ".join(buf), "source": "docx"})
    return entries



def extract_pptx(path):
    """Pull text from a .pptx. Covers title + body text frames, tables,
    grouped shapes (recursively), and speaker notes — lecture decks are often
    image/diagram heavy, so we mine every text-bearing element to avoid the
    'too little text extracted' failure. Each shape is guarded so one odd
    shape can't abort a whole slide."""
    from pptx import Presentation
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:
        MSO_SHAPE_TYPE = None
    prs = Presentation(path)
    entries = []

    def collect(shape, parts):
        # Grouped shapes → recurse into children
        try:
            if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in shape.shapes:
                    collect(sub, parts)
                return
        except Exception:
            pass
        # Tables → join each row's non-empty cells
        try:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
                return
        except Exception:
            pass
        # Text frames (put the title placeholder first)
        try:
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if not txt:
                    return
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:        # idx 0 = title placeholder
                    parts.insert(0, txt)
                else:
                    parts.append(txt)
        except Exception:
            pass

    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            collect(shape, parts)
        # Speaker notes — often where the actual explanation lives
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append("Notes: " + notes)
        except Exception:
            pass
        if parts:
            entries.append({"slide": i + 1, "text": "\n".join(parts), "source": "pptx"})
    return entries



DOC_EXTRACTORS = {
    'pdf':  extract_pdf,
    'docx': extract_docx,
    'pptx': extract_pptx,
    'ppt':  extract_pptx,
}



print("✅ Document extraction ready.")

# ══════════════════════════════════════════════════════════════
# RAG ENGINE  (FAISS-based retrieval)
# ══════════════════════════════════════════════════════════════



class LectureRAG:
    """Encodes text chunks with SentenceTransformer, indexes with FAISS,
    and retrieves top-k relevant chunks for a query."""

    def __init__(self):
        self.embedder   = None
        self.index      = None
        self.chunks     = []
        self.timestamps = []

    def build(self, entries, text_key="text", ts_key="timestamp_str", source_key="tag"):
        self.embedder   = get_embedder()
        self.chunks     = [e[text_key] for e in entries]
        self.timestamps = [
            {"timestamp": e.get(ts_key, ""), "source": e.get(source_key, "")}
            for e in entries
        ]
        vecs = self.embedder.encode(
            self.chunks, batch_size=64, show_progress_bar=False,
        ).astype('float32')
        self.index = faiss.IndexFlatIP(vecs.shape[1])
        faiss.normalize_L2(vecs)
        self.index.add(vecs)
        print(f"  RAG index built — {len(self.chunks)} chunks.")

    def query(self, question, top_k=5):
        if not self.index:
            return "No index built yet.", []
        vec = self.embedder.encode([question]).astype('float32')
        faiss.normalize_L2(vec)
        _, idxs = self.index.search(vec, top_k)
        ctx = "\n\n".join(self.chunks[i] for i in idxs[0] if i < len(self.chunks))
        tss = [self.timestamps[i] for i in idxs[0] if i < len(self.timestamps)]
        return ctx, tss



rag = LectureRAG()



def build_rag(entries, text_key="text", ts_key="timestamp_str", src_key="tag"):
    """Rebuild the global RAG index from a new set of entries."""
    global rag
    rag = LectureRAG()
    rag.build(entries, text_key=text_key, ts_key=ts_key, source_key=src_key)



print("✅ RAG engine ready.")

# ══════════════════════════════════════════════════════════════
# ANALYSIS  (Summary + Exam Hints + Topic Ranking)
# ══════════════════════════════════════════════════════════════



def generate_summary(board_text, transcript, segment_list, board_entries):
    """Merge OCR (slides + boards, all themes) + speech into one summary."""
    visual = format_visual_lecture_text(board_text, board_entries, 10000)
    merged = (
        "Instructions: Combine everything below. OCR lines are from projector slides, "
        "digital whiteboards, chalkboards, and on-screen text (any color theme).\n\n"
        f"=== SLIDES + BOARDS (OCR) ===\n{visual}\n\n"
        f"=== SPEECH TRANSCRIPT ===\n{transcript[:10000]}"
    )
    return call_ollama(
        f"Create a comprehensive, structured lecture summary:\n\n{merged}",
        system=(
            "You are an expert educational content summarizer. "
            "Treat slide/board OCR and spoken words as one lecture — cite both where useful."
        ),
    )



EMPHASIS_SIGNALS = [
    "this is important", "remember", "this will be on", "exam", "key",
    "critical", "note", "highlight", "must know", "definition",
]



def detect_exam_hints(segment_list, board_entries, chunk_size=30):
    """Two-pass LLM analysis: find emphasis in transcript + key phrases on board."""
    all_hints = []

    def has_signal(seg):
        lower = seg.get("text", "").lower()
        return any(sig in lower for sig in EMPHASIS_SIGNALS)

    # Pass 1 — transcript emphasis moments
    for i in range(0, len(segment_list), chunk_size):
        chunk    = segment_list[i:i + chunk_size]
        filtered = [s for s in chunk if has_signal(s)]
        if not filtered:
            continue
        text = "\n".join(f"[{s['start_str']}] {s['text']}" for s in filtered)
        prompt = (
            "Identify sentences where the teacher explicitly signals exam importance. "
            "Return JSON list: [{\"text\",\"timestamp\",\"emphasis_type\",\"signal_phrase\"}]\n\n"
            + text
        )
        result = call_ollama_json(prompt, fallback=[], model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
        if isinstance(result, list):
            all_hints.extend(result)

    # Pass 2 — board/slide key phrases
    if board_entries:
        board_text = "\n".join(
            f"[{e['timestamp_str']}] {e['text'][:200]}" for e in board_entries[:40]
        )
        prompt = (
            "Extract key terms and formulas from these board/slide notes. "
            "Return JSON list: [{\"text\",\"timestamp\",\"emphasis_type\",\"signal_phrase\"}]\n\n"
            + board_text
        )
        result = call_ollama_json(prompt, fallback=[], model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
        if isinstance(result, list):
            all_hints.extend(result)

    return all_hints



def analyze_topics(summary, hints):
    """Ask LLM to rank topics by exam likelihood."""
    hint_text = "\n".join(h.get("text", "")[:200] for h in hints[:20])
    prompt = (
        "Based on the summary and teacher emphasis moments below, identify the 5-8 "
        "most exam-likely topics. "
        "Return JSON: {\"ai_important_topics\":[{\"topic\",\"confidence\",\"reason\","
        "\"evidence\",\"suggested_note\"}],\"overall_exam_focus\":\"string\"}\n\n"
        f"SUMMARY:\n{summary[:3000]}\n\nHINTS:\n{hint_text}"
    )
    return call_ollama_json(prompt, fallback={}, model=QUIZ_MODEL, num_ctx=QUIZ_CTX)



print("✅ Analysis functions ready.")

# ══════════════════════════════════════════════════════════════
# QUIZ ENGINE  —  Generators
# ══════════════════════════════════════════════════════════════



# ── Shared JSON schemas for all quiz types ──────────────────────
QUIZ_SCHEMAS = {
    "mcq":   ('{"questions":[{"question":"","options":{"A":"","B":"","C":"","D":""},'
               '"correct_answer":"A","explanation":"","topic":"","difficulty":"",'
               '"bloom_level":"","source_timestamp":""}]}'),
    "tf":    ('{"questions":[{"statement":"","answer":true,"explanation":"",'
               '"topic":"","difficulty":"","source_timestamp":""}]}'
               '  NOTE: answer must be JSON boolean true or false, never a string.'),
    "fill":  ('{"questions":[{"question":"The ___ is responsible for X.","answer":"term","hint":"starts with T",'
               '"topic":"","difficulty":"","source_timestamp":""}]}'),
    "short": ('{"questions":[{"question":"","model_answer":"","key_points":[],'
               '"marks":4,"topic":"","difficulty":"","source_timestamp":""}]}'),
}


def generate_mcq(board_text, transcript, summary, num_questions=12,
                 doc_text="", board_entries=None, segment_list=None,
                 difficulty_hint=""):
    board_entries = board_entries or []
    segment_list  = segment_list or []
    if doc_text:
        content = doc_text[:8000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 8000)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate {num_questions} high-quality MCQ from this lecture content.\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"question\",\"options\":{\"A\",\"B\",\"C\",\"D\"},"
        "\"correct_answer\",\"explanation\",\"topic\",\"difficulty\",\"bloom_level\","
        "\"source_timestamp\",\"source_type\"}]}\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, segments=segment_list)
    return qs



def generate_tf(transcript, summary, doc_text="", segment_list=None,
                difficulty_hint="", board_text="", board_entries=None):
    segment_list = segment_list or []
    if doc_text:
        content = doc_text[:6000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 7500)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate 10 True/False questions.\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"statement\",\"answer\",\"explanation\","
        "\"topic\",\"difficulty\",\"source_timestamp\"}]}\n"
        "RULES: 'answer' MUST be a JSON boolean (true or false) — never a string. "
        "Mix true and false answers roughly equally.\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, segments=segment_list)
    return qs



def generate_fill(transcript, summary, doc_text="", board_entries=None,
                  difficulty_hint="", board_text=""):
    board_entries = board_entries or []
    if doc_text:
        content = doc_text[:6000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 7500)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate 10 fill-in-the-blank questions.\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"question\",\"answer\",\"hint\","
        "\"topic\",\"difficulty\",\"source_timestamp\"}]}\n"
        "RULES: Each 'question' MUST contain exactly '___' (three underscores) "
        "to mark the blank. 'answer' is the word/phrase that fills the blank. "
        "'hint' gives a clue without revealing the answer.\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, board_entries=board_entries)
    return qs



def generate_short(transcript, summary, doc_text="", segment_list=None,
                   difficulty_hint="", board_text="", board_entries=None):
    segment_list = segment_list or []
    if doc_text:
        content = doc_text[:6000]
    else:
        content = pack_video_sources(board_text, transcript, board_entries, 7500)
    diff_line = f"\nDifficulty focus: {difficulty_hint}\n" if difficulty_hint else ""

    prompt = (
        f"Generate 8 short-answer questions (4-6 marks each).\n"
        f"{diff_line}"
        "STRICT JSON: {\"questions\":[{\"question\",\"model_answer\",\"key_points\","
        "\"marks\",\"topic\",\"difficulty\",\"source_timestamp\"}]}\n\n"
        f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
    )
    data = call_ollama_json_quiz(prompt, key="questions")
    qs   = data.get("questions", []) if isinstance(data, dict) else []
    fill_missing_timestamps(qs, segments=segment_list)
    return qs



def generate_flashcards(summary, doc_text="", board_text="", board_entries=None):
    if doc_text:
        content = (doc_text + "\n\n" + summary[:2500])[:6000]
    else:
        vis = format_visual_lecture_text(board_text, board_entries or [], 4000)
        parts = []
        if vis:
            parts.append(f"=== SLIDES + BOARDS (OCR) ===\n{vis}")
        parts.append(f"=== SUMMARY ===\n{summary[:3500]}")
        content = "\n\n".join(parts)[:6500]
    prompt = (
        "Create 15 flashcards for key terms/concepts from slides, boards, and summary. "
        "STRICT JSON: {\"cards\":[{\"front\",\"back\",\"topic\",\"difficulty\"}]}\n\n"
        f"CONTENT:\n{content}"
    )
    data = call_ollama_json_quiz(prompt, key="cards", force_primary=True)
    return data.get("cards", []) if isinstance(data, dict) else []

# ══════════════════════════════════════════════════════════════
# QUIZ ENGINE  —  Graders
# ══════════════════════════════════════════════════════════════



def grade_mcq(questions, answers):
    results, score = [], 0
    for i, q in enumerate(questions):
        user    = str(answers.get(str(i), "")).strip().upper()
        correct = str(q.get("correct_answer", "")).strip().upper()
        ok = user == correct
        if ok:
            score += 1
        results.append({**q, "user_ans": user, "correct_ans": correct, "correct": ok})

    pct, grade = compute_grade(score, len(questions))
    return {"results": results, "score": score, "total": len(questions),
            "percentage": pct, "grade": grade}



def _norm_tf(v):
    """Normalise any truthy/falsy representation to the strings 'true' or 'false'."""
    if isinstance(v, bool):
        return "true" if v else "false"
    s = str(v).strip().lower()
    if s in ("1", "yes", "true",  "t"):  return "true"
    if s in ("0", "no",  "false", "f"):  return "false"
    return s  # pass-through so exact mismatch is still caught



def grade_tf(questions, answers):
    results, score = [], 0
    for i, q in enumerate(questions):
        user    = _norm_tf(answers.get(str(i), ""))
        correct = _norm_tf(q.get("answer", False))
        ok = user == correct
        if ok:
            score += 1
        results.append({**q, "user_ans": user, "correct_ans": correct, "correct": ok})

    pct, grade = compute_grade(score, len(questions))
    return {"results": results, "score": score, "total": len(questions),
            "percentage": pct, "grade": grade}



def grade_fill(questions, answers):
    results, score = [], 0
    for i, q in enumerate(questions):
        user    = str(answers.get(str(i), "")).strip().lower()
        correct = str(q.get("answer", "")).strip().lower()

        # Substring check only when correct answer is long enough to be meaningful
        # (avoids "a" matching inside any word, "DNA" matching "banana", etc.)
        substring_ok = len(correct) >= 4 and correct in user

        # Word-boundary check: correct answer appears as a distinct word/phrase in user text
        user_words   = re.split(r'\W+', user)
        word_ok      = correct in user_words

        # Fuzzy similarity as the final fallback
        fuzzy_ok     = difflib.SequenceMatcher(None, user, correct).ratio() > 0.80

        ok = word_ok or substring_ok or fuzzy_ok
        if ok:
            score += 1
        results.append({**q, "user_ans": user, "correct_ans": correct, "correct": ok})

    pct, grade = compute_grade(score, len(questions))
    return {"results": results, "score": score, "total": len(questions),
            "percentage": pct, "grade": grade}



def grade_short(questions, answers):
    """AI-powered short-answer grading — sends each Q+A pair to the LLM."""
    results, total_score, total_marks = [], 0, 0
    for i, q in enumerate(questions):
        user_ans = str(answers.get(str(i), "")).strip()
        marks    = int(q.get("marks", 4))
        total_marks += marks

        if not user_ans:
            results.append({**q, "user_ans": "", "score": 0, "max_marks": marks,
                            "correct": False, "correct_ans": q.get("model_answer", ""),
                            "feedback": "No answer provided.", "key_points_hit": []})
            continue

        prompt = (
            f"Grade this short answer out of {marks} marks.\n"
            f"Question: {q.get('question', '')}\n"
            f"Model answer: {q.get('model_answer', '')}\n"
            f"Key points expected: {json.dumps(q.get('key_points', []))}\n"
            f"Student answer: {user_ans}\n\n"
            f"Return JSON: {{\"score\": 0-{marks}, \"feedback\": \"...\", "
            f"\"key_points_hit\": [\"...\"]}}. Be fair but strict."
        )
        graded = call_ollama_json(prompt, fallback={"score": 0, "feedback": "Grading failed."})
        sc = min(int(graded.get("score", 0)), marks)
        total_score += sc
        results.append({
            **q, "user_ans": user_ans, "score": sc, "max_marks": marks,
            "correct": sc >= marks * 0.5,
            "correct_ans": q.get("model_answer", ""),
            "feedback": graded.get("feedback", ""),
            "key_points_hit": graded.get("key_points_hit", []),
        })

    pct, grade = compute_grade(total_score, total_marks)
    return {"results": results, "score": total_score, "total": total_marks,
            "percentage": pct, "grade": grade}



def explain_wrong(wrong_questions):
    """Generate AI explanations for any list of wrong questions."""
    if not wrong_questions:
        return {}
    text = "\n".join(
        f"Q: {q.get('question', q.get('statement', ''))}  "
        f"Correct: {q.get('correct_answer', q.get('correct_ans', q.get('answer', '')))}"
        for q in wrong_questions[:8]
    )
    prompt = (
        "For each wrong question below, explain the underlying concept simply "
        "so the student can learn from their mistake.\n"
        "Return JSON: {\"explanations\":{\"<question_text>\":\"<explanation>\"}}\n\n" + text
    )
    result = call_ollama_json(prompt, fallback={"explanations": {}},
                                model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
    return result.get("explanations", {}) if isinstance(result, dict) else {}



print("✅ Quiz engine ready.")

# ══════════════════════════════════════════════════════════════
# AGENT FRAMEWORK
# ══════════════════════════════════════════════════════════════



class AgentContext:
    """Shared memory for a single pipeline run.  Every agent reads/writes
    here instead of passing dozens of parameters between functions."""

    def __init__(self, content_type, file_path):
        self.content_type  = content_type        # "video" | "document"
        self.file_path     = file_path
        self.source_name   = Path(file_path).name

        # populated by extraction agents
        self.board_text    = ""
        self.board_entries = []
        self.transcript    = ""
        self.segment_list  = []
        self.doc_text      = ""
        self.doc_entries   = []

        # populated by analysis agents
        self.summary       = ""

        # agent decision log
        self.agent_log     = []



class BaseAgent:
    """Every specialist agent inherits from this.
    Provides the agentic loop: should_run → plan → execute → validate → retry."""

    name        = "BaseAgent"
    max_retries = 2

    def should_run(self, ctx):
        """Return False to skip this agent entirely."""
        return True

    def plan(self, ctx):
        """Return a short string describing what this agent will do."""
        return "execute default task"

    def execute(self, ctx):
        """Do the actual work.  Must be implemented by subclasses."""
        raise NotImplementedError

    def validate(self, ctx):
        """Check output quality.  Return (ok: bool, issues: list[str])."""
        return True, []

    def run(self, ctx):
        """Full agent loop with retry on validation failure. Every attempt and
        the overall run are timed so every agent's speed shows up in the agent
        log for free (no per-subclass changes needed)."""
        if not self.should_run(ctx):
            self._log(ctx, "SKIP", "Not needed for this content type")
            return

        self._log(ctx, "PLAN", self.plan(ctx))

        run_start   = time.time()
        last_issues = []
        for attempt in range(1, self.max_retries + 1):
            attempt_start = time.time()
            try:
                self.execute(ctx)
                attempt_sec = time.time() - attempt_start
                ok, issues = self.validate(ctx)
                if ok:
                    total_sec = time.time() - run_start
                    self._log(
                        ctx, "DONE",
                        f"Completed on attempt {attempt} in {attempt_sec:.1f}s "
                        f"(total {total_sec:.1f}s)",
                        duration_sec=round(total_sec, 2),
                    )
                    return
                last_issues = issues
                self._log(
                    ctx, "RETRY",
                    f"Quality issues (attempt {attempt}, {attempt_sec:.1f}s): {issues}",
                    duration_sec=round(attempt_sec, 2),
                )
            except Exception as e:
                attempt_sec = time.time() - attempt_start
                self._log(
                    ctx, "ERROR",
                    f"Attempt {attempt} failed after {attempt_sec:.1f}s: {e}",
                    duration_sec=round(attempt_sec, 2),
                )
                if attempt == self.max_retries:
                    total_sec = time.time() - run_start
                    self._log(
                        ctx, "FAIL",
                        f"Gave up after {self.max_retries} attempts "
                        f"({total_sec:.1f}s total): {e}",
                        duration_sec=round(total_sec, 2),
                    )
                    return

        total_sec = time.time() - run_start
        self._log(
            ctx, "DONE", f"Completed with warnings in {total_sec:.1f}s: {last_issues}",
            duration_sec=round(total_sec, 2),
        )

    def _log(self, ctx, level, msg, duration_sec=None):
        entry = {
            "agent": self.name, "level": level,
            "time": datetime.now().isoformat(), "message": msg,
        }
        if duration_sec is not None:
            entry["duration_sec"] = duration_sec
        ctx.agent_log.append(entry)
        icons = {"SKIP": "⏭", "PLAN": "📋", "DONE": "✅",
                 "RETRY": "🔄", "ERROR": "❌", "FAIL": "💀"}
        print(f"  [{self.name}] {icons.get(level, '')} {msg}")



print("✅ Agent framework ready.")

# ══════════════════════════════════════════════════════════════
# SPECIALIST AGENTS
# ══════════════════════════════════════════════════════════════



class ExtractorAgent(BaseAgent):
    """Runs EasyOCR on video frames to extract board/slide text."""
    name = "Extractor"

    def should_run(self, ctx):
        return ctx.content_type == "video"

    def plan(self, ctx):
        return f"OCR scan frames from {ctx.source_name}"

    def execute(self, ctx):
        ctx.board_text, ctx.board_entries = extract_board_text(ctx.file_path)

    def validate(self, ctx):
        if not ctx.board_entries:
            return True, ["no board text found (may be camera-only video)"]
        return True, []



class TranscriberAgent(BaseAgent):
    """Transcribes audio from video using Faster-Whisper."""
    name = "Transcriber"

    def should_run(self, ctx):
        return ctx.content_type == "video"

    def plan(self, ctx):
        return f"Whisper transcription of {ctx.source_name}"

    def execute(self, ctx):
        global _whisper_model
        try:
            ctx.transcript, ctx.segment_list = transcribe_video(ctx.file_path)
        except Exception as e:
            err = str(e).lower()
            if "out of memory" in err or "cuda" in err:
                self._log(ctx, "RETRY", f"GPU OOM — resetting Whisper to try smaller model")
                _whisper_model = None  # force re-init with fallback cascade
                ctx.transcript, ctx.segment_list = transcribe_video(ctx.file_path)
            else:
                raise

    def validate(self, ctx):
        if not ctx.transcript or ctx.transcript == "No speech detected.":
            return False, ["empty transcript — audio extraction may have failed"]
        if len(ctx.transcript) < 50:
            return False, [f"transcript too short ({len(ctx.transcript)} chars)"]
        return True, []



class DocParserAgent(BaseAgent):
    """Extracts text from PDF / DOCX / PPTX documents."""
    name = "DocParser"

    def should_run(self, ctx):
        return ctx.content_type == "document"

    def plan(self, ctx):
        ext = Path(ctx.file_path).suffix.lower().lstrip('.')
        return f"Parse {ext.upper()} document: {ctx.source_name}"

    def execute(self, ctx):
        ext = Path(ctx.file_path).suffix.lower().lstrip('.')
        extractor = DOC_EXTRACTORS.get(ext)
        if not extractor:
            ctx.doc_entries = []
            ctx.doc_text = ""
            return
        try:
            ctx.doc_entries = extractor(ctx.file_path)
        except Exception as e:
            self._log(ctx, "ERROR", f"Failed to parse {ext.upper()}: {e}")
            ctx.doc_entries = []
        ctx.doc_text = "\n\n".join(e['text'] for e in ctx.doc_entries)

    def validate(self, ctx):
        if not ctx.doc_entries:
            return False, ["no text extracted — file may be corrupt or password-protected"]
        if len(ctx.doc_text) < 100:
            return False, [f"extracted text too short ({len(ctx.doc_text)} chars)"]
        return True, []



class SummaryAgent(BaseAgent):
    """Generates a structured summary using the LLM.  Validates length and quality."""
    name = "Summarizer"

    def plan(self, ctx):
        src = "board text + transcript" if ctx.content_type == "video" else "document text"
        return f"Generate structured summary from {src}"

    def execute(self, ctx):
        if ctx.content_type == "video":
            ctx.summary = generate_summary(
                ctx.board_text, ctx.transcript, ctx.segment_list, ctx.board_entries,
            )
        else:
            ctx.summary = call_ollama(
                f"Create a detailed, structured summary of this document:\n\n"
                f"{ctx.doc_text[:12000]}",
                system="You are an expert educational content summarizer.",
            )
        save_text("summary.txt", ctx.summary)

    def validate(self, ctx):
        issues = []
        if len(ctx.summary) < 200:
            issues.append(f"summary too short ({len(ctx.summary)} chars)")
        low_quality = ["i don't have enough", "i cannot", "no content provided"]
        if any(phrase in ctx.summary.lower() for phrase in low_quality):
            issues.append("summary appears to be a refusal or low-quality response")
        return len(issues) == 0, issues



class RAGBuilderAgent(BaseAgent):
    """Builds the FAISS search index from extracted content."""
    name = "RAGBuilder"

    def plan(self, ctx):
        return "Build FAISS vector index for RAG retrieval"

    def execute(self, ctx):
        if ctx.content_type == "video":
            combined = list(ctx.segment_list) + [
                {
                    "timestamp_str": e["timestamp_str"],
                    "text": f"[{e.get('tag', 'Visual')}] {e['text']}",
                    "start_str": e["timestamp_str"],
                    "tag": e.get("tag", "Visual"),
                }
                for e in (ctx.board_entries or [])
            ]
            build_rag(combined, ts_key="start_str", src_key="tag")
        else:
            build_rag(ctx.doc_entries, text_key="text")



class HintsAgent(BaseAgent):
    """Detects teacher emphasis moments and ranks topics by exam likelihood."""
    name = "HintsDetector"

    def should_run(self, ctx):
        return bool(ctx.summary)

    def plan(self, ctx):
        return f"Analyze {len(ctx.segment_list)} segments + summary for exam emphasis signals"

    def execute(self, ctx):
        hints     = detect_exam_hints(ctx.segment_list or [], ctx.board_entries or [])
        ai_topics = analyze_topics(ctx.summary, hints)
        state.quiz["hints"] = {
            "teacher_hints": hints,
            "ai_analysis": ai_topics,
            "total_hint_sentences": len(hints),
        }



# ══════════════════════════════════════════════════════════════
# TRUE AGENTIC QUIZ SYSTEM
# Planner → Generator → Validator → Refiner loop
# All agents push live thinking entries via thinking_cb.
# ══════════════════════════════════════════════════════════════


# ── Quiz Planner Agent ──────────────────────────────────────────

class QuizPlannerAgent(BaseAgent):
    """Reads user instructions + student performance + content summary
    and produces a structured plan (topics, counts, difficulties)
    that downstream Generator/Validator/Refiner agents follow."""
    name = "QuizPlannerAgent"
    max_retries = 2

    def plan(self, ctx):
        return "Analyzing instructions and performance to build a quiz plan"

    def execute(self, ctx):
        quiz_type  = getattr(ctx, '_quiz_type', 'mcq')
        num        = getattr(ctx, '_quiz_count', 10)
        mode       = getattr(ctx, '_mode', 'adaptive')
        difficulty = getattr(ctx, '_manual_difficulty', 'medium')
        user_instr = getattr(ctx, '_user_instructions', '')
        plan_mode  = getattr(ctx, '_plan_mode', 'auto')
        manual_topics = [t.strip() for t in (getattr(ctx, '_manual_topics', []) or []) if str(t).strip()]
        perf       = difficulty_adapter.get_all_stats()
        weak       = difficulty_adapter.get_weak_topics()

        # Manual topic mode: bypass LLM planner and force the selected topics.
        if plan_mode == "manual_topics" and manual_topics:
            topics = manual_topics[:8]
            base   = num // len(topics)
            rem    = num - base * len(topics)
            counts = {t: base for t in topics}
            counts[topics[0]] += rem
            diffs  = {t: (difficulty if mode == "manual" else "medium") for t in topics}
            ctx._quiz_plan = {
                "topics": topics,
                "per_topic_count": counts,
                "per_topic_difficulty": diffs,
                "bloom_targets": ["Remember", "Understand", "Apply"],
                "focus_notes": "Manual topic plan selected by user",
            }
            return

        perf_summary = ""
        if perf:
            lines = [
                f"  - {t}: {int(s['average']*100)}% ({s['level']})"
                for t, s in list(perf.items())[:8]
            ]
            perf_summary = "Student performance history:\n" + "\n".join(lines)

        summary_snip = (ctx.summary or state.status.get("summary", ""))[:2000]

        mode_note = (
            f"Target difficulty: {difficulty.upper()} (uniform, manual mode)"
            if mode == "manual"
            else "Mode: ADAPTIVE — assign lower difficulty to weak topics, higher to strong ones"
        )

        prompt = (
            f"You are a quiz planning agent. Decide how to distribute {num} "
            f"{quiz_type.upper()} questions across topics for a student.\n\n"
            f"Quiz type: {quiz_type.upper()}\n"
            f"{mode_note}\n"
            f"Total questions: {num}\n"
            f"Student instructions: {user_instr or 'None'}\n\n"
            f"{perf_summary}\n\n"
            f"Lecture summary (excerpt):\n{summary_snip}\n\n"
            "Return STRICT JSON (no extra text):\n"
            '{"topics":["topic1","topic2"],'
            '"per_topic_count":{"topic1":5,"topic2":5},'
            '"per_topic_difficulty":{"topic1":"easy","topic2":"medium"},'
            '"bloom_targets":["Remember","Apply"],'
            '"focus_notes":"what to emphasize or avoid based on instructions"}\n\n'
            "Rules:\n"
            f"- Sum of per_topic_count MUST equal {num}\n"
            "- ADAPTIVE: weak topics get more questions at easier difficulty\n"
            f"- MANUAL: all topics use {difficulty} difficulty\n"
            "- 'focus on X' in instructions → allocate more to X\n"
            "- 'avoid Y' in instructions → set Y count to 0\n"
            "- Use 2-4 topics maximum"
        )

        plan = call_ollama_json(
            prompt,
            system="You are an expert quiz planning agent. Return only valid JSON.",
            fallback=None,
            model=QUIZ_MODEL, num_ctx=QUIZ_CTX,
        )

        if not isinstance(plan, dict) or not plan.get("topics"):
            plan = self._fallback_plan(num, mode, difficulty, weak)

        # Ensure per_topic_count sums to num
        counts = plan.get("per_topic_count", {})
        if counts:
            total = sum(counts.values())
            if total != num:
                tlist = list(counts.keys())
                counts[tlist[-1]] = max(1, counts[tlist[-1]] + (num - total))
                plan["per_topic_count"] = counts

        ctx._quiz_plan = plan

    def _fallback_plan(self, num, mode, difficulty, weak):
        topics = list(weak.keys())[:3] if weak else ["Core Concepts", "Key Principles", "Applications"]
        topics = topics[:min(len(topics), 3)]
        base   = num // len(topics)
        rem    = num - base * len(topics)
        counts = {t: base for t in topics}
        counts[topics[0]] += rem
        diffs  = {
            t: (difficulty if mode == "manual"
                else ("easy" if weak.get(t, 1.0) < 0.4 else "medium"))
            for t in topics
        }
        return {
            "topics": topics,
            "per_topic_count": counts,
            "per_topic_difficulty": diffs,
            "bloom_targets": ["Remember", "Understand", "Apply"],
            "focus_notes": "Fallback plan — planner LLM returned invalid JSON",
        }

    def validate(self, ctx):
        if not getattr(ctx, '_quiz_plan', None) or not ctx._quiz_plan.get("topics"):
            return False, ["Planner produced no valid plan"]
        return True, []


# ── Quiz Validator Agent ─────────────────────────────────────────

class QuizValidatorAgent(BaseAgent):
    """Validates a batch of generated questions on 3 dimensions:
    (1) difficulty accuracy, (2) content grounding, (3) instruction compliance.
    Pushes per-question verdicts as thinking entries."""
    name = "QuizValidatorAgent"

    def validate_batch(self, questions, quiz_type, content_snippet,
                       target_difficulty, user_instructions, thinking_cb=None):
        """Returns (passed_questions, failed_with_reasons)."""
        if not questions:
            return [], []

        q_texts = []
        for i, q in enumerate(questions[:10]):
            txt = q.get("question") or q.get("statement") or ""
            ans = str(
                q.get("correct_answer") or q.get("answer") or q.get("model_answer") or ""
            )[:80]
            q_texts.append(f"{i+1}. Q: {txt[:200]}\n   A: {ans}")

        batch_text = "\n\n".join(q_texts)
        instr_note = (
            f"\nUser instructions to comply with: {user_instructions}"
            if user_instructions else ""
        )

        prompt = (
            f"You are a quiz quality validator. For each question apply 3 checks:\n"
            f"1. DIFFICULTY: Is it genuinely {target_difficulty.upper()}?\n"
            f"2. GROUNDED: Does it relate to the content snippet below?\n"
            f"3. INSTRUCTION: Does it follow the user's instructions?{instr_note}\n\n"
            f"Content snippet:\n{content_snippet[:700]}\n\n"
            f"Questions to validate:\n{batch_text}\n\n"
            'Return STRICT JSON:\n'
            '{"verdicts":[{"num":1,"difficulty":"PASS","grounded":"PASS",'
            '"instruction":"PASS","overall":"PASS","reason":"brief 1-line fix hint if FAIL"}]}'
        )

        result = call_ollama_json(
            prompt,
            system="You are a strict quiz validator. Return only JSON.",
            fallback={"verdicts": []},
            model=QUIZ_MODEL, num_ctx=QUIZ_CTX,
        )
        verdicts = result.get("verdicts", []) if isinstance(result, dict) else []

        passed, failed = [], []
        for i, q in enumerate(questions):
            v = next((v for v in verdicts if v.get("num") == i + 1), None)
            if v:
                overall = str(v.get("overall", "PASS")).upper()
                reason  = v.get("reason", "")
                entry = {
                    "agent":        self.name,
                    "step":         "validate_question",
                    "verdict":      overall,
                    "question_num": i + 1,
                    "reason":       reason,
                    "checks": {
                        "difficulty":  v.get("difficulty",  "PASS"),
                        "grounded":    v.get("grounded",    "PASS"),
                        "instruction": v.get("instruction", "PASS"),
                    },
                }
                if thinking_cb:
                    thinking_cb(entry)
                if overall == "PASS":
                    passed.append(q)
                else:
                    failed.append({**q, "_fail_reason": reason})
            else:
                # No verdict returned for this question → keep it
                passed.append(q)

        return passed, failed


# ── Quiz Refiner Agent ───────────────────────────────────────────

class QuizRefinerAgent(BaseAgent):
    """Re-writes questions that failed validation with targeted correction prompts.
    Only re-generates the failures, not the whole quiz."""
    name = "QuizRefinerAgent"

    def refine_batch(self, failed_questions, quiz_type, content_snippet,
                     target_difficulty, user_instructions, thinking_cb=None):
        """Returns refined question list replacing the failed ones."""
        if not failed_questions:
            return []

        schema     = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
        instr_note = f"\nStudent instructions: {user_instructions}" if user_instructions else ""

        items = []
        for i, q in enumerate(failed_questions):
            txt    = q.get("question") or q.get("statement") or ""
            reason = q.get("_fail_reason", "does not meet quality standard")
            items.append(f"{i+1}. [FAILED] Q: {txt[:200]}\n   Fix needed: {reason}")

        if thinking_cb:
            thinking_cb({
                "agent":   self.name,
                "step":    "refining",
                "verdict": "OVERRIDE",
                "reason":  f"Rewriting {len(failed_questions)} failed question(s) with targeted fixes…",
            })

        prompt = (
            f"Rewrite each failed question to fix the stated issue.\n\n"
            f"Requirements:\n"
            f"- Difficulty: {target_difficulty.upper()}\n"
            f"- Must be grounded in lecture content: {content_snippet[:500]}\n"
            f"{instr_note}\n\n"
            f"Failed questions ({len(failed_questions)}):\n"
            + "\n\n".join(items) + "\n\n"
            f"Return STRICT JSON with exactly {len(failed_questions)} rewritten question(s):\n{schema}"
        )

        data    = call_ollama_json_quiz(prompt, key="questions")
        refined = data.get("questions", []) if isinstance(data, dict) else []

        if thinking_cb:
            thinking_cb({
                "agent":   self.name,
                "step":    "refined",
                "verdict": "DONE",
                "reason":  f"Produced {len(refined)} refined replacement(s)",
            })

        return refined


# ── Per-question agentic loop ──────────────────────────────────
# Plan -> for each slot: Generate ONE -> Validate -> Retry until pass -> Next.
# Every transition pushes a structured event so the UI can render real-time.

DIFF_DESCRIPTIONS = {
    "easy":   "tests basic recall, definitions, single-concept understanding",
    "medium": "requires understanding + application, may combine 2 concepts",
    "hard":   "requires deep analysis, multi-step reasoning, edge cases, synthesis",
}


def _validate_one_question(quiz_type, question, content_snip, target_difficulty,
                            user_instr, existing_questions=None):
    """Run the validator on a SINGLE question.
    Returns (passed: bool, fail_reason: str, checks: dict).
    On failure, checks["action"] carries the validator's recommended
    remediation — one of 'regenerate' | 'fetch_context' | 'replan_topic' —
    which the generation loop acts on instead of always blindly regenerating."""
    existing_questions = existing_questions or []
    qtxt = question.get("question") or question.get("statement") or ""
    ans  = str(
        question.get("correct_answer") or question.get("answer") or
        question.get("model_answer") or ""
    )[:120]

    # Fill-blank: enforce blank marker before doing anything else
    if quiz_type == "fill" and "___" not in qtxt:
        return False, "fill question missing '___' blank marker", {
            "uniqueness": "PASS", "difficulty": "PASS",
            "grounded": "PASS", "instruction": "PASS", "action": "regenerate",
        }

    # cheap deterministic uniqueness check first
    qtxt_low = qtxt.strip().lower()
    if qtxt_low:
        for ex in existing_questions:
            ex_text = (ex.get("question") or ex.get("statement") or "").strip().lower()
            if not ex_text:
                continue
            if ex_text == qtxt_low:
                return False, "duplicate of an existing question", {
                    "uniqueness": "FAIL", "difficulty": "PASS",
                    "grounded": "PASS", "instruction": "PASS", "action": "regenerate",
                }
            ratio = difflib.SequenceMatcher(None, ex_text, qtxt_low).ratio()
            if ratio > 0.85:
                return False, f"too similar ({int(ratio*100)}%) to an existing question", {
                    "uniqueness": "FAIL", "difficulty": "PASS",
                    "grounded": "PASS", "instruction": "PASS", "action": "regenerate",
                }

    diff_desc  = DIFF_DESCRIPTIONS.get(target_difficulty, "")
    instr_note = (
        f"\nMust comply with student instructions: {user_instr}"
        if user_instr else ""
    )

    prompt = (
        "You are a strict but fair quiz quality validator. Apply 3 checks to ONE question:\n"
        f"1. DIFFICULTY: Is it genuinely {target_difficulty.upper()}? ({diff_desc})\n"
        "2. GROUNDED: Is it answerable from the content snippet below?\n"
        f"3. INSTRUCTION: Does it follow the student instructions?{instr_note}\n\n"
        f"Content snippet:\n{content_snip[:2500]}\n\n"
        "Question to validate:\n"
        f"Q: {qtxt[:300]}\n"
        f"A: {ans}\n\n"
        "If it FAILS, also choose the single best remediation 'action':\n"
        "- 'fetch_context': it fails because the snippet lacks the needed facts "
        "(GROUNDED fail) — more lecture context would let a good question be written.\n"
        "- 'replan_topic': the topic itself does not appear in the lecture at all — "
        "no amount of rewriting or context will help; a different topic is needed.\n"
        "- 'regenerate': wording/clarity/difficulty issues a plain rewrite can fix.\n\n"
        "Return STRICT JSON:\n"
        '{"difficulty":"PASS|FAIL","grounded":"PASS|FAIL",'
        '"instruction":"PASS|FAIL","overall":"PASS|FAIL",'
        '"action":"regenerate|fetch_context|replan_topic",'
        '"reason":"1-line concrete fix hint if FAIL"}'
    )

    result = call_ollama_json(
        prompt,
        system="You are a strict but fair quiz validator. Return only JSON.",
        fallback={"overall": "PASS"},
        model=QUIZ_MODEL, num_ctx=QUIZ_CTX,
    )
    if not isinstance(result, dict):
        return True, "", {
            "difficulty": "PASS", "grounded": "PASS",
            "instruction": "PASS", "uniqueness": "PASS", "action": "regenerate",
        }

    overall = str(result.get("overall", "PASS")).upper()
    reason  = result.get("reason", "")
    checks  = {
        "difficulty":  str(result.get("difficulty",  "PASS")).upper(),
        "grounded":    str(result.get("grounded",    "PASS")).upper(),
        "instruction": str(result.get("instruction", "PASS")).upper(),
        "uniqueness":  "PASS",
    }

    # Resolve the remediation action: trust the validator's pick when valid, but
    # apply deterministic guards so the action matches the actual failure type.
    action = str(result.get("action", "")).strip().lower()
    if action not in ("regenerate", "fetch_context", "replan_topic"):
        action = ""
    if checks["grounded"] == "FAIL":
        action = action or "fetch_context"      # grounding gap → get more context
    elif checks["difficulty"] == "FAIL":
        action = "regenerate"                    # difficulty is a rewrite problem
    checks["action"] = action or "regenerate"

    # Enforce hard blocking checks deterministically:
    # - grounded must pass
    # - instruction must pass only when user gave instructions
    hard_fail = []
    if checks["grounded"] == "FAIL":
        hard_fail.append("grounded")
    if user_instr and checks["instruction"] == "FAIL":
        hard_fail.append("instruction")
    if checks["uniqueness"] == "FAIL":
        hard_fail.append("uniqueness")
    if hard_fail:
        return False, reason or f"failed checks: {', '.join(hard_fail)}", checks

    # Difficulty is subjective and the small validator model is noisy on it,
    # which used to be the main cause of questions burning all 3 attempts. Treat
    # a difficulty mismatch as NON-blocking at every level (not just HARD) as
    # long as the question is grounded and instruction-safe — accept it with a
    # note rather than exhausting retries on a judgement call.
    if checks["difficulty"] == "FAIL":
        return True, "difficulty borderline; accepted (grounded + instruction-safe)", checks

    # Reaching here means every concrete check passed (grounded + instruction +
    # uniqueness, with difficulty handled above). Trust those objective gates
    # rather than the model's vague holistic "overall" flag, which on the small
    # validator is noisy and caused needless retries when it contradicted the
    # per-dimension verdicts.
    return True, reason, checks


def _fetch_more_context(topic, question_text, top_k=5):
    """Remediation tool: pull extra lecture context from the RAG index for a
    topic the validator flagged as under-grounded. Returns a context block, or
    '' if no useful index/context is available."""
    try:
        query = (str(topic) + " " + (question_text or "")).strip()
        if not query:
            return ""
        ctx_str, _ = rag.query(query, top_k=top_k)
        if not ctx_str or ctx_str.strip().lower().startswith("no index"):
            return ""
        return "=== RETRIEVED CONTEXT (for grounding) ===\n" + ctx_str[:2500]
    except Exception:
        return ""


def _pick_grounded_topic(quiz_type, content_snip, summary, avoid_topics):
    """Remediation tool: ask the planner model for a different topic that is
    actually covered in the lecture, when the validator says the current topic
    isn't in the material. Returns a topic string, or '' on failure."""
    avoid = ", ".join(t for t in avoid_topics if t) or "none"
    prompt = (
        f"Pick ONE specific topic that is clearly and explicitly covered in the "
        f"lecture content below and is well-suited for a {quiz_type.upper()} question.\n"
        f"Avoid these already-used or unsuitable topics: {avoid}.\n"
        'Return STRICT JSON: {"topic":"..."}\n\n'
        f"CONTENT:\n{content_snip[:3500]}\n\nSUMMARY:\n{summary[:800]}"
    )
    data = call_ollama_json(prompt, fallback={}, model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
    if isinstance(data, dict):
        return (data.get("topic") or "").strip()
    return ""


def _topic_coverage(topic, content_snip):
    """Fraction of a topic's significant words that appear in the content — a
    cheap gate so the tool-use grounding check only fires when the handed-in
    content looks thin on the topic (avoids an LLM call for well-covered ones)."""
    words = [w for w in re.split(r'\W+', (topic or '').lower()) if len(w) > 3]
    if not words:
        return 1.0
    low = content_snip.lower()
    return sum(1 for w in words if w in low) / len(words)


def _gather_grounding(quiz_type, topic, content_snip, summary, max_tool_calls=2,
                      thinking_cb=None, qid=None, agent_name="QuizGenerator"):
    """Tool-use loop run BEFORE drafting: the generator inspects the content it
    was handed and decides for itself whether it has enough material on the
    topic. If not, it calls its search_lecture(query) tool (RAG) — choosing the
    query — to pull more context first. Returns the (possibly enriched) content.

    This is genuine model-directed tool use: the LLM decides whether to call the
    tool and with what arguments, rather than always being fed a fixed snippet."""
    if rag is None or getattr(rag, "index", None) is None:
        return content_snip                                   # no RAG tool available
    if _topic_coverage(topic, content_snip) >= 0.6:
        return content_snip                                   # already well covered

    enriched = content_snip
    for _ in range(max(1, max_tool_calls)):
        prompt = (
            f"You will write a {quiz_type.upper()} question on the topic '{topic}'.\n"
            "Tool available — search_lecture(query): returns more text from THIS lecture.\n"
            "Does the CONTENT below already contain enough specific material on that "
            "topic to write a well-grounded question?\n"
            '- If yes: return {"action":"ready"}\n'
            '- If not: return {"action":"search","query":"<short search query>"}\n\n'
            f"CONTENT:\n{enriched[:2800]}\n\nSUMMARY:\n{summary[:500]}\n\n"
            'Return STRICT JSON: {"action":"ready|search","query":"..."}'
        )
        data   = call_ollama_json(prompt, fallback={"action": "ready"},
                                  model=QUIZ_MODEL, num_ctx=QUIZ_CTX)
        action = str((data or {}).get("action", "ready")).strip().lower()
        query  = str((data or {}).get("query", "")).strip()
        if action != "search" or not query:
            break
        extra = _fetch_more_context(query, "")
        if thinking_cb:
            thinking_cb({
                "agent":       agent_name,
                "step":        "tool_call",
                "phase":       "GENERATE",
                "question_id": qid,
                "verdict":     "OVERRIDE",
                "action":      "rag_query",
                "reason":      (f"Generator searched the lecture for \"{query[:60]}\" "
                                "before drafting" + ("" if extra else " — nothing new found")),
            })
        if not extra or extra[:160] in enriched:
            break
        enriched = extra + "\n\n" + enriched
    return enriched


def _generate_one_question_with_retry(
    quiz_type, topic, difficulty, content_snip, summary,
    user_instr, existing_questions, question_idx,
    max_attempts=3, thinking_cb=None, agent_name="QuizGenerator",
):
    """Generate ONE question, validate it, and on failure let the VALIDATOR
    choose the remediation — regenerate, fetch more RAG context, or replan the
    topic — rather than blindly retrying. Pushes per-question events so the UI
    can render an attempt-by-attempt view."""
    qid                  = f"q{question_idx + 1}"
    q_start              = time.time()
    last_failure_reason  = ""
    last_question        = None
    best_question        = None   # best of the failed attempts, if none ever passes
    best_score           = -1
    best_reason          = ""
    schema               = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
    diff_desc            = DIFF_DESCRIPTIONS.get(difficulty, "")
    strategy_tracks      = {
        "mcq": [
            "Concept check style: test one core concept precisely with plausible distractors.",
            "Scenario style: short practical situation and ask the best next reasoning choice.",
            "Compare-and-justify style: ask which option is best and why others are weaker.",
        ],
        "tf": [
            "Direct claim style: one clear factual statement from content.",
            "Conditional claim style: statement involving if/when constraints from lecture.",
            "Misconception trap style: plausible but subtle statement requiring careful reading.",
        ],
        "fill": [
            "Keyword recall style: blank out a key term or phrase from lecture concepts.",
            "Process step style: blank out a critical step/condition in a method.",
            "Contrast style: blank that distinguishes two similar concepts.",
        ],
        "short": [
            "Explain-why style: ask for concise reasoning of a concept decision.",
            "Apply-to-case style: apply concept to a concrete mini-case.",
            "Trade-off style: discuss a clear trade-off grounded in lecture points.",
        ],
    }
    strategies = strategy_tracks.get(quiz_type, strategy_tracks["mcq"])

    # Tool-use: before drafting, let the generator decide whether it has enough
    # grounding for this topic and, if not, issue its own RAG search to pull more
    # context. The validator can still request a reactive fetch later, so this
    # proactive pass does NOT consume the fetched_context budget.
    working_content = _gather_grounding(
        quiz_type, topic, content_snip, summary,
        max_tool_calls=2, thinking_cb=thinking_cb, qid=qid, agent_name=agent_name,
    )

    # Remediation state — between attempts the validator can enrich the content
    # (fetch_context) or switch the topic (replan_topic) instead of just retrying.
    fetched_context = False
    replanned       = False

    for attempt in range(1, max_attempts + 1):
        strategy = strategies[(attempt - 1) % len(strategies)]
        if thinking_cb:
            thinking_cb({
                "agent":        agent_name,
                "step":         "question_generating",
                "phase":        "GENERATE",
                "question_id":  qid,
                "question_idx": question_idx,
                "topic":        topic,
                "difficulty":   difficulty,
                "attempt":      attempt,
                "max_attempts": max_attempts,
                "feedback":     last_failure_reason if attempt > 1 else "",
                "verdict":      "RETRY" if attempt > 1 else "PASS",
                "reason": (
                    f"Generating Q{question_idx + 1} on '{topic}' "
                    f"at {difficulty.upper()} difficulty…"
                    if attempt == 1
                    else f"Retry {attempt}/{max_attempts} — fixing: {last_failure_reason}"
                ),
                "strategy": strategy,
            })

        feedback_block = ""
        if last_failure_reason and attempt > 1:
            feedback_block = (
                f"\nPREVIOUS ATTEMPT FAILED: {last_failure_reason}\n"
                "You MUST fix this exact issue in the new question.\n"
            )

        existing_summary = ""
        if existing_questions:
            existing_texts = [
                (q.get("question") or q.get("statement", ""))[:90]
                for q in existing_questions
            ]
            existing_summary = (
                "\nDO NOT duplicate or paraphrase these accepted questions:\n- "
                + "\n- ".join(existing_texts) + "\n"
            )

        instr_line = f"\nStudent instructions: {user_instr}\n" if user_instr else ""

        prompt = (
            f"Generate exactly ONE {quiz_type.upper()} question wrapped in a JSON array.\n"
            f"Topic: '{topic}'\n"
            f"Difficulty: {difficulty.upper()} — {diff_desc}\n"
            f"Question strategy for this attempt: {strategy}\n"
            "GROUNDING RULE: the question and its answer MUST be answerable using ONLY "
            "the CONTENT below. Do NOT introduce formulas, methods, numbers, or concepts "
            "that are not present in it — not even to make it harder. Make it harder by "
            "deeper reasoning over the SAME material, never by adding outside topics.\n"
            f"{feedback_block}"
            f"{existing_summary}"
            f"{instr_line}"
            f"Output STRICT JSON in this exact shape (a 'questions' array containing exactly 1 object):\n"
            f"{schema}\n\n"
            f"CONTENT:\n{working_content[:5500]}\n\n"
            f"SUMMARY:\n{summary[:1200]}"
        )

        # Speed: the first attempt uses the fast quiz model (4B) and a capped
        # context; only escalate to the larger primary model on a retry, where
        # the fast model has already failed validation and the extra quality is
        # worth the latency. A 32K context on a T4 is both slow and OOM-prone
        # next to Whisper/OCR, and this prompt is small, so cap it to QUIZ_CTX.
        gen_model = QUIZ_MODEL if attempt == 1 else state.active_model
        data = call_ollama_json(prompt, fallback={"questions": []},
                                model=gen_model, num_ctx=QUIZ_CTX)
        qs   = data.get("questions", []) if isinstance(data, dict) else []
        if not qs:
            repaired = call_ollama_json_quiz(prompt, key="questions")
            qs = repaired.get("questions", []) if isinstance(repaired, dict) else []

        if not qs:
            last_failure_reason = "model returned no question"
            if thinking_cb:
                thinking_cb({
                    "agent":       "QuizValidatorAgent",
                    "step":        "question_validated",
                    "phase":       "VALIDATE",
                    "question_id": qid,
                    "verdict":     "FAIL",
                    "attempt":     attempt,
                    "reason":      last_failure_reason,
                })
            continue

        q = qs[0]
        if not q.get("topic"):      q["topic"]      = topic
        if not q.get("difficulty"): q["difficulty"] = difficulty
        last_question = q

        passed, fail_reason, checks = _validate_one_question(
            quiz_type, q, working_content, difficulty, user_instr,
            existing_questions=existing_questions,
        )

        qtxt = q.get("question") or q.get("statement") or ""

        if passed:
            if thinking_cb:
                thinking_cb({
                    "agent":       "QuizValidatorAgent",
                    "step":        "question_validated",
                    "phase":       "VALIDATE",
                    "question_id": qid,
                    "verdict":     "PASS",
                    "attempt":     attempt,
                    "checks":      checks,
                    "reason":      "All quality checks passed",
                })
                q_duration_sec = round(time.time() - q_start, 2)
                thinking_cb({
                    "agent":            agent_name,
                    "step":             "question_finalized",
                    "phase":            "GENERATE",
                    "question_id":      qid,
                    "question_idx":     question_idx,
                    "verdict":          "DONE",
                    "status":           "PASS",
                    "total_attempts":   attempt,
                    "topic":            topic,
                    "difficulty":       difficulty,
                    "question_preview": qtxt[:160],
                    "duration_sec":     q_duration_sec,
                    "reason":           f"Q{question_idx + 1} accepted on attempt {attempt} ({q_duration_sec:.1f}s)",
                })
            return q

        last_failure_reason = fail_reason or "validation failed"
        # Track the best of the failed attempts (most checks passing) so that if
        # we do exhaust retries we keep the strongest candidate, not just the
        # last one generated.
        score = sum(1 for k in ("grounded", "instruction", "difficulty", "uniqueness")
                    if (checks or {}).get(k) == "PASS")
        if score > best_score:
            best_score, best_question, best_reason = score, q, last_failure_reason
        if thinking_cb:
            thinking_cb({
                "agent":       "QuizValidatorAgent",
                "step":        "question_validated",
                "phase":       "VALIDATE",
                "question_id": qid,
                "verdict":     "FAIL",
                "attempt":     attempt,
                "checks":      checks,
                "reason":      last_failure_reason,
            })

        # ── Validator-chosen remediation ──────────────────────────
        # The validator decides HOW to fix, not just that it failed:
        #   fetch_context → pull more lecture context from RAG, then retry
        #   replan_topic  → switch to a topic that's actually in the lecture
        #   regenerate    → plain rewrite (feedback already in the prompt)
        # Only act if another attempt remains.
        if attempt < max_attempts:
            action = (checks or {}).get("action", "regenerate")
            # Escalate when a remediation has already been tried and still fails.
            if action == "fetch_context" and fetched_context:
                action = "replan_topic"
            if action == "replan_topic" and replanned:
                action = "regenerate"

            if action == "fetch_context":
                extra = _fetch_more_context(topic, qtxt)
                fetched_context = True   # attempted — escalate to replan if it fails again
                if extra and extra[:160] not in working_content:
                    working_content = (extra + "\n\n" + working_content)
                    if thinking_cb:
                        thinking_cb({
                            "agent":       "QuizValidatorAgent",
                            "step":        "remediation",
                            "phase":       "VALIDATE",
                            "question_id": qid,
                            "attempt":     attempt,
                            "verdict":     "OVERRIDE",
                            "action":      "fetch_context",
                            "reason":      f"Under-grounded — fetched extra RAG context for '{topic}', retrying",
                        })

            elif action == "replan_topic":
                avoid = [topic] + [(q2.get("topic") or "") for q2 in (existing_questions or [])]
                new_topic = _pick_grounded_topic(quiz_type, content_snip, summary, avoid)
                if new_topic and new_topic.strip().lower() != (topic or "").strip().lower():
                    if thinking_cb:
                        thinking_cb({
                            "agent":       "QuizValidatorAgent",
                            "step":        "remediation",
                            "phase":       "VALIDATE",
                            "question_id": qid,
                            "attempt":     attempt,
                            "verdict":     "OVERRIDE",
                            "action":      "replan_topic",
                            "reason":      f"Topic '{topic}' not in lecture — replanning to '{new_topic}'",
                        })
                    topic               = new_topic
                    replanned           = True
                    last_failure_reason = ""   # fresh topic — don't carry stale feedback

    # Max attempts reached without a clean pass: keep the BEST of the attempts
    # (most checks passing) so the set still fills, but label it honestly as an
    # OVERRIDE — NOT a PASS — so "max tries reached" is never silently passed off
    # as if validation succeeded. The UI shows these amber and the metrics count
    # them as overrides.
    kept = best_question if best_question is not None else last_question
    if kept is not None:
        if thinking_cb:
            q_duration_sec = round(time.time() - q_start, 2)
            thinking_cb({
                "agent":            agent_name,
                "step":             "question_finalized",
                "phase":            "GENERATE",
                "question_id":      qid,
                "question_idx":     question_idx,
                "verdict":          "OVERRIDE",
                "status":           "OVERRIDE",
                "total_attempts":   max_attempts,
                "topic":            topic,
                "difficulty":       difficulty,
                "question_preview": (kept.get("question") or kept.get("statement") or "")[:160],
                "duration_sec":     q_duration_sec,
                "reason":           f"Validator not fully satisfied after {max_attempts} attempts "
                                    f"({q_duration_sec:.1f}s) — kept best candidate. Last issue: {best_reason or last_failure_reason}",
            })
        return kept

    # Rare case: model produced no question in all attempts. Do a strict JSON recovery
    # and accept the recovered output as final attempt content.
    recovery_prompt = (
        f"Return exactly ONE valid {quiz_type.upper()} question in STRICT JSON:\n"
        f"{schema}\n"
        f"Topic: {topic}\nDifficulty: {difficulty.upper()}\n"
        f"Do not return explanations outside JSON.\n"
        f"CONTENT:\n{content_snip[:4500]}\n\nSUMMARY:\n{summary[:1000]}"
    )
    for _ in range(3):
        repaired = call_ollama_json_quiz(recovery_prompt, key="questions")
        rec_qs = repaired.get("questions", []) if isinstance(repaired, dict) else []
        if not rec_qs:
            continue
        q = rec_qs[0]
        if not q.get("topic"):      q["topic"] = topic
        if not q.get("difficulty"): q["difficulty"] = difficulty
        if thinking_cb:
            q_duration_sec = round(time.time() - q_start, 2)
            thinking_cb({
                "agent":            agent_name,
                "step":             "question_finalized",
                "phase":            "GENERATE",
                "question_id":      qid,
                "question_idx":     question_idx,
                "verdict":          "OVERRIDE",
                "status":           "OVERRIDE",
                "total_attempts":   max_attempts,
                "topic":            topic,
                "difficulty":       difficulty,
                "question_preview": (q.get("question") or q.get("statement") or "")[:160],
                "duration_sec":     q_duration_sec,
                "reason":           f"Recovered after {max_attempts} attempts returned no question — "
                                    f"kept without full validation ({q_duration_sec:.1f}s)",
            })
        return q

    raise RuntimeError(
        f"Failed to generate any {quiz_type.upper()} question text after {max_attempts} attempts for slot {question_idx + 1}"
    )


def _run_quiz_agentic_loop(ctx, quiz_type, content, summary,
                            user_instr, target_difficulty, thinking_cb=None,
                            max_attempts_per_q=3):
    """Per-question agentic loop:
    Phase 1 PLAN -> Phase 2 GENERATE one + Phase 3 VALIDATE + retry -> Phase 4 DONE.
    Each event has phase/question_id/attempt/verdict so the UI can show a live pipeline."""
    agent_name = getattr(ctx, '_agent_name', 'QuizAgent')
    num        = getattr(ctx, '_quiz_count', 10)

    # ── PHASE 1: PLAN ──────────────────────────────────────────
    if thinking_cb:
        thinking_cb({
            "agent":   "Orchestrator",
            "step":    "phase_start",
            "phase":   "PLAN",
            "verdict": "PASS",
            "reason":  "Phase 1/4 — Planner deciding topic distribution & difficulty…",
        })

    plan = getattr(ctx, '_quiz_plan', None)
    if not plan:
        QuizPlannerAgent().run(ctx)
        plan = getattr(ctx, '_quiz_plan', None) or {}

    topics      = plan.get("topics", [])
    topic_cnt   = plan.get("per_topic_count",      {})
    topic_diff  = plan.get("per_topic_difficulty", {})
    focus_notes = plan.get("focus_notes", "")

    if thinking_cb:
        thinking_cb({
            "agent":   "QuizPlannerAgent",
            "step":    "plan_ready",
            "phase":   "PLAN",
            "verdict": "DONE",
            "reason":  f"Plan ready: {len(topics)} topic(s) over {num} question(s)",
            "data": {
                "topics":               topics,
                "per_topic_count":      topic_cnt,
                "per_topic_difficulty": topic_diff,
                "focus_notes":          focus_notes,
            },
        })

    # build flat sequence of (topic, difficulty) slots
    slots = []
    if topics and topic_cnt:
        for topic in topics:
            count = topic_cnt.get(topic, 0)
            diff  = topic_diff.get(topic, target_difficulty)
            for _ in range(count):
                slots.append((topic, diff))
    if not slots:
        slots = [("General", target_difficulty)] * num

    # respect num exactly: trim or pad
    if len(slots) > num:
        slots = slots[:num]
    while len(slots) < num:
        last = slots[-1] if slots else ("General", target_difficulty)
        slots.append(last)

    # ── PHASE 2 + 3: GENERATE + VALIDATE per question ──────────
    if thinking_cb:
        thinking_cb({
            "agent":   agent_name,
            "step":    "phase_start",
            "phase":   "GENERATE",
            "verdict": "PASS",
            "reason": (
                f"Phase 2/4 — Generating {len(slots)} question(s) one at a time, "
                f"validating each, retrying on fail (max {max_attempts_per_q} attempts each)…"
            ),
            "data":    {"total_slots": len(slots)},
        })

    final_questions = []
    content_snip    = content[:6000]

    for idx, (topic, diff) in enumerate(slots):
        # One impossible slot must not abort the whole batch — isolate failures
        # per question so the rest of the set still gets generated.
        try:
            q = _generate_one_question_with_retry(
                quiz_type=quiz_type, topic=topic, difficulty=diff,
                content_snip=content_snip, summary=summary,
                user_instr=user_instr,
                existing_questions=final_questions,
                question_idx=idx,
                max_attempts=max_attempts_per_q,
                thinking_cb=thinking_cb,
                agent_name=agent_name,
            )
        except Exception as e:
            q = None
            if thinking_cb:
                thinking_cb({
                    "agent":        agent_name,
                    "step":         "question_finalized",
                    "phase":        "GENERATE",
                    "question_id":  f"q{idx + 1}",
                    "question_idx": idx,
                    "verdict":      "FAIL",
                    "status":       "FAIL",
                    "topic":        topic,
                    "difficulty":   diff,
                    "reason":       f"Q{idx + 1} skipped — generation failed: {e}",
                })
        if q:
            final_questions.append(q)

    # ── PHASE 4: DONE ─────────────────────────────────────────
    if thinking_cb:
        thinking_cb({
            "agent":   agent_name,
            "step":    "phase_complete",
            "phase":   "DONE",
            "verdict": "DONE",
            "reason":  f"All phases complete — {len(final_questions)} question(s) ready",
            "data":    {"total": len(final_questions)},
        })

    return final_questions


# ── Manual mode difficulty agents ────────────────────────────────

class _DifficultySubAgent(BaseAgent):
    """Base for Easy/Medium/Hard quiz agents.
    Delegates to the per-question agentic loop at a fixed target difficulty."""
    target_difficulty = "medium"
    max_retries       = 2
    _thinking_log     = []

    def _get_content(self, ctx):
        if ctx.content_type == "video":
            return pack_video_sources(ctx.board_text, ctx.transcript,
                                      ctx.board_entries or [], 7500)
        return (ctx.doc_text or "")[:7500]

    def _build_prompt(self, quiz_type, num, content, summary, user_instructions=""):
        """Backward-compat helper kept for the legacy refine path in /quiz/generate."""
        diff   = self.target_difficulty
        instr  = f"\nAdditional instructions: {user_instructions}\n" if user_instructions else ""
        schema = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
        return (
            f"Generate exactly {num} {quiz_type.upper()} questions at {diff.upper()} difficulty.\n"
            f"{instr}"
            f"STRICT JSON: {schema}\n\n"
            f"CONTENT:\n{content}\n\nSUMMARY:\n{summary[:1500]}"
        )

    def _validate_difficulty(self, questions, quiz_type):
        """Backward-compat helper — uses the per-question validator."""
        if not questions:
            return [], []
        kept, thinking = [], []
        for i, q in enumerate(questions):
            ok, reason, checks = _validate_one_question(
                quiz_type, q, "", self.target_difficulty, "",
                existing_questions=questions[:i],
            )
            thinking.append({
                "agent":   self.name,
                "step":    "difficulty_check",
                "verdict": "PASS" if ok else "FAIL",
                "checks":  checks,
                "reason":  reason or "passed",
            })
            if ok:
                kept.append(q)
        if not kept:
            thinking.append({
                "agent": self.name, "step": "override",
                "verdict": "OVERRIDE",
                "reason": "All questions rejected by validator — keeping originals",
            })
            return questions, thinking
        return kept, thinking


class EasyQuizAgent(_DifficultySubAgent):
    name = "EasyQuizAgent"
    target_difficulty = "easy"

    def plan(self, ctx):
        return "Generate easy-level questions (recall, definitions, basics)"

    def should_run(self, ctx):
        return getattr(ctx, '_manual_difficulty', '') == 'easy'


class MediumQuizAgent(_DifficultySubAgent):
    name = "MediumQuizAgent"
    target_difficulty = "medium"

    def plan(self, ctx):
        return "Generate medium-level questions (understanding, application)"

    def should_run(self, ctx):
        return getattr(ctx, '_manual_difficulty', '') == 'medium'


class HardQuizAgent(_DifficultySubAgent):
    name = "HardQuizAgent"
    target_difficulty = "hard"

    def plan(self, ctx):
        return "Generate hard-level questions (analysis, synthesis, edge cases)"

    def should_run(self, ctx):
        return getattr(ctx, '_manual_difficulty', '') == 'hard'


# ── Adaptive Quiz Agent ──────────────────────────────────────────

class AdaptiveQuizAgent(BaseAgent):
    """Reads student performance + content, then runs the per-question agentic
    loop with planner-driven topic + difficulty allocation."""
    name          = "AdaptiveQuizAgent"
    max_retries   = 2
    _thinking_log = []

    def _add_thinking(self, ctx, entry):
        self._thinking_log.append(entry)
        cb = getattr(ctx, "thinking_cb", None)
        if cb:
            cb(entry)

    def plan(self, ctx):
        weak = difficulty_adapter.get_weak_topics()
        if weak:
            topics = ", ".join(f"{t} ({int(s*100)}%)" for t, s in weak.items())
            return f"Adaptive quiz targeting weak topics: {topics}"
        return "No history yet — generating balanced quiz across all topics"

    def execute(self, ctx):
        self._thinking_log = []
        quiz_type  = getattr(ctx, '_quiz_type', 'mcq')
        user_instr = getattr(ctx, '_user_instructions', '')
        perf       = difficulty_adapter.get_all_stats()
        weak       = difficulty_adapter.get_weak_topics()

        self._add_thinking(ctx, {
            "agent":   self.name,
            "step":    "analyze_performance",
            "phase":   "PLAN",
            "verdict": "PASS",
            "data":    {"all_stats": perf, "weak_topics": weak},
            "reason": (
                f"Found {len(weak)} weak topic(s) out of {len(perf)} tracked"
                if perf else "No performance history yet — will generate balanced quiz"
            ),
        })

        content = (
            pack_video_sources(ctx.board_text, ctx.transcript, ctx.board_entries or [], 7500)
            if ctx.content_type == "video"
            else (ctx.doc_text or "")[:7500]
        )
        summary = ctx.summary or state.status.get("summary", "")

        ctx._mode       = "adaptive"
        ctx._agent_name = self.name

        def thinking_cb(entry):
            self._add_thinking(ctx, entry)

        final = _run_quiz_agentic_loop(
            ctx, quiz_type, content, summary, user_instr,
            target_difficulty="medium",
            thinking_cb=thinking_cb,
        )

        state.quiz[quiz_type]      = final
        state._quiz_agent_thinking = self._thinking_log

    def validate(self, ctx):
        quiz_type = getattr(ctx, '_quiz_type', 'mcq')
        qs = state.quiz.get(quiz_type, [])
        if not qs:
            return False, [f"Adaptive agent produced 0 {quiz_type} questions"]
        return True, []



class FlashcardAgent(BaseAgent):
    """Agentic flashcard generator:
    1. Generates cards in batches of 8
    2. Validates each card (front != back, non-empty, not duplicate)
    3. Retries rejected cards up to max_retries times
    4. Stops when 15 unique valid cards are collected or max attempts hit."""
    name        = "FlashcardGen"
    max_retries = 3
    TARGET      = 15

    def plan(self, ctx):
        return f"Agentic flashcard pipeline: generate → validate → retry until {self.TARGET} unique cards"

    def _content(self, ctx):
        if ctx.content_type == "video":
            vis = format_visual_lecture_text(ctx.board_text, ctx.board_entries or [], 4000)
            parts = []
            if vis:
                parts.append(f"=== SLIDES + BOARDS ===\n{vis}")
            parts.append(f"=== SUMMARY ===\n{(ctx.summary or '')[:3500]}")
            return "\n\n".join(parts)[:6500]
        doc = (ctx.doc_text or "")[:4000]
        return (doc + "\n\n=== SUMMARY ===\n" + (ctx.summary or "")[:2500])[:6500]

    def _validate_card(self, card, seen_fronts):
        """Return (ok, reason). Deterministic checks only."""
        if not isinstance(card, dict):
            return False, "not an object"
        front = (card.get("front") or "").strip()
        back  = (card.get("back")  or "").strip()
        if not front or not back:
            return False, "empty front or back"
        if front.lower() == back.lower():
            return False, "front == back"
        if len(back) < 8:
            return False, "back too short to be useful"
        if front.lower() in seen_fronts:
            return False, "duplicate front"
        return True, ""

    def execute(self, ctx):
        content     = self._content(ctx)
        seen_fronts = set()
        accepted    = []
        batch_num   = 0
        # Flashcards are a bulk JSON-array task the small 4B model can't do, so we
        # generate straight on the primary (12B) model, which handles a full batch
        # in one shot — a couple of batches is plenty.
        max_batches = self.max_retries + 2

        while len(accepted) < self.TARGET and batch_num < max_batches:
            batch_num   += 1
            need         = self.TARGET - len(accepted)
            ask          = min(need + 3, 10)    # ask a few extra to absorb rejects
            existing_txt = (
                "\nALREADY ACCEPTED (do not repeat):\n- "
                + "\n- ".join(c["front"] for c in accepted)
                if accepted else ""
            )
            prompt = (
                f"Generate exactly {ask} flashcards for key concepts from this lecture.\n"
                f"Batch {batch_num} — need {need} more accepted cards.\n"
                f"{existing_txt}\n\n"
                "STRICT JSON: {\"cards\":[{\"front\":\"term or question\","
                "\"back\":\"clear explanation (≥10 words)\","
                "\"topic\":\"\",\"difficulty\":\"easy|medium|hard\"}]}\n"
                "RULES: front != back. back must be at least 10 words. "
                "Each card must cover a different concept.\n\n"
                f"CONTENT:\n{content}"
            )
            data  = call_ollama_json_quiz(prompt, key="cards", force_primary=True)
            batch = data.get("cards", []) if isinstance(data, dict) else []

            for card in batch:
                ok, reason = self._validate_card(card, seen_fronts)
                if ok:
                    seen_fronts.add(card["front"].strip().lower())
                    accepted.append(card)
                    if len(accepted) >= self.TARGET:
                        break
                else:
                    print(f"  [FlashcardGen] Rejected card: {reason} — '{card.get('front','')[:50]}'")

        state.quiz["flash"] = accepted
        print(f"  [FlashcardGen] ✅ {len(accepted)} unique valid flashcards generated in {batch_num} batch(es)")

    def validate(self, ctx):
        cards = state.quiz.get("flash", [])
        if not cards:
            return False, ["FlashcardAgent produced 0 cards"]
        issues = []
        if len(cards) < 8:
            issues.append(f"only {len(cards)} cards — expected ~{self.TARGET}")
        return True, issues



class SuggestedQuestionsAgent(BaseAgent):
    """Generates suggested chat questions for the tutor."""
    name = "SuggestQuestions"

    def plan(self, ctx):
        return "Generate 5 suggested study questions for the chat tutor"

    def execute(self, ctx):
        if ctx.content_type == "video":
            head = pack_video_sources(
                ctx.board_text, ctx.transcript, ctx.board_entries or [], 5500,
            )
            head = f"{head}\n\n=== SUMMARY ===\n{ctx.summary[:2200]}"
        else:
            doc = (ctx.doc_text or "")[:4000]
            head = f"{doc}\n\n=== SUMMARY ===\n{ctx.summary[:2200]}" if doc else ctx.summary[:4000]
        sq = call_ollama_json_list_quiz(
            "Generate 5 study questions a student might ask about this lecture "
            "(slides, boards, speech). Return a JSON list of strings:\n"
            f"{head[:7800]}",
            fallback_list=[],
            force_primary=True,
        )
        state.suggested_questions = sq if isinstance(sq, list) else []



print("✅ Specialist agents ready.")

# ══════════════════════════════════════════════════════════════
# EVALUATION HARNESS  —  LLM-as-Judge + Agentic vs Non-Agentic
# ══════════════════════════════════════════════════════════════
#
# Compares the per-question agentic pipeline (_run_quiz_agentic_loop) against the
# single-shot non-agentic baseline generators (generate_mcq/tf/fill/short) on
# identical lecture content. Every generated question, from both pipelines, is
# scored by an independent judge LLM (the larger OLLAMA_MODEL, distinct from the
# QUIZ_MODEL used for all generation/validation). Call run_quiz_evaluation() from
# a Colab cell after processing a real lecture.



JUDGE_SCHEMA = (
    '{"correctness":1,"clarity":1,"difficulty_match":1,"distractor_quality":1,'
    '"overall":1,"accept":true,"answer_correct":true,"reason":""}'
)


def judge_question(quiz_type, question, content_snip, target_difficulty, judge_model=None):
    """LLM-as-judge: score ONE generated question on a 1-5 rubric using an
    independent (larger) model. Fails closed (accept=False) on any parse error,
    since fail-open would bias a measurement harness."""
    qtxt = question.get("question") or question.get("statement") or ""
    ans  = str(
        question.get("correct_answer") or question.get("answer") or
        question.get("model_answer") or ""
    )[:200]
    diff_desc       = DIFF_DESCRIPTIONS.get(target_difficulty, "")
    options_line    = ""
    distractor_line = ""
    if quiz_type == "mcq":
        options_line    = f"Options: {json.dumps(question.get('options', {}), ensure_ascii=False)}\n"
        distractor_line = "4. distractor_quality — wrong options are plausible but clearly incorrect\n"

    prompt = (
        f"You are an impartial, strict but fair exam-quality judge. Score this ONE "
        f"{quiz_type.upper()} question on a 1-5 scale for each dimension.\n\n"
        f"Target difficulty: {target_difficulty.upper()} — {diff_desc}\n\n"
        f"Content snippet (ground truth):\n{content_snip[:4500]}\n\n"
        f"Question:\nQ: {qtxt[:400]}\n{options_line}A: {ans}\n\n"
        "Score 1 (worst) to 5 (best) on:\n"
        "1. correctness — factually correct and answerable from the content above\n"
        "2. clarity — unambiguous, well-formed wording\n"
        "3. difficulty_match — actual difficulty matches the target difficulty\n"
        f"{distractor_line}"
        "5. overall — holistic quality\n\n"
        "BE DISCRIMINATING — use the FULL range, do not default to 5:\n"
        "  5 = flawless, exam-ready as-is; 4 = good, minor nit; 3 = acceptable but "
        "ordinary; 2 = weak (generic wording, shallow recall, weak distractors); 1 = poor.\n"
        "Most decent questions are 3-4. Reserve 5 for genuinely excellent ones. "
        "Penalise generic phrasing, near-duplicate framing, obvious/implausible "
        "distractors, and shallow rote recall.\n\n"
        "Also set answer_correct: is the marked answer 'A:' actually the correct "
        "answer to this question given the content? (true/false)\n"
        "Then decide accept (true) or reject (false) for inclusion in a real exam.\n\n"
        f"Return STRICT JSON: {JUDGE_SCHEMA}"
    )

    fallback = {
        "correctness": 0, "clarity": 0, "difficulty_match": 0,
        "distractor_quality": None, "overall": 0,
        "accept": False, "answer_correct": False, "reason": "judge_parse_error",
    }
    result = call_ollama_json(
        prompt,
        system="You are an impartial, strict but fair exam-quality judge. Return only JSON.",
        fallback=fallback,
        model=judge_model or OLLAMA_MODEL, num_ctx=QUIZ_CTX,
    )
    if not isinstance(result, dict) or "overall" not in result:
        return dict(fallback)

    out = {
        "correctness":       int(result.get("correctness", 0) or 0),
        "clarity":            int(result.get("clarity", 0) or 0),
        "difficulty_match":   int(result.get("difficulty_match", 0) or 0),
        "distractor_quality": (int(result["distractor_quality"])
                                if quiz_type == "mcq" and result.get("distractor_quality") is not None
                                else None),
        "overall":            int(result.get("overall", 0) or 0),
        "accept":             bool(result.get("accept", False)),
        "answer_correct":     bool(result.get("answer_correct", False)),
        "reason":             str(result.get("reason", ""))[:300],
    }
    # Keep accept/overall internally consistent — a low overall score cannot be "accepted".
    if out["overall"] < 3:
        out["accept"] = False
    return out


def _build_eval_ctx(quiz_type, num_questions):
    """Build the same duck-typed ctx object quiz_generate_route uses (see
    colab.py's '/quiz/generate' route), so the agentic and non-agentic paths see
    byte-identical lecture content. Returns (ctx, content, summary), or
    (None, "", "") if no lecture has been processed yet."""
    summary = state.status.get("summary", "")
    if not summary:
        return None, "", ""

    ctx = type('Ctx', (), {
        'content_type':       'video' if state.video_path else 'document',
        'board_text':         state.status.get("board_text", ""),
        'transcript':         state.status.get("transcript", ""),
        'board_entries':      load_board_entries_from_disk(),
        'doc_text':           state.status.get("summary", ""),
        'summary':            summary,
        'agent_log':          [],
        '_quiz_type':         quiz_type,
        '_quiz_count':        num_questions,
        '_user_instructions': "",
        '_plan_mode':         'auto',
        '_manual_topics':     [],
        '_manual_difficulty': 'medium',
        '_mode':              'manual',
        '_quiz_plan':         None,
        '_agent_name':        'EvalAgent',
        'segment_list':       [],
        'file_path':          state.video_path or state.doc_path or "",
        'source_name':        Path(state.video_path or state.doc_path or "unknown").name,
        'thinking_cb':        None,
    })()

    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    if os.path.exists(seg_path):
        try:
            with open(seg_path, "r", encoding="utf-8") as f:
                ctx.segment_list = json.load(f)
        except Exception:
            ctx.segment_list = []

    if state.doc_path:
        doc_path = os.path.join(OUTPUT_DIR, "doc_text.txt")
        if os.path.exists(doc_path):
            with open(doc_path, "r", encoding="utf-8") as f:
                ctx.doc_text = f.read()

    content = (
        pack_video_sources(ctx.board_text, ctx.transcript, ctx.board_entries or [], 7500)
        if ctx.content_type == "video"
        else (ctx.doc_text or "")[:7500]
    )
    return ctx, content, summary


def _collect_agentic_results(quiz_type, ctx, content, summary, num_questions,
                              target_difficulty="medium", max_attempts_per_q=3,
                              on_event=None):
    """Runs the production agentic loop unmodified, observing it via thinking_cb.
    If on_event is given, each agentic event is also forwarded live (for the UI).
    Returns {questions, attempts_by_qid, validator_verdict_by_qid,
    durations_by_qid, approx_llm_calls, wall_clock_sec}."""
    events = []
    def cb(e):
        events.append(e)
        if on_event:
            on_event(e)
    # Honor the requested count. _run_quiz_agentic_loop reads ctx._quiz_count and
    # caches ctx._quiz_plan, and the same ctx is reused across difficulty groups
    # in run_bulk_evaluation — so force the count and clear any stale plan here,
    # or every group would generate just 1 question (or reuse the prior plan).
    ctx._quiz_count = num_questions
    ctx._quiz_plan  = None
    t0 = time.time()
    questions = _run_quiz_agentic_loop(
        ctx, quiz_type, content, summary, "",
        target_difficulty=target_difficulty,
        thinking_cb=cb,
        max_attempts_per_q=max_attempts_per_q,
    )
    wall_clock_sec = time.time() - t0

    attempts_by_qid          = {}
    durations_by_qid         = {}
    validator_verdict_by_qid = {}
    last_verdict_seen        = {}
    for e in events:
        qid = e.get("question_id")
        if not qid:
            continue
        if e.get("step") == "question_validated":
            last_verdict_seen[qid] = e.get("verdict", "FAIL")
        elif e.get("step") == "question_finalized":
            # A finalized event with status FAIL means the slot was skipped (no
            # question kept) — don't count it, or it would misalign the
            # per-question metrics/F1 against the actual `questions` list.
            if e.get("status") == "FAIL":
                continue
            attempts_by_qid[qid]  = e.get("total_attempts", 1)
            durations_by_qid[qid] = e.get("duration_sec", 0.0)
            # Use the real internal-validator verdict on the kept attempt, NOT
            # question_finalized.status — that field is hardcoded "PASS" even
            # when the loop keeps a question after exhausting all retries on a
            # FAIL (see _generate_one_question_with_retry's max-attempts path).
            if qid in last_verdict_seen:
                validator_verdict_by_qid[qid] = last_verdict_seen[qid]
            # else: question came from the deep-fallback repair path and never
            # went through _validate_one_question — no verdict to record.

    approx_llm_calls = 2 * sum(attempts_by_qid.values())  # ~1 generate + 1 validate per attempt
    return {
        "questions":                questions,
        "attempts_by_qid":          attempts_by_qid,
        "durations_by_qid":         durations_by_qid,
        "validator_verdict_by_qid": validator_verdict_by_qid,
        "approx_llm_calls":         approx_llm_calls,
        "wall_clock_sec":           wall_clock_sec,
    }


def _collect_nonagentic_results(quiz_type, ctx, num_questions, target_difficulty="medium"):
    """Calls the matching legacy single-shot bulk generator ONCE, unmodified.
    Returns {questions, wall_clock_sec, llm_call_count}.
    NOTE: generate_tf/generate_fill/generate_short hardcode their own question
    counts (10/10/8) and have no num_questions parameter at all — only
    generate_mcq honors a count. The returned list length may therefore differ
    from num_questions for those three; this is a reported methodology caveat,
    not something this harness patches around."""
    t0 = time.time()
    if quiz_type == "mcq":
        questions = generate_mcq(
            ctx.board_text, ctx.transcript, ctx.summary, num_questions=num_questions,
            doc_text=ctx.doc_text, board_entries=ctx.board_entries,
            segment_list=ctx.segment_list, difficulty_hint=target_difficulty,
        )
    elif quiz_type == "tf":
        questions = generate_tf(
            ctx.transcript, ctx.summary, doc_text=ctx.doc_text,
            segment_list=ctx.segment_list, difficulty_hint=target_difficulty,
            board_text=ctx.board_text, board_entries=ctx.board_entries,
        )
    elif quiz_type == "fill":
        questions = generate_fill(
            ctx.transcript, ctx.summary, doc_text=ctx.doc_text,
            board_entries=ctx.board_entries, difficulty_hint=target_difficulty,
            board_text=ctx.board_text,
        )
    else:
        questions = generate_short(
            ctx.transcript, ctx.summary, doc_text=ctx.doc_text,
            segment_list=ctx.segment_list, difficulty_hint=target_difficulty,
            board_text=ctx.board_text, board_entries=ctx.board_entries,
        )
    wall_clock_sec = time.time() - t0
    return {
        "questions":      questions or [],
        "wall_clock_sec": wall_clock_sec,
        "llm_call_count": 1,
    }


def _duplicate_rate(questions):
    """Pairwise near-duplicate rate using the same difflib approach as
    _validate_one_question (ratio > 0.85)."""
    texts = [
        (q.get("question") or q.get("statement") or "").strip().lower()
        for q in questions
    ]
    if not texts:
        return 0.0
    dup_count = 0
    for i, t in enumerate(texts):
        if not t:
            continue
        for prev in texts[:i]:
            if prev and difflib.SequenceMatcher(None, prev, t).ratio() > 0.85:
                dup_count += 1
                break
    return dup_count / len(texts)


def _qtext(q):
    return (q.get("question") or q.get("statement") or "")


def _structural_metrics(quiz_type, questions):
    """Deterministic, NO-LLM quality signals about a question set: length,
    topic spread, and per-type structural validity (valid MCQ keys, T/F balance,
    fill blanks present, short-answer key points)."""
    n = len(questions)
    lens   = [len(_qtext(q).split()) for q in questions if _qtext(q)]
    topics = [(q.get("topic") or "").strip() for q in questions if (q.get("topic") or "").strip()]
    m = {
        "avg_question_len_words": round(sum(lens) / len(lens), 1) if lens else None,
        "distinct_topics":        len(set(topics)),
    }
    if not n:
        return m
    if quiz_type == "mcq":
        valid = sum(1 for q in questions
                    if q.get("correct_answer") in (q.get("options", {}) or {}))
        optn  = [len(q.get("options", {}) or {}) for q in questions]
        m["valid_answer_key_rate"] = round(valid / n, 3)
        m["avg_options"]           = round(sum(optn) / len(optn), 1) if optn else None
    elif quiz_type == "tf":
        trues = sum(1 for q in questions
                    if str(q.get("answer")).strip().lower() in ("true", "1", "yes", "t"))
        m["true_fraction"] = round(trues / n, 3)   # ~0.5 == well balanced
    elif quiz_type == "fill":
        m["has_blank_rate"] = round(sum(1 for q in questions if "___" in _qtext(q)) / n, 3)
    elif quiz_type == "short":
        m["has_key_points_rate"] = round(sum(1 for q in questions if q.get("key_points")) / n, 3)
    return m


def _grounding_metrics(questions, content):
    """Embedding-based, NO-LLM-judge signals using the SentenceTransformer that
    already powers RAG:
      semantic_grounding     — mean cosine of each question to its best-matching
                               chunk of the lecture (higher = better grounded).
      inter_question_similarity — mean nearest-neighbour cosine between questions
                               (lower = more diverse / less redundant)."""
    try:
        texts = [_qtext(q) for q in questions if _qtext(q).strip()]
        if not texts or not (content or "").strip():
            return {"semantic_grounding": None, "inter_question_similarity": None}
        emb    = get_embedder()
        chunks = [content[i:i + 500] for i in range(0, min(len(content), 8000), 500)] or [content[:500]]
        qv = np.asarray(emb.encode(texts), dtype="float32")
        cv = np.asarray(emb.encode(chunks), dtype="float32")

        def _norm(mtx):
            nrm = np.linalg.norm(mtx, axis=1, keepdims=True)
            nrm[nrm == 0] = 1.0
            return mtx / nrm

        qn, cn = _norm(qv), _norm(cv)
        grounding = float((qn @ cn.T).max(axis=1).mean())
        if len(texts) > 1:
            qq = qn @ qn.T
            np.fill_diagonal(qq, 0.0)
            inter = float(qq.max(axis=1).mean())
        else:
            inter = 0.0
        return {"semantic_grounding": round(grounding, 3),
                "inter_question_similarity": round(inter, 3)}
    except Exception:
        return {"semantic_grounding": None, "inter_question_similarity": None}


def _compute_pipeline_metrics(questions, judged_scores, attempts_by_qid=None,
                               llm_call_count=1, wall_clock_sec=0.0, durations_by_qid=None,
                               quiz_type=None, content=""):
    """judged_scores: list of judge_question() output dicts, same order/length as
    questions. Returns a flat metrics dict for one pipeline on one quiz_type,
    combining (a) LLM-judge scores, (b) deterministic structural metrics, and
    (c) embedding-based grounding/diversity metrics — so the comparison isn't
    judge-only."""
    n = len(judged_scores)

    def _mean(key, only_not_none=False):
        vals = [s.get(key) for s in judged_scores]
        if only_not_none:
            vals = [v for v in vals if v is not None]
        return round(sum(vals) / len(vals), 3) if vals else None

    accept_count   = sum(1 for s in judged_scores if s.get("accept"))
    attempts_vals  = list((attempts_by_qid or {}).values())
    duration_vals  = list((durations_by_qid or {}).values())

    metrics = {
        "n_questions":                          n,
        # ── LLM-judge metrics ──
        "mean_overall":                         _mean("overall"),
        "mean_correctness":                     _mean("correctness"),
        "mean_clarity":                         _mean("clarity"),
        "mean_difficulty_match":                _mean("difficulty_match"),
        "mean_distractor_quality":              _mean("distractor_quality", only_not_none=True),
        "judge_accept_rate":                    round(accept_count / n, 3) if n else 0.0,
        # ── diversity / redundancy ──
        "duplicate_rate":                       round(_duplicate_rate(questions), 3),
        # ── cost / latency ──
        "avg_attempts":                         round(sum(attempts_vals) / len(attempts_vals), 2) if attempts_vals else 1.0,
        "avg_question_duration_sec":            round(sum(duration_vals) / len(duration_vals), 2) if duration_vals else None,
        "llm_call_count":                       llm_call_count,
        "avg_llm_calls_per_accepted_question":  round(llm_call_count / max(accept_count, 1), 2),
        "wall_clock_sec":                       round(wall_clock_sec, 2),
    }
    # ── deterministic structural + embedding-based metrics (no judge) ──
    metrics.update(_structural_metrics(quiz_type or "", questions))
    metrics.update(_grounding_metrics(questions, content))
    return metrics


def _compute_validator_judge_f1(validator_verdicts, judge_accepts):
    """validator_verdicts: list of 'PASS'/'FAIL' (internal validator, agentic
    only). judge_accepts: list of bool (judge accept/reject), same order/length,
    already filtered to only the questions that actually went through the
    validator. Treats judge accept as the ground-truth label, validator PASS as
    the predicted label. Returns {precision, recall, f1, accuracy, tp, fp, fn, tn, n}."""
    pairs = list(zip(validator_verdicts, judge_accepts))
    tp = sum(1 for v, j in pairs if v == "PASS" and j)
    fp = sum(1 for v, j in pairs if v == "PASS" and not j)
    fn = sum(1 for v, j in pairs if v == "FAIL" and j)
    tn = sum(1 for v, j in pairs if v == "FAIL" and not j)

    precision = tp / (tp + fp) if (tp + fp) else 0.0
    recall    = tp / (tp + fn) if (tp + fn) else 0.0
    f1        = (2 * precision * recall / (precision + recall)) if (precision + recall) else 0.0
    accuracy  = (tp + tn) / len(pairs) if pairs else 0.0

    return {
        "precision": round(precision, 3), "recall": round(recall, 3),
        "f1": round(f1, 3), "accuracy": round(accuracy, 3),
        "tp": tp, "fp": fp, "fn": fn, "tn": tn, "n": len(pairs),
    }


def run_quiz_evaluation(quiz_types=("mcq", "tf", "fill", "short"),
                         num_questions=6, repeats=2, judge_model=None,
                         target_difficulty="medium", save_outputs=True,
                         progress_cb=None):
    """Top-level evaluation harness — call this from a Colab cell after
    processing a real lecture. Runs both the agentic pipeline and the
    non-agentic baseline on identical content for each quiz_type (repeated
    `repeats` times), judges generated questions with an independent LLM judge,
    aggregates comparison metrics (incl. validator-vs-judge F1), renders charts,
    and (if save_outputs) writes eval_results.json + PNGs to OUTPUT_DIR.
    If progress_cb is given it is called with live event dicts for the UI.
    Always returns the full results dict."""
    def emit(kind, **kw):
        if progress_cb:
            try:
                progress_cb({"kind": kind, **kw})
            except Exception:
                pass

    ctx0, _, _ = _build_eval_ctx(quiz_types[0], 1)
    if ctx0 is None:
        emit("error", reason="No content processed yet — upload and process a lecture first.")
        return {"ok": False, "error": "No content processed yet — upload and process a lecture first."}

    total_units = len(quiz_types) * max(1, repeats)
    unit        = 0
    # Cap how many questions per pipeline get judged, so the slow 12B judge
    # doesn't have to score the non-agentic generators' 10/10/8 hardcoded
    # outputs — keeps both pipelines comparable and the run bounded.
    judge_cap = max(1, num_questions)
    emit("start", total_units=total_units, quiz_types=list(quiz_types),
         num_questions=num_questions, repeats=repeats, judge_cap=judge_cap)

    runs_by_type = {qt: [] for qt in quiz_types}

    for qt in quiz_types:
        for r in range(max(1, repeats)):
            unit += 1
            ctx, content, summary = _build_eval_ctx(qt, num_questions)
            emit("unit_start", quiz_type=qt, repeat=r + 1, unit=unit, total_units=total_units)

            # ── Agentic pipeline (live sub-events forwarded) ──
            emit("agentic_start", quiz_type=qt, unit=unit)
            def _fwd(e, _qt=qt, _u=unit):
                emit("agentic_event", quiz_type=_qt, unit=_u,
                     step=e.get("step"), agent=e.get("agent"),
                     verdict=e.get("verdict"), reason=e.get("reason", ""))
            agentic = _collect_agentic_results(qt, ctx, content, summary,
                                               num_questions, target_difficulty,
                                               on_event=_fwd)
            emit("agentic_done", quiz_type=qt, unit=unit,
                 count=len(agentic["questions"]),
                 secs=round(agentic["wall_clock_sec"], 1))

            # show a couple of the agentic questions in the live feed
            for q in agentic["questions"][:3]:
                emit("sample", quiz_type=qt, unit=unit, pipeline="agentic",
                     text=(_qtext(q) or "")[:140])

            # ── Non-agentic baseline (single shot) ──
            emit("nonagentic_start", quiz_type=qt, unit=unit)
            nonagentic = _collect_nonagentic_results(qt, ctx, num_questions, target_difficulty)
            emit("nonagentic_done", quiz_type=qt, unit=unit,
                 count=len(nonagentic["questions"]),
                 secs=round(nonagentic["wall_clock_sec"], 1))
            # show a couple of the non-agentic questions so the student can see
            # exactly what the baseline produced
            for q in nonagentic["questions"][:3]:
                emit("sample", quiz_type=qt, unit=unit, pipeline="nonagentic",
                     text=(_qtext(q) or "")[:140])

            # ── Judge both (capped) ──
            content_snip  = content[:6000]
            ag_qs = agentic["questions"][:judge_cap]
            na_qs = nonagentic["questions"][:judge_cap]
            emit("judge_start", quiz_type=qt, unit=unit,
                 n_agentic=len(ag_qs), n_nonagentic=len(na_qs))
            agentic_judged = []
            for i, q in enumerate(ag_qs):
                agentic_judged.append(judge_question(qt, q, content_snip, target_difficulty, judge_model))
                emit("judge_progress", quiz_type=qt, unit=unit, pipeline="agentic",
                     done=i + 1, total=len(ag_qs))
            nonagentic_judged = []
            for i, q in enumerate(na_qs):
                nonagentic_judged.append(judge_question(qt, q, content_snip, target_difficulty, judge_model))
                emit("judge_progress", quiz_type=qt, unit=unit, pipeline="nonagentic",
                     done=i + 1, total=len(na_qs))

            agentic_metrics = _compute_pipeline_metrics(
                ag_qs, agentic_judged, agentic["attempts_by_qid"],
                agentic["approx_llm_calls"], agentic["wall_clock_sec"],
                durations_by_qid=agentic["durations_by_qid"],
                quiz_type=qt, content=content_snip,
            )
            nonagentic_metrics = _compute_pipeline_metrics(
                na_qs, nonagentic_judged, None,
                nonagentic["llm_call_count"], nonagentic["wall_clock_sec"],
                quiz_type=qt, content=content_snip,
            )

            # F1 only over questions that actually went through the internal
            # validator (excludes the rare deep-fallback recovery path).
            validator_verdicts, judge_accepts_for_f1 = [], []
            for idx, qid in enumerate(list(agentic["attempts_by_qid"].keys())[:judge_cap]):
                if qid in agentic["validator_verdict_by_qid"] and idx < len(agentic_judged):
                    validator_verdicts.append(agentic["validator_verdict_by_qid"][qid])
                    judge_accepts_for_f1.append(agentic_judged[idx]["accept"])
            f1_metrics = _compute_validator_judge_f1(validator_verdicts, judge_accepts_for_f1)

            emit("unit_done", quiz_type=qt, unit=unit, total_units=total_units,
                 agentic_score=agentic_metrics.get("mean_overall"),
                 nonagentic_score=nonagentic_metrics.get("mean_overall"),
                 agentic_accept=agentic_metrics.get("judge_accept_rate"),
                 nonagentic_accept=nonagentic_metrics.get("judge_accept_rate"),
                 f1=f1_metrics.get("f1"))

            def _q_with_score(qs, judged):
                out = []
                for i, q in enumerate(qs):
                    j = judged[i] if i < len(judged) else {}
                    out.append({
                        "question": (_qtext(q) or "")[:300],
                        "answer":   str(q.get("correct_answer") or q.get("answer")
                                        or q.get("model_answer") or "")[:160],
                        "options":  q.get("options") if qt == "mcq" else None,
                        "topic":    q.get("topic", ""),
                        "difficulty": q.get("difficulty", ""),
                        "judge_overall": j.get("overall"),
                        "judge_accept":  j.get("accept"),
                        "judge_reason":  j.get("reason", ""),
                    })
                return out

            runs_by_type[qt].append({
                "agentic_metrics":    agentic_metrics,
                "nonagentic_metrics": nonagentic_metrics,
                "f1_metrics":         f1_metrics,
                # the ACTUAL questions each pipeline produced (with judge scores),
                # so the report/UI can show what non-agentic vs agentic generated
                "agentic_questions":    _q_with_score(ag_qs, agentic_judged),
                "nonagentic_questions": _q_with_score(na_qs, nonagentic_judged),
            })

    def _aggregate(metric_dicts, numeric_keys):
        agg = {}
        for k in numeric_keys:
            vals = [m[k] for m in metric_dicts if m.get(k) is not None]
            agg[k] = round(sum(vals) / len(vals), 3) if vals else None
        return agg

    metric_keys = [
        "mean_overall", "mean_correctness", "mean_clarity", "mean_difficulty_match",
        "mean_distractor_quality", "judge_accept_rate", "duplicate_rate",
        "avg_attempts", "avg_question_duration_sec",
        "avg_llm_calls_per_accepted_question", "wall_clock_sec",
        # deterministic + embedding metrics (no judge)
        "avg_question_len_words", "distinct_topics", "semantic_grounding",
        "inter_question_similarity", "valid_answer_key_rate", "avg_options",
        "true_fraction", "has_blank_rate", "has_key_points_rate",
    ]
    f1_keys = ["precision", "recall", "f1", "accuracy"]

    aggregated = {}
    for qt in quiz_types:
        runs = runs_by_type[qt]
        aggregated[qt] = {
            "agentic":    _aggregate([r["agentic_metrics"] for r in runs], metric_keys),
            "nonagentic": _aggregate([r["nonagentic_metrics"] for r in runs], metric_keys),
            "f1_metrics": _aggregate([r["f1_metrics"] for r in runs], f1_keys),
        }

    all_runs = [r for qt in quiz_types for r in runs_by_type[qt]]
    aggregated["overall"] = {
        "agentic":    _aggregate([r["agentic_metrics"] for r in all_runs], metric_keys),
        "nonagentic": _aggregate([r["nonagentic_metrics"] for r in all_runs], metric_keys),
        "f1_metrics": _aggregate([r["f1_metrics"] for r in all_runs], f1_keys),
    }

    results = {
        "ok": True,
        "generated_at": datetime.now().isoformat(),
        "config": {
            "quiz_types": list(quiz_types), "num_questions": num_questions,
            "repeats": repeats, "target_difficulty": target_difficulty,
            "judge_model": judge_model or OLLAMA_MODEL, "generator_model": QUIZ_MODEL,
        },
        "runs_by_quiz_type": runs_by_type,
        "aggregated": aggregated,
    }

    if save_outputs:
        save_json("eval_results.json", results)
        _render_eval_charts(results, OUTPUT_DIR)

    ov = aggregated.get("overall", {})
    emit("done",
         agentic_score=ov.get("agentic", {}).get("mean_overall"),
         nonagentic_score=ov.get("nonagentic", {}).get("mean_overall"),
         f1=ov.get("f1_metrics", {}).get("f1"))
    return results


def _generate_nonagentic_on_topic(quiz_type, topic, content, summary, model=None):
    """The NON-AGENTIC way: one single-shot LLM call to write ONE question on a
    given topic — no planning, no validation, no retry. Used to produce a
    head-to-head counterpart for an already-generated agentic question on the
    SAME topic."""
    schema = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
    prompt = (
        f"Generate exactly ONE {quiz_type.upper()} question on the topic '{topic}'.\n"
        "Single attempt, no self-checking.\n"
        f"Output STRICT JSON (a 'questions' array with exactly 1 object):\n{schema}\n\n"
        f"CONTENT:\n{content[:5000]}\n\nSUMMARY:\n{summary[:1000]}"
    )
    data  = call_ollama_json(prompt, fallback={"questions": []},
                             model=model or QUIZ_MODEL, num_ctx=QUIZ_CTX)
    items = _extract_items(data, "questions")
    if not items:
        # fast model returned nothing usable — fall back to the primary model so
        # the non-agentic side still produces a question to compare against
        data  = call_ollama_json(prompt, fallback={"questions": []})
        items = _extract_items(data, "questions")
    if items:
        q = items[0]
        if not q.get("topic"):
            q["topic"] = topic
        return q
    return None


def run_quiz_comparison(quiz_type, judge_model=None, target_difficulty="medium",
                        progress_cb=None):
    """Evaluate the ALREADY-GENERATED questions of one quiz type against a
    non-agentic baseline built on the SAME topics, then judge both with an
    independent LLM and report a 3-segment rating (overall / in-context /
    difficulty), F1, and other metrics. Anchors on existing questions — it does
    NOT regenerate the agentic side."""
    def emit(kind, **kw):
        if progress_cb:
            try:
                progress_cb({"kind": kind, **kw})
            except Exception:
                pass

    agentic_qs = state.quiz.get(quiz_type, [])
    if not agentic_qs:
        emit("error", reason=f"No {quiz_type.upper()} questions generated yet.")
        return {"ok": False, "error": f"No {quiz_type.upper()} questions generated yet — generate first."}

    ctx, content, summary = _build_eval_ctx(quiz_type, len(agentic_qs))
    content_snip     = (content or summary or state.status.get("summary", ""))[:6000]
    nonagentic_model = QUIZ_MODEL
    judge            = judge_model or OLLAMA_MODEL

    emit("compare_start", quiz_type=quiz_type, n=len(agentic_qs),
         nonagentic_model=nonagentic_model, judge_model=judge,
         agentic_model=QUIZ_MODEL)

    # 1) Non-agentic generates one question per existing topic
    emit("nonagentic_model", model=nonagentic_model)
    pairs = []
    for i, aq in enumerate(agentic_qs):
        topic = (aq.get("topic") or "General")
        emit("nonagentic_gen", topic=topic, idx=i + 1, total=len(agentic_qs))
        nq = _generate_nonagentic_on_topic(quiz_type, topic, content_snip, summary, nonagentic_model)
        emit("nonagentic_q", topic=topic, text=((_qtext(nq) if nq else "(generation failed)") or "")[:140])
        pairs.append((aq, nq, topic))

    # 2) Judge both sides + re-validate the agentic side for F1
    emit("judge_model", model=judge)
    rows = []
    ag_overall, ag_ctx, ag_diff, ag_accept = [], [], [], 0
    na_overall, na_ctx, na_diff, na_accept = [], [], [], 0
    val_verdicts, judge_accepts = [], []
    for i, (aq, nq, topic) in enumerate(pairs):
        diff = aq.get("difficulty") or target_difficulty
        ja   = judge_question(quiz_type, aq, content_snip, diff, judge)
        passed, _, _ = _validate_one_question(quiz_type, aq, content_snip, diff, "")
        val_verdicts.append("PASS" if passed else "FAIL")
        judge_accepts.append(bool(ja.get("accept")))
        jn = (judge_question(quiz_type, nq, content_snip, diff, judge) if nq else
              {"overall": 0, "correctness": 0, "difficulty_match": 0, "accept": False, "reason": "no question produced"})

        ag_overall.append(ja.get("overall") or 0); ag_ctx.append(ja.get("correctness") or 0)
        ag_diff.append(ja.get("difficulty_match") or 0); ag_accept += 1 if ja.get("accept") else 0
        na_overall.append(jn.get("overall") or 0); na_ctx.append(jn.get("correctness") or 0)
        na_diff.append(jn.get("difficulty_match") or 0); na_accept += 1 if jn.get("accept") else 0

        rows.append({
            "topic": topic,
            "agentic": {
                "question": (_qtext(aq) or "")[:240], "overall": ja.get("overall"),
                "in_context": ja.get("correctness"), "difficulty": ja.get("difficulty_match"),
                "accept": bool(ja.get("accept")), "reason": ja.get("reason", ""),
            },
            "nonagentic": {
                "question": (_qtext(nq) if nq else "") or "", "overall": jn.get("overall"),
                "in_context": jn.get("correctness"), "difficulty": jn.get("difficulty_match"),
                "accept": bool(jn.get("accept")), "reason": jn.get("reason", ""),
            },
        })
        emit("pair_done", idx=i + 1, total=len(pairs), topic=topic,
             agentic_overall=ja.get("overall"), nonagentic_overall=jn.get("overall"))

    def _avg(x):
        return round(sum(x) / len(x), 2) if x else None

    n = len(pairs)
    f1 = _compute_validator_judge_f1(val_verdicts, judge_accepts)
    ag_list = [p[0] for p in pairs]
    na_list = [p[1] for p in pairs if p[1]]
    agentic_metrics = {
        "accept_rate":    round(ag_accept / n, 2) if n else 0.0,
        "duplicate_rate": round(_duplicate_rate(ag_list), 3),
        **_grounding_metrics(ag_list, content_snip),
    }
    nonagentic_metrics = {
        "accept_rate":    round(na_accept / n, 2) if n else 0.0,
        "duplicate_rate": round(_duplicate_rate(na_list), 3),
        **_grounding_metrics(na_list, content_snip),
    }

    # Deterministic tiebreaker: when the judge averages are equal, decide the
    # segment on real metrics — reward grounding, penalise repetition — so the
    # comparison gives an actual verdict instead of collapsing to "tie".
    ag_comp = (agentic_metrics.get("semantic_grounding") or 0)    - (agentic_metrics.get("duplicate_rate") or 0)
    na_comp = (nonagentic_metrics.get("semantic_grounding") or 0) - (nonagentic_metrics.get("duplicate_rate") or 0)

    segments = {
        "overall":    {"label": "Overall quality",    "agentic": _avg(ag_overall), "nonagentic": _avg(na_overall)},
        "in_context": {"label": "In-context / grounded", "agentic": _avg(ag_ctx),  "nonagentic": _avg(na_ctx)},
        "difficulty": {"label": "Difficulty match",    "agentic": _avg(ag_diff),    "nonagentic": _avg(na_diff)},
    }
    for seg in segments.values():
        a, b = seg["agentic"] or 0, seg["nonagentic"] or 0
        if a > b:
            seg["winner"] = "agentic"
        elif b > a:
            seg["winner"] = "nonagentic"
        elif ag_comp > na_comp:
            seg["winner"], seg["tiebreak"] = "agentic", "grounding/diversity"
        elif na_comp > ag_comp:
            seg["winner"], seg["tiebreak"] = "nonagentic", "grounding/diversity"
        else:
            seg["winner"] = "tie"

    result = {
        "ok": True, "quiz_type": quiz_type, "n": n,
        "agentic_model": QUIZ_MODEL, "nonagentic_model": nonagentic_model, "judge_model": judge,
        "segments": segments, "f1": f1,
        "agentic_metrics": agentic_metrics, "nonagentic_metrics": nonagentic_metrics,
        "rows": rows, "generated_at": datetime.now().isoformat(),
    }
    save_json("eval_compare.json", result)
    emit("compare_done",
         overall_agentic=segments["overall"]["agentic"],
         overall_nonagentic=segments["overall"]["nonagentic"], f1=f1.get("f1"))
    return result


# ══════════════════════════════════════════════════════════════
# BULK EVALUATION  —  N MCQ x difficulties, agentic vs non-agentic,
# with answer-correctness (independent re-solve + judge-verify) and
# gold-answer calibration. Resumable background run.
# ══════════════════════════════════════════════════════════════

BULK_ROWS_PATH      = os.path.join(OUTPUT_DIR, "bulk_eval_rows.jsonl")
BULK_QUESTIONS_PATH = os.path.join(OUTPUT_DIR, "bulk_eval_questions.json")
BULK_RESULTS_PATH   = os.path.join(OUTPUT_DIR, "bulk_eval_results.json")
BULK_CSV_PATH       = os.path.join(OUTPUT_DIR, "bulk_eval.csv")


def solve_mcq_independently(question, content_snip, judge_model=None):
    """Independent re-solve: answer an MCQ from the content WITHOUT seeing the
    marked answer. Returns the chosen option letter (upper) or ''."""
    opts = question.get("options", {}) or {}
    if not opts:
        return ""
    opt_lines = "\n".join(f"{k}. {v}" for k, v in opts.items())
    prompt = (
        "Answer this multiple-choice question using ONLY the content below. "
        "Pick the single best option.\n\n"
        f"CONTENT:\n{(content_snip or '')[:4500]}\n\n"
        f"Question: {(question.get('question') or '')[:400]}\n{opt_lines}\n\n"
        'Return STRICT JSON: {"answer":"A"}'
    )
    data = call_ollama_json(prompt, fallback={}, model=judge_model or OLLAMA_MODEL, num_ctx=QUIZ_CTX)
    if isinstance(data, dict):
        return str(data.get("answer", "")).strip().upper()[:1]
    return ""


def load_gold_mcqs(path):
    """Load gold MCQs from JSON (list or {questions:[...]}) or CSV with columns
    question,A,B,C,D,answer[,difficulty,passage]. Returns a normalised list."""
    if not path or not os.path.exists(path):
        return []
    gold = []
    if path.lower().endswith(".csv"):
        import csv
        with open(path, "r", encoding="utf-8") as f:
            for r in csv.DictReader(f):
                opts = {k: r[k] for k in ("A", "B", "C", "D") if r.get(k)}
                gold.append({"question": r.get("question", ""), "options": opts,
                             "answer": str(r.get("answer", "")).strip().upper()[:1],
                             "difficulty": r.get("difficulty", ""),
                             "passage": r.get("passage", "")})
        return gold
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    items = data.get("questions", data) if isinstance(data, dict) else data
    for q in (items or []):
        gold.append({"question": q.get("question", ""),
                     "options": q.get("options", {}) or {},
                     "answer": str(q.get("answer") or q.get("correct_answer") or "").strip().upper()[:1],
                     "difficulty": q.get("difficulty", ""),
                     "passage": q.get("passage", "")})
    return gold


def evaluate_solver_on_gold(gold, models=None, progress_cb=None):
    """Cross-check answer correctness against KNOWN gold answers. Each model in
    `models` independently answers every gold MCQ and we report its accuracy vs
    the gold answer (overall + per difficulty). Default models = the agentic
    generator model (QUIZ_MODEL, 4B) and the non-agentic/primary model
    (OLLAMA_MODEL, 12B), so BOTH are cross-checked against ground truth.
    This validates the answer-checkers and gives a true answer-accuracy number."""
    if models is None:
        models = []
        for m in (QUIZ_MODEL, OLLAMA_MODEL):
            if m and m not in models:
                models.append(m)
    n = len(gold)
    per_model = {}
    for m in models:
        correct, per_diff, rows = 0, {}, []
        for i, g in enumerate(gold):
            choice = solve_mcq_independently(
                {"question": g["question"], "options": g["options"]},
                g.get("passage", ""), m,
            )
            ok = bool(choice) and choice == g.get("answer", "")
            correct += 1 if ok else 0
            d = g.get("difficulty") or "all"
            per_diff.setdefault(d, [0, 0]); per_diff[d][0] += 1 if ok else 0; per_diff[d][1] += 1
            rows.append({"question": (g["question"] or "")[:200], "gold": g.get("answer"),
                         "choice": choice, "correct": ok, "difficulty": d})
            if progress_cb:
                progress_cb({"kind": "gold_progress", "model": m, "done": i + 1, "total": n})
        per_model[m] = {
            "accuracy": round(correct / n, 3) if n else None,
            "by_difficulty": {d: round(c / t, 3) for d, (c, t) in per_diff.items() if t},
            "rows": rows,
        }
    result = {
        "ok": True, "n": n,
        "agentic_model": QUIZ_MODEL, "nonagentic_model": OLLAMA_MODEL,
        "per_model": per_model, "generated_at": datetime.now().isoformat(),
    }
    save_json("gold_calibration.json", result)
    return result


def _bulk_split(total, difficulties):
    base = total // len(difficulties)
    rem  = total - base * len(difficulties)
    counts = {d: base for d in difficulties}
    counts[difficulties[0]] += rem
    return counts


def _bulk_load_done():
    """Return (rows, done_set) from the resumable jsonl, keyed by pipeline/diff/idx."""
    rows, done = [], set()
    if os.path.exists(BULK_ROWS_PATH):
        with open(BULK_ROWS_PATH, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    r = json.loads(line)
                except Exception:
                    continue
                rows.append(r)
                done.add((r["pipeline"], r["difficulty"], r["idx"]))
    return rows, done


def run_bulk_evaluation(total=200, difficulties=("easy", "medium", "hard"),
                        compare=True, judge_model=None, fresh=False,
                        gold_path="", progress_cb=None, max_attempts_per_q=3):
    """Generate ~`total` MCQ split across `difficulties` with the agentic pipeline
    (and the non-agentic baseline if `compare`), judge each (quality +
    answer_correct) and independently re-solve it, and report per-difficulty
    comparison metrics including answer correctness. Resumable: generated
    questions and judged rows are written to disk and reused on re-invocation
    unless `fresh=True`. Returns the aggregated results dict."""
    def emit(kind, **kw):
        if progress_cb:
            try:
                progress_cb({"kind": kind, **kw})
            except Exception:
                pass

    ctx, content, summary = _build_eval_ctx("mcq", 1)
    if ctx is None:
        emit("error", reason="No content processed yet — upload and process a lecture first.")
        return {"ok": False, "error": "No content processed yet — upload and process a lecture first."}
    content_snip = content[:6000]

    if fresh:
        for p in (BULK_ROWS_PATH, BULK_QUESTIONS_PATH):
            if os.path.exists(p):
                os.remove(p)

    counts    = _bulk_split(total, difficulties)
    pipelines = ["agentic"] + (["nonagentic"] if compare else [])
    emit("bulk_start", total=total, counts=counts, pipelines=pipelines,
         judge_model=judge_model or OLLAMA_MODEL, generator_model=QUIZ_MODEL)

    # ── Phase 1: generation (saved per group; reused on resume) ──
    questions = {}
    verdicts  = {}            # group key -> [validator PASS/FAIL per question] (agentic)
    if not fresh and os.path.exists(BULK_QUESTIONS_PATH):
        try:
            with open(BULK_QUESTIONS_PATH, "r", encoding="utf-8") as f:
                saved = json.load(f)
            questions = saved.get("questions", {})
            verdicts  = saved.get("verdicts", {})
        except Exception:
            questions, verdicts = {}, {}

    def _save_questions():
        save_json("bulk_eval_questions.json", {"questions": questions, "verdicts": verdicts})

    for d in difficulties:
        k = counts[d]
        # agentic group
        key = f"agentic:{d}"
        if key not in questions:
            emit("group_gen", pipeline="agentic", difficulty=d, count=k)
            ag = _collect_agentic_results("mcq", ctx, content, summary, k,
                                          target_difficulty=d, max_attempts_per_q=max_attempts_per_q)
            qs = ag["questions"][:k]
            questions[key] = qs
            vv = ag.get("validator_verdict_by_qid", {})
            verdicts[key] = [vv.get(f"q{i+1}", "") for i in range(len(qs))]
            _save_questions()
            emit("group_done", pipeline="agentic", difficulty=d, count=len(qs))
        # non-agentic group
        if compare:
            key = f"nonagentic:{d}"
            if key not in questions:
                emit("group_gen", pipeline="nonagentic", difficulty=d, count=k)
                collected = []
                while len(collected) < k:
                    need  = min(10, k - len(collected))
                    batch = generate_mcq(ctx.board_text, ctx.transcript, summary,
                                         num_questions=need, doc_text=ctx.doc_text,
                                         board_entries=ctx.board_entries,
                                         segment_list=ctx.segment_list, difficulty_hint=d)
                    if not batch:
                        break
                    collected.extend(batch)
                questions[key] = collected[:k]
                _save_questions()
                emit("group_done", pipeline="nonagentic", difficulty=d, count=len(questions[key]))

    # ── Phase 2: judge + independent solve (per-question resumable) ──
    rows, done = ([], set()) if fresh else _bulk_load_done()
    total_to_judge = sum(len(questions.get(f"{p}:{d}", [])) for p in pipelines for d in difficulties)
    judged = len(rows)
    emit("judge_start", total=total_to_judge, already=judged)

    rf = open(BULK_ROWS_PATH, "a", encoding="utf-8")
    try:
        for p in pipelines:
            for d in difficulties:
                key = f"{p}:{d}"
                for idx, q in enumerate(questions.get(key, [])):
                    if (p, d, idx) in done:
                        continue
                    jr     = judge_question("mcq", q, content_snip, d, judge_model)
                    choice = solve_mcq_independently(q, content_snip, judge_model)
                    marked = str(q.get("correct_answer", "")).strip().upper()[:1]
                    row = {
                        "pipeline": p, "difficulty": d, "idx": idx,
                        "question": (q.get("question") or "")[:300],
                        "options": q.get("options", {}),
                        "topic": q.get("topic", ""),
                        "marked_answer": marked,
                        "judge_overall": jr.get("overall"),
                        "judge_correctness": jr.get("correctness"),
                        "judge_clarity": jr.get("clarity"),
                        "judge_difficulty_match": jr.get("difficulty_match"),
                        "judge_accept": bool(jr.get("accept")),
                        "answer_correct": bool(jr.get("answer_correct")),      # judge-verify
                        "independent_choice": choice,
                        "independent_match": bool(choice) and choice == marked, # re-solve
                        "validator_verdict": (verdicts.get(key, [None] * (idx + 1))[idx]
                                              if p == "agentic" else None),
                    }
                    rf.write(json.dumps(row, ensure_ascii=False) + "\n"); rf.flush()
                    rows.append(row); judged += 1
                    emit("judge_progress", done=judged, total=total_to_judge,
                         pipeline=p, difficulty=d)
    finally:
        rf.close()

    # ── Aggregate ──
    def _agg(pipeline, difficulty):
        sub = [r for r in rows if r["pipeline"] == pipeline and r["difficulty"] == difficulty]
        n = len(sub)
        if not n:
            return None
        def mean(key):
            vals = [r[key] for r in sub if isinstance(r.get(key), (int, float))]
            return round(sum(vals) / len(vals), 3) if vals else None
        qlist = [{"question": r["question"], "options": r["options"]} for r in sub]
        m = {
            "n": n,
            "mean_overall": mean("judge_overall"),
            "mean_correctness": mean("judge_correctness"),
            "mean_clarity": mean("judge_clarity"),
            "mean_difficulty_match": mean("judge_difficulty_match"),
            "accept_rate": round(sum(1 for r in sub if r["judge_accept"]) / n, 3),
            "verified_correct_rate": round(sum(1 for r in sub if r["answer_correct"]) / n, 3),
            "independent_match_rate": round(sum(1 for r in sub if r["independent_match"]) / n, 3),
            "duplicate_rate": round(_duplicate_rate(qlist), 3),
        }
        m.update(_grounding_metrics(qlist, content_snip))
        # effective accept (duplicate-penalised) and answer accuracy
        m["effective_accept_rate"] = round(m["accept_rate"] * (1 - m["duplicate_rate"]), 3)
        if pipeline == "agentic":
            vv = [r["validator_verdict"] for r in sub]
            ja = [r["judge_accept"] for r in sub]
            pairs = [(v, j) for v, j in zip(vv, ja) if v in ("PASS", "FAIL")]
            if pairs:
                m["validator_judge_f1"] = _compute_validator_judge_f1(
                    [v for v, _ in pairs], [j for _, j in pairs])
        return m

    aggregated = {}
    for d in difficulties:
        aggregated[d] = {p: _agg(p, d) for p in pipelines}
    aggregated["overall"] = {
        p: _agg_overall(rows, p, content_snip) for p in pipelines
    }

    gold_result = None
    if gold_path:
        gold = load_gold_mcqs(gold_path)
        if gold:
            emit("gold_start", n=len(gold))
            gold_result = evaluate_solver_on_gold(gold, progress_cb=progress_cb)

    results = {
        "ok": True, "generated_at": datetime.now().isoformat(),
        "config": {"total": total, "difficulties": list(difficulties), "compare": compare,
                   "quiz_type": "mcq", "judge_model": judge_model or OLLAMA_MODEL,
                   "generator_model": QUIZ_MODEL},
        "aggregated": aggregated,
        "gold_calibration": gold_result,
        "n_rows": len(rows),
    }
    save_json("bulk_eval_results.json", results)
    _write_bulk_csv(rows)
    try:
        _render_bulk_charts(results, OUTPUT_DIR)
    except Exception as e:
        print(f"  [bulk] chart render skipped: {e}")

    emit("bulk_done", n_rows=len(rows))
    return results


def _agg_overall(rows, pipeline, content_snip):
    sub = [r for r in rows if r["pipeline"] == pipeline]
    n = len(sub)
    if not n:
        return None
    qlist = [{"question": r["question"], "options": r["options"]} for r in sub]
    def mean(key):
        vals = [r[key] for r in sub if isinstance(r.get(key), (int, float))]
        return round(sum(vals) / len(vals), 3) if vals else None
    dup = round(_duplicate_rate(qlist), 3)
    acc = round(sum(1 for r in sub if r["judge_accept"]) / n, 3)
    out = {
        "n": n, "mean_overall": mean("judge_overall"),
        "accept_rate": acc,
        "verified_correct_rate": round(sum(1 for r in sub if r["answer_correct"]) / n, 3),
        "independent_match_rate": round(sum(1 for r in sub if r["independent_match"]) / n, 3),
        "duplicate_rate": dup,
        "effective_accept_rate": round(acc * (1 - dup), 3),
    }
    out.update(_grounding_metrics(qlist, content_snip))
    return out


def _write_bulk_csv(rows):
    import csv
    cols = ["pipeline", "difficulty", "idx", "topic", "question", "marked_answer",
            "judge_overall", "judge_correctness", "judge_clarity",
            "judge_difficulty_match", "judge_accept", "answer_correct",
            "independent_choice", "independent_match", "validator_verdict"]
    with open(BULK_CSV_PATH, "w", encoding="utf-8", newline="") as f:
        w = csv.DictWriter(f, fieldnames=cols, extrasaction="ignore")
        w.writeheader()
        for r in rows:
            w.writerow(r)


def _render_bulk_charts(results, output_dir):
    """Bar charts by difficulty (agentic vs non-agentic) for the headline rates."""
    agg   = results["aggregated"]
    diffs = results["config"]["difficulties"]
    pipes = ["agentic"] + (["nonagentic"] if results["config"]["compare"] else [])

    def chart(metric, title, ylabel, fname):
        fig, ax = plt.subplots(figsize=(7, 4.5))
        x, width = range(len(diffs)), 0.35
        for pi, p in enumerate(pipes):
            vals = [((agg.get(d, {}).get(p) or {}).get(metric) or 0) for d in diffs]
            off  = (pi - (len(pipes) - 1) / 2) * width
            ax.bar([i + off for i in x], vals, width, label=p)
        ax.set_xticks(list(x)); ax.set_xticklabels([d.upper() for d in diffs])
        ax.set_ylabel(ylabel); ax.set_title(title); ax.legend()
        fig.tight_layout()
        fig.savefig(os.path.join(output_dir, fname), bbox_inches="tight")
        plt.show(); plt.close(fig)

    chart("independent_match_rate", "Answer Accuracy (independent re-solve)", "match rate (0-1)", "bulk_answer_accuracy.png")
    chart("verified_correct_rate", "Answer Correct (judge-verified)", "rate (0-1)", "bulk_answer_verified.png")
    chart("accept_rate", "Judge Accept Rate by Difficulty", "accept rate (0-1)", "bulk_accept_rate.png")
    chart("duplicate_rate", "Duplicate Rate by Difficulty", "duplicate rate (0-1)", "bulk_duplicate_rate.png")
    chart("mean_overall", "Mean Judge Score by Difficulty", "overall (1-5)", "bulk_mean_overall.png")


def _render_eval_charts(results, output_dir):
    """Renders 4 matplotlib bar charts comparing agentic vs non-agentic
    pipelines, saves them as PNGs to output_dir, and calls plt.show() so Colab
    displays them inline. Returns the list of saved file paths."""
    quiz_types = results["config"]["quiz_types"]
    agg        = results["aggregated"]
    saved      = []

    def _save(fig, name):
        path = os.path.join(output_dir, name)
        fig.savefig(path, bbox_inches="tight")
        plt.show()
        plt.close(fig)
        saved.append(path)

    def _grouped_bar(metric_key, title, ylabel, fname):
        agentic_vals    = [agg[qt]["agentic"].get(metric_key) or 0 for qt in quiz_types]
        nonagentic_vals = [agg[qt]["nonagentic"].get(metric_key) or 0 for qt in quiz_types]
        x     = range(len(quiz_types))
        width = 0.35
        fig, ax = plt.subplots(figsize=(7, 4.5))
        ax.bar([i - width / 2 for i in x], agentic_vals, width, label="Agentic")
        ax.bar([i + width / 2 for i in x], nonagentic_vals, width, label="Non-agentic")
        ax.set_xticks(list(x))
        ax.set_xticklabels([qt.upper() for qt in quiz_types])
        ax.set_ylabel(ylabel)
        ax.set_title(title)
        ax.legend()
        fig.tight_layout()
        _save(fig, fname)

    _grouped_bar("mean_overall", "Mean Judge Score — Agentic vs Non-Agentic",
                 "Mean overall score (1-5)", "eval_chart_mean_judge_score.png")
    _grouped_bar("judge_accept_rate", "Judge Accept Rate — Agentic vs Non-Agentic",
                 "Accept rate (0-1)", "eval_chart_accept_rate.png")
    _grouped_bar("duplicate_rate", "Duplicate Rate — Agentic vs Non-Agentic",
                 "Duplicate rate (0-1)", "eval_chart_duplicate_rate.png")

    # Validator-vs-judge agreement — agentic only, no non-agentic counterpart.
    fig, ax = plt.subplots(figsize=(7, 4.5))
    x          = range(len(quiz_types))
    width      = 0.25
    precisions = [agg[qt]["f1_metrics"].get("precision") or 0 for qt in quiz_types]
    recalls    = [agg[qt]["f1_metrics"].get("recall") or 0 for qt in quiz_types]
    f1s        = [agg[qt]["f1_metrics"].get("f1") or 0 for qt in quiz_types]
    ax.bar([i - width for i in x], precisions, width, label="Precision")
    ax.bar(list(x), recalls, width, label="Recall")
    ax.bar([i + width for i in x], f1s, width, label="F1")
    ax.set_xticks(list(x))
    ax.set_xticklabels([qt.upper() for qt in quiz_types])
    ax.set_ylabel("Score (0-1)")
    ax.set_title("Validator-vs-Judge Agreement (Agentic only)")
    ax.legend()
    fig.tight_layout()
    _save(fig, "eval_chart_validator_judge_f1.png")

    return saved


print("✅ Evaluation harness ready.")

# ══════════════════════════════════════════════════════════════
# ORCHESTRATOR AGENT
# ══════════════════════════════════════════════════════════════



class OrchestratorAgent:
    """The brain — decides which agents to run, in what order,
    with what parameters, and handles failures gracefully.

    Agents inside the same step-list entry run in parallel threads.
    Each step waits for the previous step to finish before starting."""

    # Each entry is a list of agent classes that can run concurrently.
    # Entries themselves execute sequentially (step 1 finishes → step 2 starts).
    PIPELINE = {
        "video": [
            [ExtractorAgent, TranscriberAgent],                # parallel: OCR + Whisper
            [SummaryAgent],
            [RAGBuilderAgent, HintsAgent],                     # parallel: index + hints
            [FlashcardAgent, SuggestedQuestionsAgent],         # parallel: cards + suggestions
        ],
        "document": [
            [DocParserAgent],
            [SummaryAgent],
            [RAGBuilderAgent, HintsAgent],                     # parallel: index + topic ranking
            [FlashcardAgent, SuggestedQuestionsAgent],         # parallel: cards + suggestions
        ],
    }

    def _run_agent(self, agent, ctx):
        """Run a single agent, catching crashes so the pipeline continues."""
        try:
            agent.run(ctx)
        except Exception as e:
            self._log(ctx, "ERROR", f"{agent.name} crashed: {e} — continuing")

    def _run_parallel_group(self, agent_classes, ctx):
        """Run a group of agents concurrently and wait for all to finish."""
        agents = [cls() for cls in agent_classes]

        if len(agents) == 1:
            self._run_agent(agents[0], ctx)
            return

        # multiple agents in this step — fan out into threads
        self._log(ctx, "PLAN",
                  f"Running in parallel: {', '.join(a.name for a in agents)}")
        with ThreadPoolExecutor(max_workers=len(agents)) as pool:
            futs = {pool.submit(self._run_agent, a, ctx): a for a in agents}
            for fut in as_completed(futs):
                fut.result()   # propagate exceptions if any

    def run(self, content_type, file_path):
        """Execute the full agentic pipeline with parallel steps. Times the
        overall run and every step so speed shows up in the agent log
        alongside the per-agent timing BaseAgent already records."""
        pipeline_start = time.time()
        ctx = AgentContext(content_type, file_path)
        state.session_id = state._make_session_id(ctx.source_name)
        state.agent_log  = []
        state.quiz_history = []

        pipeline = self.PIPELINE.get(content_type, [])
        total_steps = len(pipeline)

        self._log(ctx, "START",
                  f"Pipeline for {content_type}: {ctx.source_name} "
                  f"({total_steps} steps, parallel where possible)")

        for step_idx, agent_group in enumerate(pipeline):
            names = "+".join(cls.name if hasattr(cls, 'name') else cls.__name__
                            for cls in agent_group)
            pct = int(((step_idx + 0.5) / total_steps) * 100)
            state.update_status(
                state="running", stage=f"{names}…", pct=min(pct, 98),
            )
            step_start = time.time()
            self._run_parallel_group(agent_group, ctx)
            step_sec = time.time() - step_start
            self._log(
                ctx, "DONE",
                f"Step {step_idx + 1}/{total_steps} ({names}) finished in {step_sec:.1f}s",
                duration_sec=round(step_sec, 2),
            )

        # finalize
        timeline = ""
        if content_type == "video":
            timeline = build_lecture_timeline_text(ctx.segment_list, ctx.board_entries)
            if timeline:
                save_text(
                    "lecture_timeline.txt",
                    "=== LECTURE TIMELINE (speech + slides & boards) ===\n\n" + timeline,
                )
        state.update_status(
            state="done", stage="Complete ✓", pct=100,
            summary=ctx.summary,
            transcript=ctx.transcript,
            board_text=ctx.board_text,
            lecture_timeline=timeline,
        )
        state.agent_log = ctx.agent_log
        state.save_session(ctx.source_name)
        total_sec = time.time() - pipeline_start
        self._log(
            ctx, "DONE",
            f"Pipeline complete in {total_sec:.1f}s — {len(ctx.agent_log)} agent decisions logged",
            duration_sec=round(total_sec, 2),
        )

    def _log(self, ctx, level, msg, duration_sec=None):
        entry = {
            "agent": "Orchestrator", "level": level,
            "time": datetime.now().isoformat(), "message": msg,
        }
        if duration_sec is not None:
            entry["duration_sec"] = duration_sec
        ctx.agent_log.append(entry)
        print(f"  [Orchestrator] {msg}")



orchestrator = OrchestratorAgent()



print("✅ Orchestrator ready.")

# ══════════════════════════════════════════════════════════════
# DIFFICULTY ADAPTER  (self-improving quiz difficulty)
# ══════════════════════════════════════════════════════════════



class DifficultyAdapter:
    """Tracks student performance per topic and generates targeted
    questions for weak areas at adjusted difficulty."""

    def record_performance(self, grade_results):
        """Called after every quiz submission.  Stores per-topic score history."""
        for r in grade_results.get("results", []):
            topic = r.get("topic", "General")
            if topic not in state.student_performance:
                state.student_performance[topic] = []
            state.student_performance[topic].append(1.0 if r.get("correct") else 0.0)

    def get_weak_topics(self, threshold=0.5):
        """Return topics where average score is below threshold."""
        weak = {}
        for topic, scores in state.student_performance.items():
            if not scores:
                continue
            avg = sum(scores) / len(scores)
            if avg < threshold:
                weak[topic] = round(avg, 2)
        return weak

    def get_all_stats(self):
        """Return full performance stats for every topic."""
        stats = {}
        for topic, scores in state.student_performance.items():
            if not scores:
                continue
            avg = sum(scores) / len(scores)
            stats[topic] = {
                "average": round(avg, 2),
                "attempts": len(scores),
                "correct": sum(1 for s in scores if s == 1.0),
                "total": len(scores),
                "level": "weak" if avg < 0.5 else "ok" if avg < 0.8 else "strong",
            }
        return stats

    def suggest_difficulty(self, topic):
        """Based on past performance, return recommended difficulty."""
        scores = state.student_performance.get(topic, [])
        if not scores:
            return "medium"
        avg = sum(scores) / len(scores)
        if avg > 0.8:
            return "hard"
        if avg < 0.4:
            return "easy"
        return "medium"

    def generate_targeted(self, quiz_type="mcq", num=5):
        """Generate new questions focused on weak topics at adjusted difficulty."""
        weak = self.get_weak_topics()
        if not weak:
            return []

        topic_list = ", ".join(weak.keys())
        diff_hints = ", ".join(
            f"{t}: {self.suggest_difficulty(t)}" for t in weak
        )
        summary = state.status.get("summary", "")
        if not summary:
            return []

        prompt = (
            f"Generate {num} {quiz_type.upper()} questions focused specifically on "
            f"these WEAK topics the student struggles with: {topic_list}.\n"
            f"Difficulty guidance per topic: {diff_hints}.\n"
            f"The student has been getting these wrong — make questions that help them learn.\n\n"
        )

        if quiz_type == "mcq":
            prompt += (
                "STRICT JSON: {\"questions\":[{\"question\",\"options\":{\"A\",\"B\",\"C\",\"D\"},"
                "\"correct_answer\",\"explanation\",\"topic\",\"difficulty\",\"bloom_level\"}]}\n\n"
            )
        elif quiz_type == "tf":
            prompt += (
                "STRICT JSON: {\"questions\":[{\"statement\",\"answer\",\"explanation\","
                "\"topic\",\"difficulty\"}]}\n\n"
            )
        elif quiz_type == "fill":
            prompt += (
                "STRICT JSON: {\"questions\":[{\"question\",\"answer\",\"hint\","
                "\"topic\",\"difficulty\"}]}\n\n"
            )
        else:
            prompt += (
                "STRICT JSON: {\"questions\":[{\"question\",\"model_answer\",\"key_points\","
                "\"marks\",\"topic\",\"difficulty\"}]}\n\n"
            )

        if state.video_path:
            be = load_board_entries_from_disk()
            bt = state.status.get("board_text", "")
            tr = state.status.get("transcript", "")
            content = (
                pack_video_sources(bt, tr, be, 5500)
                + "\n\n=== SUMMARY ===\n"
                + summary[:4000]
            )
        else:
            content = summary[:3000]
        prompt += f"CONTENT:\n{content[:8000]}"
        data = call_ollama_json_quiz(prompt, key="questions")
        return data.get("questions", []) if isinstance(data, dict) else []



difficulty_adapter = DifficultyAdapter()



print("✅ Difficulty adapter ready.")

# ══════════════════════════════════════════════════════════════
# AUTO-RESTORE  (load last session from Drive on startup)
# ══════════════════════════════════════════════════════════════



def _auto_restore():
    """Check Google Drive for saved sessions. Load the latest one
    so quizzes/summary/flashcards are available immediately."""
    saved = AppState.list_sessions()
    if not saved:
        print("📂 No saved sessions found — start fresh.")
        return

    latest = saved[0]
    state.load_session(latest["session_id"])
    print(f"🔄 Auto-loaded last session: {latest['source_file']} ({latest['session_id']})")
    print(f"   📊 {latest['mcq_count']} MCQ · {latest['tf_count']} T/F · {latest['flash_count']} flashcards")

    # Rebuild the RAG search index from saved segment/board files
    rag_entries = []
    seg_path   = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    board_path = os.path.join(OUTPUT_DIR, "board_entries.json")
    if os.path.exists(seg_path):
        with open(seg_path, "r", encoding="utf-8") as f:
            rag_entries += json.load(f)
    if os.path.exists(board_path):
        with open(board_path, "r", encoding="utf-8") as f:
            for e in json.load(f):
                rag_entries.append({
                    "text": e["text"], "start_str": e["timestamp_str"],
                    "tag": e.get("tag", ""),
                })
    if rag_entries:
        build_rag(rag_entries, ts_key="start_str", src_key="tag")
        print(f"   🔍 RAG index rebuilt with {len(rag_entries)} chunks")



_auto_restore()

# ══════════════════════════════════════════════════════════════
# FLASK API
# ══════════════════════════════════════════════════════════════



app = Flask(__name__)
CORS(app)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024   # 500 MB

# ── Health / Status ───────────────────────────────────────────



@app.route('/')
def index():
    return jsonify(
        app="LectureForge v1.1",
        status="API running",
        model=state.active_model,
        message="Open quizforge_ui.html locally to use the UI.",
    )



@app.route('/status')
def status_route():
    """Pipeline status + live agentic-quiz thinking events.
    The UI uses `quiz_thinking_session` to detect a fresh generation run."""
    payload = dict(state.status)
    payload["model"]                  = state.active_model
    payload["session_id"]             = state.session_id
    payload["quiz_thinking"]          = state._quiz_agent_thinking
    payload["quiz_thinking_session"]  = state._quiz_thinking_session
    payload["quiz_thinking_active"]   = state._quiz_thinking_active
    return jsonify(payload)

# ── Upload ────────────────────────────────────────────────────



@app.route('/upload/video', methods=['POST'])
def upload_video():
    f = request.files.get('video')
    if not f:
        return jsonify(ok=False, error="No file sent.")
    ext = Path(f.filename).suffix.lower().lstrip('.')
    if ext not in ALLOWED_VIDEO:
        return jsonify(ok=False, error=f"Unsupported type: {ext}")
    path = os.path.join(UPLOAD_DIR, secure_filename(f.filename))
    f.save(path)
    state.video_path = path
    state.reset_status()
    threading.Thread(target=orchestrator.run, args=("video", path), daemon=True).start()
    return jsonify(ok=True, filename=f.filename)



@app.route('/upload/document', methods=['POST'])
def upload_document():
    f = request.files.get('document')
    if not f:
        return jsonify(ok=False, error="No file sent.")
    ext = Path(f.filename).suffix.lower().lstrip('.')
    if ext not in ALLOWED_DOCS:
        return jsonify(ok=False, error=f"Unsupported type: {ext}")
    path = os.path.join(UPLOAD_DIR, secure_filename(f.filename))
    f.save(path)
    state.doc_path = path
    state.reset_status()
    threading.Thread(target=orchestrator.run, args=("document", path), daemon=True).start()
    return jsonify(ok=True, filename=f.filename)

# ── Quiz endpoints ────────────────────────────────────────────



@app.route('/quiz/mcq')
def quiz_mcq_route():
    qs = state.quiz.get("mcq", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/true_false')
def quiz_tf_route():
    qs = state.quiz.get("tf", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/fill_blank')
def quiz_fill_route():
    qs = state.quiz.get("fill", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/short_answer')
def quiz_short_route():
    qs = state.quiz.get("short", [])
    return jsonify(questions=qs, total=len(qs))



@app.route('/quiz/history')
def quiz_history_route():
    """Return past graded quiz submissions saved into this session."""
    return jsonify(history=state.quiz_history, total=len(state.quiz_history))


def _new_quiz_run_id():
    return f"run_{int(time.time() * 1000)}_{os.urandom(4).hex()}"


def _append_generated_quiz_run(source, quiz_type, mode, difficulty, questions):
    """Store a snapshot of a generated quiz so the UI can reload past sets."""
    if not questions:
        return
    with state.lock:
        run = {
            "id":          _new_quiz_run_id(),
            "created_at":  datetime.now().isoformat(),
            "source":      source,          # "async" | "upload"
            "quiz_type":   quiz_type,
            "mode":        mode,            # auto | adaptive | manual
            "difficulty":  difficulty or "",
            "count":       len(questions),
            "questions":   copy.deepcopy(questions),
        }
        state.generated_quiz_runs.insert(0, run)
        state.generated_quiz_runs = state.generated_quiz_runs[:50]


@app.route('/quiz/generated_runs')
def quiz_generated_runs_route():
    with state.lock:
        runs = list(state.generated_quiz_runs)
    light = []
    for r in runs:
        light.append({
            "id":         r.get("id"),
            "created_at": r.get("created_at"),
            "source":     r.get("source"),
            "quiz_type":  r.get("quiz_type"),
            "mode":       r.get("mode"),
            "difficulty": r.get("difficulty"),
            "count":      r.get("count", len(r.get("questions", []))),
        })
    return jsonify(runs=light, total=len(light))


@app.route('/quiz/generated_runs/restore', methods=['POST'])
def quiz_restore_generated_run_route():
    data   = request.get_json(silent=True) or {}
    run_id = (data.get("run_id") or "").strip()
    if not run_id:
        return jsonify(ok=False, error="run_id required"), 400
    with state.lock:
        found = None
        for r in state.generated_quiz_runs:
            if r.get("id") == run_id:
                found = r
                break
        if not found:
            return jsonify(ok=False, error="Run not found"), 404
        qt = found.get("quiz_type")
        qs = found.get("questions") or []
        if qt not in state.quiz:
            return jsonify(ok=False, error="Invalid quiz type in run"), 400
        state.quiz[qt] = copy.deepcopy(qs)
    return jsonify(ok=True, quiz_type=qt, total=len(qs))



def _load_segments_and_board_artifacts():
    segs = []
    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    if os.path.exists(seg_path):
        try:
            with open(seg_path, "r", encoding="utf-8") as f:
                segs = json.load(f)
        except Exception:
            segs = []
    be = load_board_entries_from_disk()
    return segs, be



def _record_quiz_submission(quiz_type, request_answers, result):
    # Ensure each graded question has a source_timestamp for UI clip links.
    segs, be = _load_segments_and_board_artifacts()
    graded_items = result.get("results", [])
    if graded_items:
        fill_missing_timestamps(graded_items, segments=segs, board_entries=be)

    entry = {
        "created_at": datetime.now().isoformat(),
        "quiz_type": quiz_type,
        "mode": "auto",
        "count": result.get("total", len(graded_items)),
        "score": result.get("score", 0),
        "total": result.get("total", len(graded_items)),
        "percentage": result.get("percentage", 0),
        "grade": result.get("grade", ""),
        "answers": request_answers or {},
        "results": graded_items,
    }
    state.quiz_history.append(entry)



@app.route('/quiz/mcq/grade', methods=['POST'])
def grade_mcq_route():
    data   = request.get_json(silent=True) or {}
    result = grade_mcq(state.quiz.get("mcq", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("mcq", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/tf/grade', methods=['POST'])
def grade_tf_route():
    data   = request.get_json(silent=True) or {}
    result = grade_tf(state.quiz.get("tf", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("tf", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/fill/grade', methods=['POST'])
def grade_fill_route():
    data   = request.get_json(silent=True) or {}
    result = grade_fill(state.quiz.get("fill", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("fill", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/short/reveal', methods=['POST'])
def reveal_short_route():
    data = request.get_json(silent=True) or {}
    qs   = state.quiz.get("short", [])
    revealed = {i: qs[i] for i in (data.get("indices") or range(len(qs))) if i < len(qs)}
    return jsonify(revealed=revealed)



@app.route('/quiz/short/grade', methods=['POST'])
def grade_short_route():
    data   = request.get_json(silent=True) or {}
    result = grade_short(state.quiz.get("short", []), data.get("answers", {}))
    difficulty_adapter.record_performance(result)
    _record_quiz_submission("short", data.get("answers", {}), result)
    return jsonify(**result)



@app.route('/quiz/explain', methods=['POST'])
def explain_mistakes_route():
    data  = request.get_json(silent=True) or {}
    wrong = data.get("wrong_questions", [])
    exps  = explain_wrong(wrong)
    return jsonify(explanations=exps)

@app.route('/quiz/generate', methods=['POST'])
def quiz_generate_route():
    """Generate quiz using the full agentic loop: Plan → Generate → Validate → Refine.
    Body: {mode: "adaptive"|"manual", difficulty: "easy"|"medium"|"hard",
           quiz_type: "mcq"|"tf"|"fill"|"short", count: 10,
           instructions: "...", action: "new"|"refine"}"""
    data       = request.get_json(silent=True) or {}
    mode       = data.get("mode", "adaptive")
    difficulty = data.get("difficulty", "medium")
    quiz_type  = data.get("quiz_type", "mcq")
    action     = data.get("action", "new")
    num        = min(int(data.get("count", 10)), 20)
    user_instr = data.get("instructions", "").strip()
    plan_mode  = data.get("plan_mode", "auto")
    manual_topics = data.get("manual_topics", []) or []

    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No content processed yet — upload a file first.",
                       questions=[], thinking=[])

    ctx = type('Ctx', (), {
        'content_type':       'video' if state.video_path else 'document',
        'board_text':         state.status.get("board_text", ""),
        'transcript':         state.status.get("transcript", ""),
        'board_entries':      load_board_entries_from_disk(),
        'doc_text':           state.status.get("summary", ""),
        'summary':            summary,
        'agent_log':          [],
        '_quiz_type':         quiz_type,
        '_quiz_count':        num,
        '_user_instructions': user_instr,
        '_plan_mode':         plan_mode,
        '_manual_topics':     manual_topics,
        '_manual_difficulty': difficulty,
        '_mode':              mode,
        '_quiz_plan':         None,
        '_agent_name':        '',
        'segment_list':       [],
        'file_path':          state.video_path or state.doc_path or "",
        'source_name':        Path(state.video_path or state.doc_path or "unknown").name,
        'thinking_cb':        None,
    })()

    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    if os.path.exists(seg_path):
        try:
            with open(seg_path, "r", encoding="utf-8") as f:
                ctx.segment_list = json.load(f)
        except Exception:
            ctx.segment_list = []

    if state.doc_path:
        doc_path = os.path.join(OUTPUT_DIR, "doc_text.txt")
        if os.path.exists(doc_path):
            with open(doc_path, "r", encoding="utf-8") as f:
                ctx.doc_text = f.read()

    thinking   = []
    def thinking_cb(entry):
        thinking.append(entry)
    ctx.thinking_cb = thinking_cb

    content = (
        pack_video_sources(ctx.board_text, ctx.transcript, ctx.board_entries or [], 7500)
        if ctx.content_type == "video"
        else (ctx.doc_text or "")[:7500]
    )

    # -- Refine mode: send existing questions through Refiner + Validator --
    if action == "refine" and state.quiz.get(quiz_type):
        selected_topics = [t.strip() for t in (manual_topics or []) if str(t).strip()] if plan_mode == "manual_topics" else []
        existing_all = state.quiz[quiz_type]
        if selected_topics:
            existing = [
                q for q in existing_all
                if (q.get("topic") or "").strip() in selected_topics
            ]
            if not existing:
                return jsonify(ok=False, error="No existing questions match selected manual topics.", questions=[], thinking=[])
        else:
            existing = existing_all
        thinking_cb({
            "agent": "QuizRefinerAgent", "step": "refine_start", "verdict": "PASS",
            "phase": "PLAN",
            "reason": f"Refining {len(existing)} existing {quiz_type.upper()} questions with new instructions...",
        })
        thinking_cb({
            "agent": "QuizRefinerAgent", "step": "phase_start", "phase": "GENERATE", "verdict": "PASS",
            "reason": f"Generating refined set for {len(existing)} selected question(s)…",
            "data": {"total_slots": len(existing)},
        })
        weak = difficulty_adapter.get_weak_topics()
        weak_str = ", ".join(weak.keys()) if weak else "none"
        topics_line = ", ".join(selected_topics) if selected_topics else "auto"
        schema = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
        prompt = (
            f"Revise this {quiz_type.upper()} quiz to better follow the student's instructions.\n"
            f"Mode: {mode}  Difficulty: {difficulty.upper()}\n"
            f"Topic scope: {topics_line}\n"
            f"Weak topics to emphasize: {weak_str}\n"
            f"Student instructions: {user_instr or 'None'}\n\n"
            f"EXISTING QUIZ:\n{json.dumps(existing, ensure_ascii=False)}\n\n"
            f"CONTENT:\n{content[:6000]}\n\nSUMMARY:\n{summary[:1500]}\n\n"
            f"STRICT JSON: {schema}"
        )
        thinking_cb({"agent": "QuizRefinerAgent", "step": "llm_prompt", "phase": "GENERATE", "verdict": "PASS",
                     "reason": "Prompt sent to LLM for refine generation",
                     "data": {"prompt_preview": prompt[:1200]}})
        thinking_cb({"agent": "QuizRefinerAgent", "step": "generating", "verdict": "PASS",
                     "reason": "Sending revision prompt to LLM..."})
        data_out = call_ollama_json_quiz(prompt, key="questions")
        qs = data_out.get("questions", []) if isinstance(data_out, dict) else []
        # validate the refined output
        passed, failed = QuizValidatorAgent().validate_batch(
            qs, quiz_type, content[:800], difficulty, user_instr, thinking_cb=thinking_cb
        )
        if failed:
            fixed = QuizRefinerAgent().refine_batch(
                failed, quiz_type, content[:800], difficulty, user_instr, thinking_cb=thinking_cb
            )
            qs = passed + fixed
        else:
            qs = passed
        fill_missing_timestamps(qs, segments=ctx.segment_list, board_entries=ctx.board_entries)
        state.quiz[quiz_type] = qs
        thinking_cb({"agent": "QuizRefinerAgent", "step": "final", "verdict": "DONE",
                     "reason": f"Refined quiz ready -- {len(qs)} questions"})

    # -- Adaptive mode --
    elif mode == "adaptive":
        ctx._mode = "adaptive"
        agent = AdaptiveQuizAgent()
        agent.run(ctx)
        thinking.extend(agent._thinking_log)
        fill_missing_timestamps(
            state.quiz.get(quiz_type, []),
            segments=ctx.segment_list, board_entries=ctx.board_entries,
        )

    # -- Manual mode --
    elif mode == "manual":
        agent_map = {"easy": EasyQuizAgent, "medium": MediumQuizAgent, "hard": HardQuizAgent}
        agent = agent_map.get(difficulty, MediumQuizAgent)()
        ctx._mode       = "manual"
        ctx._agent_name = agent.name
        thinking_cb({"agent": agent.name, "step": "start", "verdict": "PASS",
                     "reason": f"Manual {difficulty.upper()} mode -- {num} {quiz_type.upper()} questions"})
        if user_instr:
            thinking_cb({"agent": agent.name, "step": "user_instructions", "verdict": "PASS",
                         "reason": f"Student instructions: {user_instr}"})

        max_attempts = 3
        final = _run_quiz_agentic_loop(
            ctx, quiz_type, content, summary, user_instr,
            target_difficulty=difficulty, thinking_cb=thinking_cb,
            max_attempts_per_q=max_attempts,
        )
        for q in final:
            q["difficulty"] = difficulty
        fill_missing_timestamps(final, segments=ctx.segment_list, board_entries=ctx.board_entries)
        state.quiz[quiz_type] = final

    else:
        return jsonify(ok=False, error=f"Unknown mode: {mode}", questions=[], thinking=[])

    state._quiz_agent_thinking = thinking
    final_qs = state.quiz.get(quiz_type, [])
    return jsonify(ok=True, questions=final_qs, total=len(final_qs),
                   thinking=thinking, mode=mode,
                   difficulty=difficulty if mode == "manual" else "adaptive")

# ── Study endpoints (flashcards, hints, chat) ─────────────────



@app.route('/flashcards')
def flashcards_route():
    cards = state.quiz.get("flash", [])
    return jsonify(cards=cards, total=len(cards))


@app.route('/flashcards/generate', methods=['POST'])
def flashcards_generate_route():
    """Generate flashcards on-demand from the Flashcards section."""
    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No processed content yet. Upload and process first."), 400

    if state.video_path:
        cards = generate_flashcards(
            summary,
            board_text=state.status.get("board_text", ""),
            board_entries=load_board_entries_from_disk(),
        )
    else:
        cards = generate_flashcards(summary, doc_text=summary)

    # Deduplicate by front text
    seen, unique = set(), []
    for c in cards or []:
        if not isinstance(c, dict):
            continue
        front = (c.get("front") or "").strip().lower()
        if front and front not in seen:
            seen.add(front)
            unique.append(c)
    state.quiz["flash"] = unique
    return jsonify(ok=True, cards=unique, total=len(unique))



@app.route('/exam_hints')
def exam_hints_route():
    return jsonify(**state.quiz.get("hints", {}))


@app.route('/quiz/topics')
def quiz_topics_route():
    """Return topic options for manual topic plan mode."""
    topics = []

    # 1) AI-ranked exam topics
    hints = state.quiz.get("hints", {}) or {}
    ai = hints.get("ai_analysis", {}) if isinstance(hints, dict) else {}
    ranked = ai.get("ai_important_topics", []) if isinstance(ai, dict) else []
    for item in ranked[:20]:
        if isinstance(item, dict):
            t = (item.get("topic") or "").strip()
            if t:
                topics.append(t)

    # 2) performance-tracked topics
    for t in list(difficulty_adapter.get_all_stats().keys())[:20]:
        if t and t not in topics:
            topics.append(t)

    # 3) weak topics
    for t in list(difficulty_adapter.get_weak_topics().keys())[:20]:
        if t and t not in topics:
            topics.append(t)

    # 4) sensible fallback
    if not topics:
        topics = ["Core Concepts", "Key Principles", "Applications", "Problem Solving"]

    return jsonify(topics=topics[:30], total=len(topics[:30]))



@app.route('/suggested_questions')
def suggested_questions_route():
    return jsonify(questions=state.suggested_questions)



GREETING_WORDS = {"hi","hello","hey","sup","yo","thanks","thank you","bye","ok","okay","good morning","good evening","hola"}



# ── Socratic tutor prompts ─────────────────────────────────────

SOCRATIC_SYSTEM = """You are a Socratic AI lecture tutor. Your role is to help students \
DISCOVER answers through guided questioning, not by giving answers directly.

STRICT RULES:
1. On the FIRST exchange about a topic: NEVER give the answer directly.
   Instead, ask 1-2 short probing questions that activate prior knowledge.
2. When a student attempts an answer: acknowledge what is correct, gently correct \
   what is wrong, then ask a follow-up question that moves them closer to the full answer.
3. Give the COMPLETE answer only when:
   - The student has made 2 or more genuine attempts, OR
   - The student says something like "just tell me", "I give up", "I don't know", \
     "tell me the answer", or "what is it".
4. Every Socratic response must end with a question mark — you are always probing.
5. Keep probing questions short (1-2 sentences) and focus on ONE concept at a time.
6. Reference lecture content naturally: "The lecture mentioned at [timestamp]…"
7. Be warm and encouraging — never make the student feel bad for not knowing.

EXAMPLE:
Student: "What is gradient descent?"
BAD response: "Gradient descent is an optimization algorithm that…"
GOOD response: "Great question! Let's think about it together. Imagine you're \
standing on a hilly landscape and want to reach the lowest point — what strategy \
would you use to get there?"
"""

DIRECT_SYSTEM = "You are a helpful lecture tutor. Answer clearly and concisely. Cite timestamps when possible."

# phrases that signal the student wants the answer directly
_GIVE_UP_PHRASES = (
    "just tell me", "tell me the answer", "i give up", "i don't know",
    "what is it", "what's the answer", "i have no idea", "please just answer",
    "stop asking", "skip the questions", "idk", "no idea", "بলো", "বলো",
)

def _student_gave_up(text):
    t = text.lower().strip()
    return any(ph in t for ph in _GIVE_UP_PHRASES)

def _build_history_block(history):
    """Convert [{role, content}] list into a readable conversation string."""
    if not history:
        return ""
    lines = []
    for turn in history[-6:]:   # last 6 turns = 3 exchanges
        role    = "Student" if turn.get("role") == "user" else "Tutor"
        content = str(turn.get("content", "")).strip()[:400]
        lines.append(f"{role}: {content}")
    return "=== RECENT CONVERSATION ===\n" + "\n".join(lines)

def _count_student_attempts(history):
    """Count how many times the student has already responded (not first-time)."""
    return sum(1 for t in history if t.get("role") == "user")



@app.route('/chat', methods=['POST'])
def chat_route():
    data    = request.get_json(silent=True) or {}
    q       = data.get("question", "").strip()
    history = data.get("history",  [])        # [{role:"user"|"assistant", content:"..."}]
    mode    = data.get("mode",     "socratic") # "socratic" | "direct"

    if not q:
        return jsonify(answer="Please ask a question.", timestamps=[], mode_used=mode, socratic=False)

    # ── greeting shortcut ──────────────────────────────────────
    if q.lower().strip("!?., ") in GREETING_WORDS or (len(q.split()) <= 2 and not any(c.isdigit() for c in q)):
        answer = call_ollama(
            f"The student said: \"{q}\"\nRespond with a brief friendly greeting. "
            "You are their Socratic lecture tutor. Keep it to 1-2 sentences.",
            system="You are a friendly lecture tutor.",
        )
        return jsonify(answer=answer, timestamps=[], mode_used="direct", socratic=False)

    # ── RAG retrieval + grounding ──────────────────────────────
    ctx, tss = rag.query(q, top_k=5)
    summary_snip = (state.status.get("summary") or "")[:2000]
    extra = []
    if summary_snip:
        extra.append(f"=== LECTURE SUMMARY ===\n{summary_snip}")
    if state.video_path:
        be  = load_board_entries_from_disk()
        bt  = state.status.get("board_text", "")
        vis = format_visual_lecture_text(bt, be, 2800)
        if vis:
            extra.append(f"=== SLIDES + BOARDS (OCR) ===\n{vis}")
    grounding = "\n\n".join(extra)

    # ── clip duration estimation ───────────────────────────────
    for i, ts in enumerate(tss):
        chunk_idx = None
        try:
            chunk_idx = rag.timestamps.index(ts)
        except (ValueError, AttributeError):
            pass
        words = len(rag.chunks[chunk_idx].split()) if (chunk_idx is not None and chunk_idx < len(rag.chunks)) else 50
        tss[i] = {**ts, "clip_dur": max(10, min(60, int(words / 2.5)))}

    # ── decide whether to be Socratic or direct ────────────────
    attempts       = _count_student_attempts(history)
    gave_up        = _student_gave_up(q)
    use_socratic   = (mode == "socratic") and (not gave_up) and (attempts < 2)
    history_block  = _build_history_block(history)

    if use_socratic:
        prompt = (
            "Use the lecture material below to guide a Socratic dialogue.\n"
            "Do NOT answer directly — ask probing questions that help the student think.\n\n"
            f"{grounding}\n\n"
            f"=== RELEVANT SNIPPETS ===\n{ctx}\n\n"
            f"{history_block}\n\n"
            f"Student's latest message: {q}\n\n"
            "Your response (Socratic — end with a question):"
        )
        system = SOCRATIC_SYSTEM
    else:
        # direct answer: student gave up, or already tried 2+ times, or mode is direct
        reveal_note = ""
        if gave_up:
            reveal_note = "The student has indicated they need the answer directly. Give a clear, complete explanation.\n\n"
        elif attempts >= 2:
            reveal_note = f"The student has attempted this {attempts} times. Now give a full, clear explanation.\n\n"

        prompt = (
            f"{reveal_note}"
            "Answer using the lecture material below. Be clear and cite timestamps.\n\n"
            f"{grounding}\n\n=== RELEVANT SNIPPETS ===\n{ctx}\n\n"
            f"{history_block}\n\n"
            f"QUESTION: {q}"
        )
        system = DIRECT_SYSTEM

    answer = call_ollama(prompt, system=system)
    return jsonify(
        answer    = answer,
        timestamps= tss,
        mode_used = "socratic" if use_socratic else "direct",
        socratic  = use_socratic,
    )

# ── Lesson Planner ────────────────────────────────────────────



@app.route('/lesson_plan', methods=['POST'])
def lesson_plan_route():
    data  = request.get_json(silent=True) or {}
    topic = data.get("topic", "")
    level = data.get("level", "undergraduate")
    dur   = data.get("duration", 45)
    if not topic:
        return jsonify(error="No topic provided.")
    prompt = (
        f"Create a complete {dur}-minute lesson plan about '{topic}' for {level} students.\n"
        "Include: Learning Objectives, Prerequisites, Introduction (5 min), "
        "Main Content with timed sub-topics, Activities/Exercises, "
        "Summary & Key Takeaways, Assessment Questions, and Homework.\n"
        "Format in clear Markdown with headers."
    )
    plan = call_ollama(prompt, system="You are an expert curriculum designer.")
    return jsonify(plan=plan, topic=topic, level=level, duration=dur)

# ── Sessions (list / load / delete past processed files) ──────



@app.route('/sessions')
def list_sessions_route():
    return jsonify(sessions=AppState.list_sessions())



@app.route('/sessions/load', methods=['POST'])
def load_session_route():
    data = request.get_json(silent=True) or {}
    sid  = data.get("session_id", "")
    if not sid:
        return jsonify(ok=False, error="No session_id provided.")
    ok = state.load_session(sid)
    if not ok:
        return jsonify(ok=False, error=f"Session '{sid}' not found.")

    # Rebuild RAG from restored transcript/board data
    rag_entries = []
    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    board_path = os.path.join(OUTPUT_DIR, "board_entries.json")
    if os.path.exists(seg_path):
        with open(seg_path, "r", encoding="utf-8") as f:
            rag_entries += json.load(f)
    if os.path.exists(board_path):
        with open(board_path, "r", encoding="utf-8") as f:
            for e in json.load(f):
                rag_entries.append({
                    "text": e["text"], "start_str": e["timestamp_str"],
                    "tag": e.get("tag", ""),
                })
    if rag_entries:
        build_rag(rag_entries, ts_key="start_str", src_key="tag")

    return jsonify(ok=True, session_id=sid, source_file=data.get("source_file", ""))



@app.route('/sessions/delete', methods=['POST'])
def delete_session_route():
    data = request.get_json(silent=True) or {}
    sid  = data.get("session_id", "")
    if not sid:
        return jsonify(ok=False, error="No session_id provided.")
    folder = os.path.join(SESSIONS_DIR, sid)
    if os.path.isdir(folder):
        shutil.rmtree(folder)
        return jsonify(ok=True, deleted=sid)
    return jsonify(ok=False, error="Session not found."), 404



@app.route('/sessions/current')
def current_session_route():
    return jsonify(
        session_id=state.session_id,
        source_file=Path(state.video_path or state.doc_path or "").name,
        has_data=state.status.get("state") == "done",
    )

# ── Agentic endpoints ─────────────────────────────────────────



@app.route('/agent_log')
def agent_log_route():
    return jsonify(log=state.agent_log, total=len(state.agent_log))



@app.route('/quiz/regenerate', methods=['POST'])
def quiz_regenerate_route():
    data      = request.get_json(silent=True) or {}
    quiz_type = data.get("type", "mcq")
    num       = min(int(data.get("count", 5)), 20)
    qs        = difficulty_adapter.generate_targeted(quiz_type=quiz_type, num=num)
    if not qs:
        weak = difficulty_adapter.get_weak_topics()
        if not weak:
            return jsonify(questions=[], message="No weak topics detected yet — take a quiz first.")
        return jsonify(questions=[], message="No summary available to generate from.")
    return jsonify(questions=qs, total=len(qs), weak_topics=difficulty_adapter.get_weak_topics())



@app.route('/quiz/regenerate_one', methods=['POST'])
def quiz_regenerate_one_route():
    """Regenerate a single question in-place using the per-question agentic loop.
    Body: {quiz_type, index, instructions?, difficulty?}.
    Streams the same Generate -> Validate -> Retry events the bulk loop emits."""
    data        = request.get_json(silent=True) or {}
    quiz_type   = data.get("quiz_type", "mcq")
    try:
        idx = int(data.get("index", 0))
    except (TypeError, ValueError):
        return jsonify(ok=False, error="Invalid index", thinking=[])
    user_instr  = (data.get("instructions") or "").strip()
    requested_diff = data.get("difficulty")

    qs = state.quiz.get(quiz_type, [])
    if not qs or idx < 0 or idx >= len(qs):
        return jsonify(ok=False, error="Invalid question index", thinking=[])

    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No content available", thinking=[])

    if state.video_path:
        be = load_board_entries_from_disk()
        bt = state.status.get("board_text", "")
        tr = state.status.get("transcript", "")
        content = pack_video_sources(bt, tr, be, 7500)
    else:
        content = summary[:7500]

    orig  = qs[idx]
    topic = orig.get("topic") or "General"
    target_difficulty = requested_diff or orig.get("difficulty") or "medium"

    thinking = []
    def push(entry): thinking.append(entry)

    push({
        "agent":       "Orchestrator",
        "step":        "regen_start",
        "phase":       "GENERATE",
        "verdict":     "PASS",
        "question_id": f"q{idx+1}",
        "reason":      f"Regenerating Q{idx+1} ({topic}, {target_difficulty}) on demand",
    })

    # exclude the question being replaced from the dedup set
    existing = [q for i, q in enumerate(qs) if i != idx]

    new_q = _generate_one_question_with_retry(
        quiz_type=quiz_type, topic=topic, difficulty=target_difficulty,
        content_snip=content, summary=summary,
        user_instr=user_instr,
        existing_questions=existing,
        question_idx=idx,
        max_attempts=3,
        thinking_cb=push,
        agent_name="QuizRegenerator",
    )

    if not new_q:
        return jsonify(ok=False, error="Failed to generate replacement",
                       thinking=thinking)

    seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
    segs = []
    if os.path.exists(seg_path):
        try:
            with open(seg_path, "r", encoding="utf-8") as f:
                segs = json.load(f)
        except Exception:
            segs = []
    be = load_board_entries_from_disk()
    fill_missing_timestamps([new_q], segments=segs, board_entries=be)

    qs[idx] = new_q
    state.quiz[quiz_type] = qs
    state._quiz_agent_thinking = thinking
    return jsonify(ok=True, question=new_q, index=idx,
                   thinking=thinking, total=len(qs))



@app.route('/quiz/thinking')
def quiz_thinking_route():
    """Return the thinking log from the last quiz generation."""
    return jsonify(thinking=state._quiz_agent_thinking,
                   total=len(state._quiz_agent_thinking))



def _summarize_quiz_thinking(events):
    """Roll the raw per-question thinking events of the LAST generation into a
    compact per-question metrics table (attempts, time, validator verdicts,
    remediations) plus totals — what the UI offers as a download."""
    per, order = {}, []
    for e in events or []:
        qid  = e.get("question_id")
        step = e.get("step")
        if not qid:
            continue
        if qid not in per:
            per[qid] = {
                "question_id": qid, "question_idx": e.get("question_idx"),
                "topic": e.get("topic", ""), "difficulty": e.get("difficulty", ""),
                "attempts": None, "duration_sec": None, "final_status": None,
                "validator_pass": 0, "validator_fail": 0,
                "tool_calls": 0, "remediations": [],
            }
            order.append(qid)
        row = per[qid]
        if e.get("topic"):      row["topic"] = e.get("topic")
        if e.get("difficulty"): row["difficulty"] = e.get("difficulty")
        if step == "question_validated":
            v = str(e.get("verdict", "")).upper()
            if   v == "PASS": row["validator_pass"] += 1
            elif v == "FAIL": row["validator_fail"] += 1
        elif step == "tool_call":
            row["tool_calls"] += 1
        elif step == "remediation":
            row["remediations"].append(e.get("action") or e.get("reason", ""))
        elif step == "question_finalized":
            row["attempts"]     = e.get("total_attempts", row["attempts"])
            row["duration_sec"] = e.get("duration_sec", row["duration_sec"])
            row["final_status"] = e.get("status", e.get("verdict"))
            if e.get("question_idx") is not None:
                row["question_idx"] = e.get("question_idx")

    rows = [per[q] for q in order]
    durs = [r["duration_sec"] for r in rows if isinstance(r["duration_sec"], (int, float))]
    atts = [r["attempts"]     for r in rows if isinstance(r["attempts"], (int, float))]
    totals = {
        "questions":          len(rows),
        "total_duration_sec": round(sum(durs), 2) if durs else 0.0,
        "avg_duration_sec":   round(sum(durs) / len(durs), 2) if durs else None,
        "avg_attempts":       round(sum(atts) / len(atts), 2) if atts else None,
        "first_try_count":    sum(1 for r in rows if r["attempts"] == 1),
        "total_tool_calls":   sum(r.get("tool_calls", 0) for r in rows),
        "total_remediations": sum(len(r["remediations"]) for r in rows),
    }
    return {"per_question": rows, "totals": totals}



@app.route('/quiz/last_run_metrics')
def quiz_last_run_metrics_route():
    """Compiled per-question metrics for the LAST quiz generation run — the
    thing the UI downloads after generating a quiz (distinct from the eval
    harness benchmark, which compares pipelines)."""
    events  = state._quiz_agent_thinking or []
    summary = _summarize_quiz_thinking(events)
    return jsonify(
        ok=True,
        generated_at=datetime.now().isoformat(),
        totals=summary["totals"],
        per_question=summary["per_question"],
        total_events=len(events),
        raw_events=events,
    )



@app.route('/quiz/score_current', methods=['POST'])
def quiz_score_current_route():
    """AI-judge ONLY the quiz questions ALREADY generated — not a benchmark, not
    new questions. quiz_type='all' (default) scores every type that has
    questions; a specific type scores just that one. Scores exactly what the
    student made and sees (1, 5, MCQ, TF, short — whatever exists)."""
    data        = request.get_json(silent=True) or {}
    quiz_type   = data.get("quiz_type", "all")
    target_diff = data.get("target_difficulty", "medium")
    judge_model = data.get("judge_model") or None

    valid = ("mcq", "tf", "fill", "short")
    types = [quiz_type] if quiz_type in valid else list(valid)

    by_type      = {}
    all_overalls = []
    all_accepts  = 0
    all_n        = 0

    for qt in types:
        qs = state.quiz.get(qt, [])
        if not qs:
            continue
        ctx, content, summary = _build_eval_ctx(qt, len(qs))
        content_snip = (content or summary or state.status.get("summary", ""))[:6000]

        scored = []
        for q in qs:
            diff = q.get("difficulty") or target_diff
            j = judge_question(qt, q, content_snip, diff, judge_model)
            scored.append({
                "question":         (q.get("question") or q.get("statement") or "")[:200],
                "topic":            q.get("topic", ""),
                "difficulty":       diff,
                "correctness":      j.get("correctness"),
                "clarity":          j.get("clarity"),
                "difficulty_match": j.get("difficulty_match"),
                "overall":          j.get("overall"),
                "accept":           bool(j.get("accept")),
                "reason":           j.get("reason", ""),
            })

        overalls = [s["overall"] for s in scored if isinstance(s["overall"], (int, float))]
        accepts  = sum(1 for s in scored if s["accept"])
        totals = {
            "n":            len(scored),
            "mean_overall": round(sum(overalls) / len(overalls), 2) if overalls else None,
            "accept_rate":  round(accepts / len(scored), 2) if scored else 0.0,
            "duplicate_rate": round(_duplicate_rate(qs), 3),
        }
        # add the same deterministic + embedding metrics the benchmark uses
        totals.update(_structural_metrics(qt, qs))
        totals.update(_grounding_metrics(qs, content_snip))
        by_type[qt] = {"scored": scored, "totals": totals}
        all_overalls += overalls
        all_accepts  += accepts
        all_n        += len(scored)

    if not by_type:
        return jsonify(ok=False, error="No questions generated yet — generate a quiz first.")

    overall = {
        "n":            all_n,
        "mean_overall": round(sum(all_overalls) / len(all_overalls), 2) if all_overalls else None,
        "accept_rate":  round(all_accepts / all_n, 2) if all_n else 0.0,
    }
    return jsonify(ok=True, by_type=by_type, overall=overall,
                   judge_model=judge_model or OLLAMA_MODEL)



def _new_job_id():
    short_hash = hashlib.md5(f"{time.time()}_{os.urandom(8)}".encode()).hexdigest()[:10]
    return f"job_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{short_hash}"



def _job_push_thinking(job_id, entry):
    with state.quiz_generation_lock:
        job = state.quiz_generation_jobs.get(job_id)
        if not job:
            return
        job["thinking"].append(entry)



@app.route('/quiz/generate_async', methods=['POST'])
def quiz_generate_async_route():
    """Kick off quiz generation in a background thread.
    Frontend polls /quiz/generate_job?job_id=... for live thinking entries.
    Uses the full Plan → Generate → Validate → Refine agentic loop."""
    data       = request.get_json(silent=True) or {}
    mode       = data.get("mode", "adaptive")
    difficulty = data.get("difficulty", "medium")
    quiz_type  = data.get("quiz_type", "mcq")
    action     = data.get("action", "new")
    num        = min(int(data.get("count", 10)), 20)
    user_instr = data.get("instructions", "").strip()
    plan_mode  = data.get("plan_mode", "auto")
    manual_topics = data.get("manual_topics", []) or []

    summary = state.status.get("summary", "")
    if not summary:
        return jsonify(ok=False, error="No content processed yet — upload a file first.", job_id=None)

    job_id = _new_job_id()
    with state.quiz_generation_lock:
        state.quiz_generation_jobs[job_id] = {
            "status": "running", "thinking": [], "result": None, "error": "",
        }

    # Bump the global thinking session so the UI also sees this run via /status
    state._quiz_thinking_session += 1
    state._quiz_thinking_active   = True
    state._quiz_agent_thinking    = []

    def worker():
        try:
            def push(entry):
                _job_push_thinking(job_id, entry)
                state._quiz_agent_thinking.append(entry)

            push({"agent": "Orchestrator", "step": "start", "verdict": "PASS",
                  "reason": f"{mode} mode started for {num} {quiz_type.upper()} questions — {difficulty.upper()} difficulty"})

            ctx = type('Ctx', (), {
                'content_type':       'video' if state.video_path else 'document',
                'board_text':         state.status.get("board_text", ""),
                'transcript':         state.status.get("transcript", ""),
                'board_entries':      load_board_entries_from_disk(),
                'doc_text':           state.status.get("summary", ""),
                'summary':            summary,
                'agent_log':          [],
                '_quiz_type':         quiz_type,
                '_quiz_count':        num,
                '_user_instructions': user_instr,
                '_plan_mode':         plan_mode,
                '_manual_topics':     manual_topics,
                '_manual_difficulty': difficulty,
                '_mode':              mode,
                '_quiz_plan':         None,
                '_agent_name':        '',
                'segment_list':       [],
                'file_path':          state.video_path or state.doc_path or "",
                'source_name':        Path(state.video_path or state.doc_path or "unknown").name,
                'thinking_cb':        push,
            })()

            seg_path = os.path.join(OUTPUT_DIR, "transcript_segments.json")
            if os.path.exists(seg_path):
                try:
                    with open(seg_path, "r", encoding="utf-8") as f:
                        ctx.segment_list = json.load(f)
                except Exception:
                    ctx.segment_list = []

            if state.doc_path:
                doc_entries_path = os.path.join(OUTPUT_DIR, "doc_text.txt")
                if os.path.exists(doc_entries_path):
                    with open(doc_entries_path, "r", encoding="utf-8") as f:
                        ctx.doc_text = f.read()

            content = (
                pack_video_sources(ctx.board_text, ctx.transcript,
                                   ctx.board_entries or [], 7500)
                if ctx.content_type == "video"
                else (ctx.doc_text or "")[:7500]
            )

            # Refine mode
            if action == "refine" and state.quiz.get(quiz_type):
                selected_topics = [t.strip() for t in (manual_topics or []) if str(t).strip()] if plan_mode == "manual_topics" else []
                existing_all = state.quiz[quiz_type]
                if selected_topics:
                    existing = [
                        q for q in existing_all
                        if (q.get("topic") or "").strip() in selected_topics
                    ]
                    if not existing:
                        raise ValueError("No existing questions match selected manual topics.")
                else:
                    existing = existing_all
                push({"agent": "QuizRefinerAgent", "step": "refine_start", "verdict": "PASS",
                      "phase": "PLAN",
                      "reason": f"Refining {len(existing)} questions with new instructions…"})
                push({"agent": "QuizRefinerAgent", "step": "phase_start", "phase": "GENERATE", "verdict": "PASS",
                      "reason": f"Generating refined set for {len(existing)} selected question(s)…",
                      "data": {"total_slots": len(existing)}})
                weak    = difficulty_adapter.get_weak_topics()
                schema  = QUIZ_SCHEMAS.get(quiz_type, QUIZ_SCHEMAS["mcq"])
                topics_line = ", ".join(selected_topics) if selected_topics else "auto"
                prompt  = (
                    f"Revise this {quiz_type.upper()} quiz per student instructions.\n"
                    f"Mode: {mode}  Difficulty: {difficulty.upper()}\n"
                    f"Topic scope: {topics_line}\n"
                    f"Weak topics: {', '.join(weak.keys()) if weak else 'none'}\n"
                    f"Instructions: {user_instr or 'None'}\n\n"
                    f"EXISTING QUIZ:\n{json.dumps(existing, ensure_ascii=False)}\n\n"
                    f"CONTENT:\n{content[:6000]}\n\nSUMMARY:\n{summary[:1500]}\n\n"
                    f"STRICT JSON: {schema}"
                )
                push({"agent": "QuizRefinerAgent", "step": "llm_prompt", "phase": "GENERATE", "verdict": "PASS",
                      "reason": "Prompt sent to LLM for refine generation",
                      "data": {"prompt_preview": prompt[:1200]}})
                push({"agent": "QuizRefinerAgent", "step": "generating", "verdict": "PASS",
                      "reason": "Sending revision prompt to LLM…"})
                data_out = call_ollama_json_quiz(prompt, key="questions")
                qs = data_out.get("questions", []) if isinstance(data_out, dict) else []
                passed, failed = QuizValidatorAgent().validate_batch(
                    qs, quiz_type, content[:800], difficulty, user_instr, thinking_cb=push
                )
                if failed:
                    fixed = QuizRefinerAgent().refine_batch(
                        failed, quiz_type, content[:800], difficulty, user_instr, thinking_cb=push
                    )
                    qs = passed + fixed
                else:
                    qs = passed
                fill_missing_timestamps(qs, segments=ctx.segment_list, board_entries=ctx.board_entries)
                state.quiz[quiz_type] = qs
                push({"agent": "QuizRefinerAgent", "step": "final", "verdict": "DONE",
                      "reason": f"Refined quiz ready — {len(qs)} questions"})

            elif mode == "adaptive":
                ctx._mode = "adaptive"
                agent = AdaptiveQuizAgent()
                agent.run(ctx)   # internally calls _run_quiz_agentic_loop via thinking_cb=push
                fill_missing_timestamps(
                    state.quiz.get(quiz_type, []),
                    segments=ctx.segment_list, board_entries=ctx.board_entries,
                )

            elif mode == "manual":
                agent_map = {"easy": EasyQuizAgent, "medium": MediumQuizAgent,
                             "hard": HardQuizAgent}
                agent = agent_map.get(difficulty, MediumQuizAgent)()
                ctx._mode       = "manual"
                ctx._agent_name = agent.name
                push({"agent": agent.name, "step": "start", "verdict": "PASS",
                      "reason": f"Manual {difficulty.upper()} mode — {num} {quiz_type.upper()} questions"})
                if user_instr:
                    push({"agent": agent.name, "step": "user_instructions", "verdict": "PASS",
                          "reason": f"Student instructions: {user_instr}"})
                max_attempts = 3
                final = _run_quiz_agentic_loop(
                    ctx, quiz_type, content, summary, user_instr,
                    target_difficulty=difficulty, thinking_cb=push,
                    max_attempts_per_q=max_attempts,
                )
                for q in final:
                    q["difficulty"] = difficulty
                fill_missing_timestamps(final, segments=ctx.segment_list, board_entries=ctx.board_entries)
                state.quiz[quiz_type] = final

            else:
                raise ValueError(f"Unknown mode: {mode}")

            final_qs = state.quiz.get(quiz_type, [])
            diff_label = difficulty if mode == "manual" else "adaptive"
            _append_generated_quiz_run("async", quiz_type, mode, diff_label, final_qs)
            with state.quiz_generation_lock:
                job = state.quiz_generation_jobs.get(job_id)
                if job:
                    job["status"] = "done"
                    job["result"] = {
                        "mode":       mode,
                        "difficulty": difficulty if mode == "manual" else "adaptive",
                        "total":      len(final_qs),
                        "questions":  final_qs,
                    }
        except Exception as e:
            with state.quiz_generation_lock:
                job = state.quiz_generation_jobs.get(job_id)
                if job:
                    job["status"] = "error"
                    job["error"]  = str(e)
        finally:
            state._quiz_thinking_active = False

    threading.Thread(target=worker, daemon=True).start()
    return jsonify(ok=True, job_id=job_id)



@app.route('/quiz/generate_job')
def quiz_generate_job_route():
    job_id = request.args.get("job_id", "")
    if not job_id:
        return jsonify(ok=False, error="Missing job_id"), 400
    with state.quiz_generation_lock:
        job = state.quiz_generation_jobs.get(job_id)
        if not job:
            return jsonify(ok=False, error="Job not found"), 404
        return jsonify(
            ok=True,
            status=job.get("status", ""),
            thinking=job.get("thinking", []),
            result=job.get("result", None),
            error=job.get("error", ""),
        )



@app.route('/performance')
def performance_route():
    stats = difficulty_adapter.get_all_stats()
    weak  = difficulty_adapter.get_weak_topics()
    return jsonify(stats=stats, weak_topics=weak, total_topics=len(stats))

# ── Video streaming (for in-UI playback at source timestamps) ─



def _stream_video_path(file_path, mimetype="video/mp4"):
    """Serve MP4 with HTTP Range so browsers can seek without downloading the whole file."""
    if not os.path.exists(file_path):
        return jsonify(error="File not found."), 404
    file_size = os.path.getsize(file_path)
    range_header = request.headers.get("Range")
    if not range_header:
        def gen_full():
            with open(file_path, "rb") as f:
                while True:
                    chunk = f.read(1024 * 1024)
                    if not chunk:
                        break
                    yield chunk

        return Response(
            gen_full(),
            status=200,
            mimetype=mimetype,
            headers={
                "Content-Length": str(file_size),
                "Accept-Ranges": "bytes",
                "Cache-Control": "no-cache",
            },
        )
    try:
        byte_range = range_header.replace("bytes=", "").split("-")
        start = int(byte_range[0]) if byte_range[0] else 0
        end = int(byte_range[1]) if byte_range[1] else file_size - 1
    except (ValueError, IndexError):
        return Response(status=416)
    end = min(max(start, end), file_size - 1)
    start = min(max(0, start), file_size - 1)
    if start > end:
        return Response(status=416)
    length = end - start + 1

    def gen_chunk():
        with open(file_path, "rb") as f:
            f.seek(start)
            remaining = length
            while remaining > 0:
                chunk = f.read(min(1024 * 1024, remaining))
                if not chunk:
                    break
                remaining -= len(chunk)
                yield chunk

    return Response(
        gen_chunk(),
        status=206,
        mimetype=mimetype,
        headers={
            "Content-Range": f"bytes {start}-{end}/{file_size}",
            "Accept-Ranges": "bytes",
            "Content-Length": str(length),
            "Cache-Control": "no-cache",
        },
    )



@app.route('/stream/video')
def stream_video_route():
    path = state.video_path
    if not path or not os.path.exists(path):
        return jsonify(error="No video available."), 404
    return _stream_video_path(path)



@app.route('/clip/video')
def clip_video_route():
    """Extract a short clip around a timestamp. ?t=MM:SS&dur=15"""
    path = state.video_path
    if not path or not os.path.exists(path):
        return jsonify(error="No video available."), 404

    ts_str = request.args.get("t", "00:00")
    dur    = min(int(request.args.get("dur", 20)), 90)

    # parse MM:SS or HH:MM:SS to seconds
    parts = ts_str.split(":")
    if len(parts) == 3:
        start = int(parts[0]) * 3600 + int(parts[1]) * 60 + int(parts[2])
    elif len(parts) == 2:
        start = int(parts[0]) * 60 + int(parts[1])
    else:
        start = int(float(ts_str))

    # Cut should start at the exact timestamp shown in the UI.
    start = max(0, start)
    # Suffix bumps cache when encode settings change (browser-safe yuv420p + main).
    clip_name = f"clip_{start}_{dur}_web.mp4"
    clip_path = os.path.join(CLIPS_DIR, clip_name)

    if not os.path.exists(clip_path):
        cmd = [
            "ffmpeg", "-y", "-ss", str(start), "-i", path,
            "-t", str(dur),
            "-c:v", "libx264", "-preset", "veryfast",
            "-profile:v", "main", "-pix_fmt", "yuv420p",
            "-c:a", "aac", "-b:a", "128k", "-ar", "44100", "-ac", "2",
            "-movflags", "+faststart",
            clip_path,
        ]
        try:
            subprocess.run(cmd, capture_output=True, timeout=120)
        except Exception as e:
            return jsonify(error=f"Clip extraction failed: {e}"), 500

    if not os.path.exists(clip_path):
        return jsonify(error="Clip extraction produced no output."), 500
    return _stream_video_path(clip_path)

# ── Downloads ─────────────────────────────────────────────────



DOWNLOAD_MAP = {
    "summary":    f"{OUTPUT_DIR}/summary.txt",
    "transcript": f"{OUTPUT_DIR}/transcript.txt",
    "board":      f"{OUTPUT_DIR}/board_text.txt",
    "timeline":   f"{OUTPUT_DIR}/lecture_timeline.txt",
    "eval":       f"{OUTPUT_DIR}/eval_results.json",   # evaluation harness results + metrics
    "bulk_eval":  f"{OUTPUT_DIR}/bulk_eval_results.json",  # bulk evaluation aggregates
    "bulk_csv":   f"{OUTPUT_DIR}/bulk_eval.csv",           # bulk eval, one row per question
    "gold":       f"{OUTPUT_DIR}/gold_calibration.json",   # solver-vs-gold calibration
}



@app.route('/download/<file_type>')
def download_route(file_type):
    path = DOWNLOAD_MAP.get(file_type)
    if path and os.path.exists(path):
        return send_file(path, as_attachment=True)
    return jsonify(error="File not found."), 404

# ── Print-friendly quiz export ────────────────────────────────



@app.route('/export/quiz')
def export_quiz_route():
    """Returns a self-contained printable HTML page of all quizzes."""
    q = state.quiz
    rows = []
    for i, m in enumerate(q.get("mcq", [])):
        opts = "".join(f"<li><b>{k}.</b> {v}</li>" for k, v in m.get("options", {}).items())
        rows.append(f"<div class='eq'><p><b>Q{i+1}.</b> {m.get('question','')}</p><ul>{opts}</ul></div>")
    for i, t in enumerate(q.get("tf", [])):
        rows.append(f"<div class='eq'><p><b>T/F {i+1}.</b> {t.get('statement','')}</p></div>")
    for i, f in enumerate(q.get("fill", [])):
        rows.append(f"<div class='eq'><p><b>Fill {i+1}.</b> {f.get('question','')}</p></div>")
    for i, s in enumerate(q.get("short", [])):
        rows.append(f"<div class='eq'><p><b>Short {i+1}.</b> ({s.get('marks',4)} marks) {s.get('question','')}</p>"
                     f"<div class='ans-space'></div></div>")
    body = "\n".join(rows) or "<p>No quiz data available.</p>"

    # answer key
    keys = []
    for i, m in enumerate(q.get("mcq", [])):
        keys.append(f"Q{i+1}: {m.get('correct_answer','')}")
    for i, t in enumerate(q.get("tf", [])):
        keys.append(f"T/F {i+1}: {t.get('answer','')}")
    for i, f in enumerate(q.get("fill", [])):
        keys.append(f"Fill {i+1}: {f.get('answer','')}")
    key_html = " &nbsp;|&nbsp; ".join(keys)

    html = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<title>LectureForge — Printable Quiz</title>
<style>
body{{font-family:Georgia,serif;max-width:750px;margin:40px auto;color:#1a1a2e;line-height:1.7}}
h1{{text-align:center;border-bottom:2px solid #333;padding-bottom:8px}}
.eq{{margin-bottom:18px;page-break-inside:avoid}}
.eq p{{margin:0 0 4px}} .eq ul{{margin:4px 0;padding-left:24px}}
.ans-space{{border-bottom:1px dashed #999;height:60px}}
.key{{margin-top:40px;padding-top:16px;border-top:2px solid #333;font-size:12px;color:#555}}
@media print{{.no-print{{display:none}}}}
</style></head><body>
<button class="no-print" onclick="window.print()" style="float:right;padding:8px 16px;cursor:pointer">🖨 Print / Save PDF</button>
<h1>LectureForge Quiz</h1>
{body}
<div class="key"><b>Answer Key:</b><br>{key_html}</div>
</body></html>"""
    return html, 200, {'Content-Type': 'text/html; charset=utf-8'}

# ── Evaluation harness ────────────────────────────────────────



@app.route('/eval/run', methods=['POST'])
def eval_run_route():
    """Trigger the agentic-vs-non-agentic evaluation harness remotely (blocking).
    Body: {quiz_types: ["mcq","tf","fill","short"], num_questions: 6, repeats: 2,
           judge_model: null, target_difficulty: "medium"}"""
    data          = request.get_json(silent=True) or {}
    quiz_types    = data.get("quiz_types") or ["mcq", "tf", "fill", "short"]
    num_questions = min(int(data.get("num_questions", 6)), 15)
    repeats       = min(int(data.get("repeats", 2)), 5)
    judge_model   = data.get("judge_model") or None
    target_diff   = data.get("target_difficulty", "medium")

    results = run_quiz_evaluation(
        quiz_types=tuple(quiz_types), num_questions=num_questions,
        repeats=repeats, judge_model=judge_model, target_difficulty=target_diff,
    )
    return jsonify(results)


# Background eval runner — the harness can take minutes (it judges many
# questions with the 12B model), so the UI kicks it off here and polls
# /eval/status for a LIVE feed rather than holding one long HTTP request open.
_eval_job      = {"status": "idle", "error": "", "started_at": "", "finished_at": "",
                  "phase": "", "log": [], "result": None, "pct": 0,
                  "_total": 0, "_done": 0}
_eval_job_lock = threading.Lock()


def _eval_to_line(ev):
    """Turn a progress event from run_quiz_evaluation into a short live line
    tagged by pipeline, for the UI feed."""
    k  = ev.get("kind")
    qt = (ev.get("quiz_type") or "").upper()
    if k == "start":
        return ("info", f"Starting benchmark — {len(ev.get('quiz_types', []))} type(s), "
                        f"{ev.get('num_questions')} Q each, judging up to {ev.get('judge_cap')} per side")
    if k == "unit_start":
        return ("info", f"{qt}  (unit {ev.get('unit')}/{ev.get('total_units')})")
    if k == "agentic_start":
        return ("agentic", f"{qt}: agentic pipeline started (plan → generate → validate → retry)")
    if k == "agentic_event":
        step = (ev.get("step") or "").replace("_", " ")
        rsn  = (ev.get("reason") or "")[:90]
        return ("agentic", f"{qt} · {step}: {rsn}")
    if k == "agentic_done":
        return ("agentic", f"{qt}: agentic produced {ev.get('count')} question(s) in {ev.get('secs')}s")
    if k == "nonagentic_start":
        return ("nonagentic", f"{qt}: non-agentic single-shot generation…")
    if k == "nonagentic_done":
        return ("nonagentic", f"{qt}: non-agentic produced {ev.get('count')} question(s) in {ev.get('secs')}s")
    if k == "sample":
        return (ev.get("pipeline", "info"),
                f"{qt} {ev.get('pipeline', '')} Q: {ev.get('text', '')}")
    if k == "judge_start":
        return ("judge", f"{qt}: judging {ev.get('n_agentic')} agentic + {ev.get('n_nonagentic')} non-agentic…")
    if k == "judge_progress":
        return ("judge", f"{qt}: judged {ev.get('pipeline')} {ev.get('done')}/{ev.get('total')}")
    if k == "unit_done":
        return ("result", f"{qt} done: agentic score {ev.get('agentic_score')} vs non-agentic "
                          f"{ev.get('nonagentic_score')} · accept {ev.get('agentic_accept')} vs "
                          f"{ev.get('nonagentic_accept')} · validator-judge F1 {ev.get('f1')}")
    if k == "done":
        return ("result", f"DONE — overall agentic {ev.get('agentic_score')} vs non-agentic "
                          f"{ev.get('nonagentic_score')} · F1 {ev.get('f1')}")
    # ── per-tab comparison (run_quiz_comparison) ──
    if k == "compare_start":
        return ("info", f"Comparing {ev.get('n')} existing {qt} question(s): agentic "
                        f"({ev.get('agentic_model')}) vs non-agentic ({ev.get('nonagentic_model')}), "
                        f"judge: {ev.get('judge_model')}")
    if k == "nonagentic_model":
        return ("nonagentic", f"Non-agentic generator model: {ev.get('model')}")
    if k == "nonagentic_gen":
        return ("nonagentic", f"Non-agentic generating on topic '{ev.get('topic')}' "
                              f"({ev.get('idx')}/{ev.get('total')})…")
    if k == "nonagentic_q":
        return ("nonagentic", f"Non-agentic Q [{ev.get('topic')}]: {ev.get('text')}")
    if k == "judge_model":
        return ("judge", f"Judge model: {ev.get('model')}")
    if k == "pair_done":
        return ("result", f"[{ev.get('topic')}] {ev.get('idx')}/{ev.get('total')}: "
                          f"agentic {ev.get('agentic_overall')}/5 vs non-agentic "
                          f"{ev.get('nonagentic_overall')}/5")
    if k == "compare_done":
        return ("result", f"DONE — agentic {ev.get('overall_agentic')}/5 vs non-agentic "
                          f"{ev.get('overall_nonagentic')}/5 · F1 {ev.get('f1')}")
    if k == "error":
        return ("error", ev.get("reason", "error"))
    return ("info", k or "…")


def _eval_push(ev):
    tag, line = _eval_to_line(ev)
    k = ev.get("kind")
    with _eval_job_lock:
        # progress tracking → pct (so the UI bar doesn't depend on glyphs)
        if k == "start":
            _eval_job["_total"] = ev.get("total_units") or 0
            _eval_job["_done"]  = 0
            _eval_job["pct"]    = 0
        elif k == "compare_start":
            _eval_job["_total"] = ev.get("n") or 0
            _eval_job["_done"]  = 0
            _eval_job["pct"]    = 0
        elif k in ("unit_done", "pair_done"):
            _eval_job["_done"] = _eval_job.get("_done", 0) + 1
            tot = _eval_job.get("_total", 0) or 0
            _eval_job["pct"] = round(_eval_job["_done"] / tot * 100) if tot else 0
        elif k in ("done", "compare_done"):
            _eval_job["pct"] = 100

        _eval_job["phase"] = line if tag in ("info", "result") else _eval_job.get("phase", "")
        _eval_job["log"].append({"tag": tag, "line": line,
                                 "t": datetime.now().strftime("%H:%M:%S")})
        if len(_eval_job["log"]) > 400:
            _eval_job["log"] = _eval_job["log"][-400:]


@app.route('/eval/run_async', methods=['POST'])
def eval_run_async_route():
    """Start the evaluation harness in a background thread with a LIVE feed.
    Poll /eval/status for status + log. Defaults kept small so it finishes."""
    data          = request.get_json(silent=True) or {}
    quiz_types    = tuple(data.get("quiz_types") or ["mcq", "tf", "fill", "short"])
    num_questions = min(int(data.get("num_questions", 3)), 15)
    repeats       = min(int(data.get("repeats", 1)), 5)
    judge_model   = data.get("judge_model") or None
    target_diff   = data.get("target_difficulty", "medium")

    if not state.status.get("summary"):
        return jsonify(ok=False, error="No content processed yet — upload and process a lecture first.")

    with _eval_job_lock:
        if _eval_job["status"] == "running":
            return jsonify(ok=False, error="An evaluation is already running.", status="running")
        _eval_job.update(status="running", error="", phase="Starting…", log=[],
                         pct=0, _total=0, _done=0,
                         started_at=datetime.now().isoformat(), finished_at="")

    def worker():
        try:
            res = run_quiz_evaluation(
                quiz_types=quiz_types, num_questions=num_questions,
                repeats=repeats, judge_model=judge_model, target_difficulty=target_diff,
                progress_cb=_eval_push,
            )
            if isinstance(res, dict) and not res.get("ok", False):
                with _eval_job_lock:
                    _eval_job.update(status="error", error=res.get("error", "evaluation failed"),
                                     finished_at=datetime.now().isoformat())
                return
            with _eval_job_lock:
                _eval_job.update(status="done", finished_at=datetime.now().isoformat())
        except Exception as e:
            with _eval_job_lock:
                _eval_job.update(status="error", error=str(e),
                                 finished_at=datetime.now().isoformat())

    threading.Thread(target=worker, daemon=True).start()
    return jsonify(ok=True, status="running")


@app.route('/eval/compare_async', methods=['POST'])
def eval_compare_async_route():
    """Per-tab evaluation: judge the ALREADY-generated questions of one quiz
    type against a non-agentic baseline built on the same topics. Runs in the
    background with the same live feed; final structured comparison is returned
    in /eval/status's `result`."""
    data        = request.get_json(silent=True) or {}
    quiz_type   = data.get("quiz_type", "mcq")
    judge_model = data.get("judge_model") or None
    target_diff = data.get("target_difficulty", "medium")

    if not state.quiz.get(quiz_type):
        return jsonify(ok=False, error=f"No {quiz_type.upper()} questions generated yet — generate first.")

    with _eval_job_lock:
        if _eval_job["status"] == "running":
            return jsonify(ok=False, error="An evaluation is already running.", status="running")
        _eval_job.update(status="running", error="", phase="Starting…", log=[], result=None,
                         pct=0, _total=0, _done=0,
                         started_at=datetime.now().isoformat(), finished_at="")

    def worker():
        try:
            res = run_quiz_comparison(quiz_type, judge_model=judge_model,
                                      target_difficulty=target_diff, progress_cb=_eval_push)
            with _eval_job_lock:
                if isinstance(res, dict) and not res.get("ok", False):
                    _eval_job.update(status="error", error=res.get("error", "comparison failed"),
                                     finished_at=datetime.now().isoformat())
                else:
                    _eval_job.update(status="done", result=res,
                                     finished_at=datetime.now().isoformat())
        except Exception as e:
            with _eval_job_lock:
                _eval_job.update(status="error", error=str(e),
                                 finished_at=datetime.now().isoformat())

    threading.Thread(target=worker, daemon=True).start()
    return jsonify(ok=True, status="running")


@app.route('/eval/status')
def eval_status_route():
    with _eval_job_lock:
        job = dict(_eval_job)
        job["log"] = list(_eval_job["log"])
    job["has_results"] = os.path.exists(os.path.join(OUTPUT_DIR, "eval_results.json"))
    return jsonify(**job)


# ── Bulk evaluation: background job + live feed (resumable) ────
_bulk_job      = {"status": "idle", "error": "", "phase": "", "log": [],
                  "result": None, "pct": 0, "started_at": "", "finished_at": ""}
_bulk_job_lock = threading.Lock()


def _bulk_to_line(ev):
    k = ev.get("kind")
    if k == "bulk_start":
        return ("info", f"Bulk eval: {ev.get('total')} MCQ across {', '.join(ev.get('counts', {}).keys())} "
                        f"| pipelines: {', '.join(ev.get('pipelines', []))} | judge {ev.get('judge_model')}")
    if k == "group_gen":
        return (ev.get("pipeline", "info"),
                f"Generating {ev.get('count')} {ev.get('pipeline')} MCQ ({ev.get('difficulty')})…")
    if k == "group_done":
        return (ev.get("pipeline", "info"),
                f"{ev.get('pipeline')} ({ev.get('difficulty')}): {ev.get('count')} questions ready")
    if k == "judge_start":
        return ("judge", f"Judging + answer-checking {ev.get('total')} questions "
                         f"(resuming from {ev.get('already')})…")
    if k == "judge_progress":
        return ("judge", f"{ev.get('done')}/{ev.get('total')} judged "
                         f"({ev.get('pipeline')} {ev.get('difficulty')})")
    if k == "gold_start":
        return ("judge", f"Gold cross-check: both models answer {ev.get('n')} known MCQ…")
    if k == "gold_progress":
        return ("judge", f"gold [{ev.get('model')}] {ev.get('done')}/{ev.get('total')}")
    if k == "bulk_done":
        return ("result", f"DONE — {ev.get('n_rows')} questions evaluated")
    if k == "error":
        return ("error", ev.get("reason", "error"))
    return ("info", k or "…")


def _bulk_push(ev):
    tag, line = _bulk_to_line(ev)
    k = ev.get("kind")
    with _bulk_job_lock:
        if k == "judge_start":
            _bulk_job["_total"] = ev.get("total") or 0
            _bulk_job["_done"]  = ev.get("already") or 0
        elif k == "judge_progress":
            _bulk_job["_done"] = ev.get("done") or _bulk_job.get("_done", 0)
            tot = _bulk_job.get("_total", 0) or 0
            _bulk_job["pct"] = round(_bulk_job["_done"] / tot * 100) if tot else 0
        elif k == "bulk_done":
            _bulk_job["pct"] = 100
        _bulk_job["phase"] = line if tag in ("info", "result") else _bulk_job.get("phase", "")
        _bulk_job["log"].append({"tag": tag, "line": line,
                                 "t": datetime.now().strftime("%H:%M:%S")})
        if len(_bulk_job["log"]) > 600:
            _bulk_job["log"] = _bulk_job["log"][-600:]


@app.route('/eval/bulk_async', methods=['POST'])
def eval_bulk_async_route():
    """Start the bulk evaluation (N MCQ x difficulties, agentic vs non-agentic,
    answer-correctness) in the background. Poll /eval/bulk_status."""
    data       = request.get_json(silent=True) or {}
    total      = min(int(data.get("total", 200)), 400)
    compare    = bool(data.get("compare", True))
    fresh      = bool(data.get("fresh", False))
    judge_model = data.get("judge_model") or None
    gold_path  = data.get("gold_path") or ""

    if not state.status.get("summary"):
        return jsonify(ok=False, error="No content processed yet — upload and process a lecture first.")

    with _bulk_job_lock:
        if _bulk_job["status"] == "running":
            return jsonify(ok=False, error="A bulk evaluation is already running.", status="running")
        _bulk_job.update(status="running", error="", phase="Starting…", log=[], result=None,
                         pct=0, _total=0, _done=0,
                         started_at=datetime.now().isoformat(), finished_at="")

    def worker():
        try:
            res = run_bulk_evaluation(total=total, compare=compare, fresh=fresh,
                                      judge_model=judge_model, gold_path=gold_path,
                                      progress_cb=_bulk_push)
            with _bulk_job_lock:
                if isinstance(res, dict) and not res.get("ok", False):
                    _bulk_job.update(status="error", error=res.get("error", "bulk eval failed"),
                                     finished_at=datetime.now().isoformat())
                else:
                    _bulk_job.update(status="done", result=res,
                                     finished_at=datetime.now().isoformat())
        except Exception as e:
            with _bulk_job_lock:
                _bulk_job.update(status="error", error=str(e),
                                 finished_at=datetime.now().isoformat())

    threading.Thread(target=worker, daemon=True).start()
    return jsonify(ok=True, status="running")


@app.route('/eval/bulk_status')
def eval_bulk_status_route():
    with _bulk_job_lock:
        job = dict(_bulk_job)
        job["log"] = list(_bulk_job["log"])
    job["has_results"] = os.path.exists(BULK_RESULTS_PATH)
    return jsonify(**job)



print("✅ Flask API defined.")

# ══════════════════════════════════════════════════════════════
# LAUNCH  (Flask)  —  Jupyter / RunPod (no ngrok; use the HTTP proxy)
# ══════════════════════════════════════════════════════════════
# In a notebook we run Flask in a BACKGROUND THREAD so the cell returns and the
# kernel stays interactive (you can then trigger the eval from the UI, or call
# run_bulk_evaluation(...) directly in Cell 3). Ollama was already started and
# the models pulled by setup_ollama() above.

import requests

def check_ollama():
    try:
        r = requests.get(f"{OLLAMA_BASE}/api/tags", timeout=3)
        if r.status_code == 200:
            tags = [m.get("name") for m in r.json().get("models", [])]
            print("✅ Ollama is RUNNING — models:", tags)
        else:
            print("⚠️ Ollama responded oddly:", r.status_code)
    except Exception as e:
        print("❌ Ollama not reachable — run `ollama serve` in the pod. Error:", e)

check_ollama()

def _run_flask():
    app.run(host='0.0.0.0', port=API_PORT, debug=False, use_reloader=False)

threading.Thread(target=_run_flask, daemon=True).start()
time.sleep(2)

print(f"""
{'='*64}
  🎓 LectureForge — RunPod backend LIVE (background thread)
{'='*64}
  Flask listening on 0.0.0.0:{API_PORT}

  Expose HTTP port {API_PORT} in your RunPod template. RunPod then gives a URL:
      https://<POD_ID>-{API_PORT}.proxy.runpod.net
  Paste THAT url into LectureAssis.html → Connect.

  Model : {state.active_model}
  Data  : {BASE_DIR}
  CORS  : enabled (local UI supported)
{'='*64}
The kernel is free — process a lecture from the UI, or run Cell 3 below.
""")


# %% [Cell 3 — OPTIONAL: run the bulk evaluation directly from the notebook]
# Requires a lecture to have been processed first (via the UI or your own call).
# results = run_bulk_evaluation(
#     total=12,                                   # smoke test; raise to 200 for the real run
#     compare=True,
#     gold_path="/workspace/LectureForge/gold_limits.json",  # optional gold cross-check
# )
# import json; print(json.dumps(results["aggregated"], indent=2))



